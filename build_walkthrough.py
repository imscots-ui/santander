"""
build_walkthrough.py
Builds Santander_LimitedCompany_Walkthrough.pptx — a ~28-slide portfolio deck
for the Santander Business Banking prototype.

Run:  python build_walkthrough.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ──────────────────────────────────────────────────────────────
RED       = RGBColor(0xC8, 0x10, 0x2E)   # Santander red
DARK      = RGBColor(0x1C, 0x19, 0x17)   # stone-900
WARM_BG   = RGBColor(0xFA, 0xF6, 0xEF)   # warm off-white
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
MID_GREY  = RGBColor(0x78, 0x71, 0x69)   # stone-500
LIGHT_BG  = RGBColor(0xF5, 0xF2, 0xED)   # stone-100
AMBER     = RGBColor(0xD9, 0x77, 0x06)   # amber-600
BLUE      = RGBColor(0x25, 0x63, 0xEB)   # blue-600
GREEN     = RGBColor(0x05, 0x96, 0x69)   # emerald-600
SECTION_MUTED = RGBColor(0xA8, 0xA2, 0x99)  # stone-400

# ── Slide dimensions (16:9) ────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Low-level drawing helpers ─────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text,
                font_size=14, bold=False, color=DARK,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb


def add_para(tf, text, font_size=13, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p   = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p


def set_para(p, text, font_size=13, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, italic=False):
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return p


# ── Structural slide builders ─────────────────────────────────────────────────

def cover_slide(prs, title, subtitle, date_str):
    """Dark full-bleed cover."""
    slide = prs.slides.add_slide(blank_layout)
    # Dark background
    add_rect(slide, 0, 0, W, H, fill=DARK)
    # Red accent bar (left edge)
    add_rect(slide, 0, 0, Inches(0.18), H, fill=RED)
    # Red accent strip (bottom)
    add_rect(slide, 0, H - Inches(0.12), W, Inches(0.12), fill=RED)

    # Classification tag (top-right)
    add_rect(slide, Inches(10.95), Inches(0.30), Inches(1.95), Inches(0.34), fill=RED)
    add_textbox(slide, Inches(10.95), Inches(0.34), Inches(1.95), Inches(0.26),
                "INTERNAL · CONFIDENTIAL", font_size=7.5, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)

    # Company wordmark area
    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(5), Inches(0.5),
                "Santander Business Banking",
                font_size=11, color=RGBColor(0xA8,0xA2,0x99), bold=False)

    # Big title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(10), Inches(2.2))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r  = p.add_run()
    r.text = title
    r.font.size  = Pt(52)
    r.font.bold  = True
    r.font.color.rgb = WHITE
    r.font.name  = "Calibri"

    # Subtitle
    add_textbox(slide, Inches(0.5), Inches(4.2), Inches(9), Inches(0.6),
                subtitle, font_size=18, color=RGBColor(0xC8,0x10,0x2E), bold=True)

    # Date
    add_textbox(slide, Inches(0.5), Inches(4.95), Inches(6), Inches(0.4),
                date_str, font_size=13, color=RGBColor(0xA8,0xA2,0x99))

    # Bottom-right: URL
    add_textbox(slide, Inches(8.5), Inches(6.8), Inches(4.5), Inches(0.4),
                "imscots-ui.github.io/santander",
                font_size=11, color=RGBColor(0x78,0x71,0x69), align=PP_ALIGN.RIGHT)
    return slide


def section_divider(prs, section_label, section_title, section_sub=""):
    """Dark section-break slide."""
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, W, H, fill=DARK)
    add_rect(slide, 0, 0, Inches(0.18), H, fill=RED)
    add_rect(slide, Inches(0.18), H - Inches(0.12), W, Inches(0.12), fill=RED)

    add_textbox(slide, Inches(0.5), Inches(2.2), Inches(12), Inches(0.5),
                section_label, font_size=11, color=RGBColor(0xC8,0x10,0x2E), bold=True)
    add_textbox(slide, Inches(0.5), Inches(2.8), Inches(11), Inches(1.4),
                section_title, font_size=40, bold=True, color=WHITE)
    if section_sub:
        add_textbox(slide, Inches(0.5), Inches(4.5), Inches(10), Inches(0.6),
                    section_sub, font_size=16, color=RGBColor(0xA8,0xA2,0x99))
    return slide


LIGHTRED = RGBColor(0xFE, 0xCA, 0xCA)   # light red for dark footer
STONE3   = RGBColor(0xD6, 0xD3, 0xD1)   # stone-300

def add_confidential_footer(slide):
    """Dark footer bar: classification (left) + credit (right). Standard for internal bank decks."""
    add_rect(slide, 0, H - Inches(0.30), W, Inches(0.30), fill=DARK)
    add_textbox(slide, Inches(0.3), H - Inches(0.275), Inches(4), Inches(0.22),
                "INTERNAL · CONFIDENTIAL", font_size=8, bold=True, color=LIGHTRED)
    add_textbox(slide, W - Inches(7.0), H - Inches(0.275), Inches(6.7), Inches(0.22),
                "Santander Business Banking · Limited Company Walkthrough · June 2026",
                font_size=8, color=STONE3, align=PP_ALIGN.RIGHT)

def content_slide(prs, section_label, title, bullets, notes=""):
    """
    Standard white content slide.
    bullets: list of strings. Prefix with '  ' (2 spaces) for sub-bullet.
    """
    slide = prs.slides.add_slide(blank_layout)
    # Warm off-white background
    add_rect(slide, 0, 0, W, H, fill=WARM_BG)
    # Red top bar
    add_rect(slide, 0, 0, W, Inches(0.08), fill=RED)
    # Light sidebar
    add_rect(slide, 0, Inches(0.08), Inches(0.06), H - Inches(0.08), fill=LIGHT_BG)

    # Section label
    add_textbox(slide, Inches(0.3), Inches(0.18), Inches(8), Inches(0.3),
                section_label.upper(), font_size=8,
                color=SECTION_MUTED, bold=True)

    # Title
    add_textbox(slide, Inches(0.3), Inches(0.52), Inches(12.5), Inches(0.85),
                title, font_size=28, bold=True, color=DARK)

    # Divider line
    line = slide.shapes.add_shape(1, Inches(0.3), Inches(1.42), Inches(12.6), Pt(1.5))
    line.fill.solid(); line.fill.fore_color.rgb = RED
    line.line.fill.background()

    # Body text box
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(1.55), Inches(12.6), Inches(5.6))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True

    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        is_sub   = b.startswith("    ")   # 4 spaces = sub-sub
        is_child = b.startswith("  ") and not is_sub  # 2 spaces = sub
        text = b.lstrip()

        if text == "":
            p.space_before = Pt(4)
            continue

        if text.startswith("##"):
            # Section heading within slide
            run = p.add_run()
            run.text = text[2:].strip()
            run.font.size  = Pt(14)
            run.font.bold  = True
            run.font.color.rgb = RED
            run.font.name  = "Calibri"
            p.space_before = Pt(8)
        elif is_sub:
            p.level = 1
            p.space_before = Pt(1)
            run = p.add_run()
            run.text = "—  " + text
            run.font.size  = Pt(11)
            run.font.color.rgb = MID_GREY
            run.font.name  = "Calibri"
        elif is_child:
            p.level = 1
            p.space_before = Pt(2)
            run = p.add_run()
            run.text = "•  " + text
            run.font.size  = Pt(12)
            run.font.color.rgb = DARK
            run.font.name  = "Calibri"
        else:
            p.space_before = Pt(4)
            run = p.add_run()
            is_bullet = text.startswith("•")
            if is_bullet:
                run.text = text
            else:
                run.text = text
            run.font.size  = Pt(13)
            run.font.bold  = False
            run.font.color.rgb = DARK
            run.font.name  = "Calibri"

    # Optional presenter notes
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    add_confidential_footer(slide)
    return slide


def two_col_slide(prs, section_label, title, left_bullets, right_bullets, notes=""):
    """Two-column layout."""
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, W, H, fill=WARM_BG)
    add_rect(slide, 0, 0, W, Inches(0.08), fill=RED)
    add_rect(slide, 0, Inches(0.08), Inches(0.06), H - Inches(0.08), fill=LIGHT_BG)

    add_textbox(slide, Inches(0.3), Inches(0.18), Inches(8), Inches(0.3),
                section_label.upper(), font_size=8, color=SECTION_MUTED, bold=True)
    add_textbox(slide, Inches(0.3), Inches(0.52), Inches(12.5), Inches(0.85),
                title, font_size=28, bold=True, color=DARK)
    line = slide.shapes.add_shape(1, Inches(0.3), Inches(1.42), Inches(12.6), Pt(1.5))
    line.fill.solid(); line.fill.fore_color.rgb = RED
    line.line.fill.background()

    # Left column
    lx, ly, lw = Inches(0.3), Inches(1.58), Inches(6.1)
    tb = slide.shapes.add_textbox(lx, ly, lw, Inches(5.5))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for b in left_bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        is_child = b.startswith("  ")
        text = b.lstrip()
        if text == "":
            p.space_before = Pt(4); continue
        if text.startswith("##"):
            run = p.add_run(); run.text = text[2:].strip()
            run.font.size = Pt(14); run.font.bold = True
            run.font.color.rgb = RED; run.font.name = "Calibri"
            p.space_before = Pt(6)
        elif is_child:
            p.level = 1; p.space_before = Pt(2)
            run = p.add_run(); run.text = "—  " + text
            run.font.size = Pt(11); run.font.color.rgb = MID_GREY; run.font.name = "Calibri"
        else:
            p.space_before = Pt(3)
            run = p.add_run(); run.text = text
            run.font.size = Pt(12.5); run.font.color.rgb = DARK; run.font.name = "Calibri"

    # Vertical divider
    div = slide.shapes.add_shape(1, Inches(6.55), Inches(1.58), Pt(1.5), Inches(5.5))
    div.fill.solid(); div.fill.fore_color.rgb = LIGHT_BG
    div.line.fill.background()

    # Right column
    rx, ry, rw = Inches(6.7), Inches(1.58), Inches(6.1)
    tb2 = slide.shapes.add_textbox(rx, ry, rw, Inches(5.5))
    tb2.word_wrap = True
    tf2 = tb2.text_frame; tf2.word_wrap = True
    first = True
    for b in right_bullets:
        p = tf2.paragraphs[0] if first else tf2.add_paragraph()
        first = False
        is_child = b.startswith("  ")
        text = b.lstrip()
        if text == "":
            p.space_before = Pt(4); continue
        if text.startswith("##"):
            run = p.add_run(); run.text = text[2:].strip()
            run.font.size = Pt(14); run.font.bold = True
            run.font.color.rgb = RED; run.font.name = "Calibri"
            p.space_before = Pt(6)
        elif is_child:
            p.level = 1; p.space_before = Pt(2)
            run = p.add_run(); run.text = "—  " + text
            run.font.size = Pt(11); run.font.color.rgb = MID_GREY; run.font.name = "Calibri"
        else:
            p.space_before = Pt(3)
            run = p.add_run(); run.text = text
            run.font.size = Pt(12.5); run.font.color.rgb = DARK; run.font.name = "Calibri"

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    add_confidential_footer(slide)
    return slide


def back_cover(prs, contact, disclaimer):
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, W, H, fill=DARK)
    add_rect(slide, 0, 0, Inches(0.18), H, fill=RED)
    add_rect(slide, 0, H - Inches(0.12), W, Inches(0.12), fill=RED)
    add_textbox(slide, Inches(0.5), Inches(2.0), Inches(12), Inches(1.2),
                "Santander Business Banking", font_size=34, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.5), Inches(3.3), Inches(12), Inches(0.5),
                "Prototype — Concept piece for internal discussion only",
                font_size=16, color=RED, bold=True)
    add_textbox(slide, Inches(0.5), Inches(4.1), Inches(11), Inches(0.45),
                contact, font_size=13, color=RGBColor(0xA8,0xA2,0x99))
    add_textbox(slide, Inches(0.5), Inches(4.7), Inches(11.5), Inches(1.4),
                disclaimer, font_size=10, color=RGBColor(0x78,0x71,0x69))
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(8), Inches(0.4),
                "imscots-ui.github.io/santander",
                font_size=11, color=RGBColor(0x78,0x71,0x69))
    add_textbox(slide, Inches(9.5), Inches(6.5), Inches(3.5), Inches(0.4),
                "June 2026", font_size=11, color=RGBColor(0x78,0x71,0x69),
                align=PP_ALIGN.RIGHT)
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# BUILD ALL 28 SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

# ── SLIDE 1: Cover ─────────────────────────────────────────────────────────────
cover_slide(prs,
    title    = "Limited Company\nWalkthrough",
    subtitle = "Santander Business Banking · Concept Prototype",
    date_str = "June 2026  ·  Created by Alan Davidson")

# ── SLIDE 2: Section A divider ─────────────────────────────────────────────────
section_divider(prs,
    section_label = "Section A  ·  Slides 2–4",
    section_title = "About & Access",
    section_sub   = "What the prototype is, who it's for, and how to open it")

# ── SLIDE 3: What this is ──────────────────────────────────────────────────────
content_slide(prs,
    section_label = "A  ·  About & Access",
    title = "What this prototype is",
    bullets = [
        "• A self-initiated concept piece exploring paperless workflows for Santander Business Banking.",
        "• Front-end only — no backend, no live data, no real system connections.",
        "• Built in React (single App.jsx, ~5,700 lines) and deployed on GitHub Pages.",
        "",
        "##Who it is for",
        "• Internal product and design teams — exploring what a modern business banking experience could look like.",
        "• Customer research sessions — showing participants realistic flows without touching live systems.",
        "• Presentations to senior stakeholders — demonstrating the art of the possible.",
        "",
        "##What it demonstrates",
        "• Paperless mandate changes, bulk wages payments, KYC/KYB verification",
        "• MTD / HMRC VAT submissions, account closures, FX international payments",
        "• Account unlinking, credit ring-fencing, dormancy management, and lending draw-down",
        "• A 'Command Centre' home screen that surfaces proactive insight across all accounts",
    ],
    notes="Emphasise: nothing here is connected to any live Santander system. All data is fabricated for demo purposes.")

# ── SLIDE 4: How to access ─────────────────────────────────────────────────────
content_slide(prs,
    section_label = "A  ·  About & Access",
    title = "How to access the prototype",
    bullets = [
        "##URL",
        "• https://imscots-ui.github.io/santander",
        "  — Works in any modern browser. No login required.",
        "",
        "##Recommended setup",
        "• Mobile browser (Chrome / Safari) — the default view is a phone shell at 390 px wide.",
        "• For a laptop demo, open on a phone or use Chrome DevTools device emulation (iPhone 14 Pro).",
        "• A Desktop mode toggle in the header switches to a full sidebar layout for larger screens.",
        "",
        "##Troubleshooting",
        "• Hard-refresh (Cmd+Shift+R / Ctrl+Shift+R) if the screen looks blank or stale.",
        "• Use an Incognito / Private window to clear any cached state.",
        "• The prototype is stateless — refresh resets everything to the opening defaults.",
        "",
        "##No installation needed",
        "• Public GitHub Pages URL — no VPN, no internal network, no credentials.",
    ])

# ── SLIDE 5: Section B divider ─────────────────────────────────────────────────
section_divider(prs,
    section_label = "Section B  ·  Slides 5–8",
    section_title = "Navigation Guide",
    section_sub   = "The mobile shell, desktop mode, demo controls, and tab overview")

# ── SLIDE 6: Mobile shell ──────────────────────────────────────────────────────
content_slide(prs,
    section_label = "B  ·  Navigation Guide",
    title = "The mobile shell — five tabs",
    bullets = [
        "The prototype opens as a mobile phone shell with a bottom navigation bar. Five tabs:",
        "",
        "##  Home — Command Centre",
        "  • The main screen. 15+ cards covering balances, forecasts, approvals, and actions.",
        "  • This is the headline feature — everything about the business on one screen.",
        "",
        "##  Approve",
        "  • Pending signature queue. Items waiting for the logged-in director to countersign.",
        "  • OTP or biometric confirmation. Shows expiry countdown per request.",
        "",
        "##  Audit",
        "  • Immutable log of every action taken in the system.",
        "  • Export options for compliance, auditors, and accountants.",
        "",
        "##  Tax (MTD)",
        "  • Making Tax Digital hub. VAT obligations, ITSA, and HMRC submission workflow.",
        "  • Live countdown to the next VAT deadline. One-tap direct submission to HMRC.",
        "",
        "##  Statements",
        "  • Full transaction history. Filter by date, category, counterparty, and payment method.",
        "  • Tabular-figure monetary amounts. Export as CSV or PDF.",
    ])

# ── SLIDE 7: Desktop mode ──────────────────────────────────────────────────────
content_slide(prs,
    section_label = "B  ·  Navigation Guide",
    title = "Desktop mode — sidebar layout",
    bullets = [
        "• A 'Desktop' pill toggle sits in the top header bar (visible in both views).",
        "• Tapping it switches from the mobile phone shell to a full-width sidebar layout.",
        "",
        "##What changes",
        "  — Navigation moves from a bottom bar to a left sidebar.",
        "  — Content columns widen to use the full viewport.",
        "  — The same state is shared — switching mid-demo preserves any workflow in progress.",
        "  — Useful for presenting on a laptop or projector without device emulation.",
        "",
        "##What stays the same",
        "  — All 15 Command Centre components are identical.",
        "  — All workflows are identical — step count, logic, and data.",
        "  — Entity switcher and demo controls remain accessible.",
        "",
        "##Presenter tip",
        "  — Start on mobile to show the customer experience, then toggle to desktop",
        "    to show a business banking manager's view on the same data.",
    ])

# ── SLIDE 8: Demo controls ─────────────────────────────────────────────────────
content_slide(prs,
    section_label = "B  ·  Navigation Guide",
    title = "Demo controls — three pills on the Home screen",
    bullets = [
        "Three small control pills appear near the top of the Home screen (Command Centre).",
        "These are for demo and research use — they would not appear in a production app.",
        "",
        "##1. Entity switcher  [Limited Company ›]",
        "  — Cycles through all 7 entity types: Limited Company, Sole Trader,",
        "    Partnership, LLP, Charity, Club, Society.",
        "  — Instantly updates: greeting, entity name, account list, mandate rules,",
        "    signatory labels (director / partner / trustee / member), and all copy.",
        "  — Companies House number shown for Limited and LLP. Charity Commission",
        "    number shown for registered charities.",
        "",
        "##2. Demo cooling-off",
        "  — Injects a mock 24-hour cooling-off request into the Command Centre.",
        "  — Shows the amber countdown card with live progress bar, Cancel, and 'Do it now'.",
        "  — Useful for demonstrating the regulatory 24h delay without running a full workflow.",
        "",
        "##3. Simulate timeout",
        "  — Injects a stalled co-signer request ('Priya's on it' card).",
        "  — Demonstrates what happens when a co-signer misses their approval window.",
        "  — Shows the RM escalation path: Priya saw this within 2 minutes, 'Call Priya' button.",
    ])

# ── SLIDE 9: Tab overview ──────────────────────────────────────────────────────
content_slide(prs,
    section_label = "B  ·  Navigation Guide",
    title = "What each tab contains — at a glance",
    bullets = [
        "##Home (Command Centre)",
        "  — 15 components. Use for: opening demos, showing the breadth of insight,",
        "    launching any workflow, switching entity type mid-session.",
        "",
        "##Approve",
        "  — Pending items awaiting the director's signature.",
        "  — Three demo items: Bulk Payment (October Wages £183k), Mandate Change",
        "    (Add Mark Patel), Account Closure (Legacy Trading ····4422).",
        "  — Use for: demonstrating dual-authorisation, OTP flow, biometric sign-off.",
        "",
        "##Audit",
        "  — Chronological log of all events. Tamper-evident by design.",
        "  — Use for: compliance conversations, showing auditability of paperless workflows.",
        "",
        "##Tax (MTD)",
        "  — VAT obligations: Q3 Jul–Sep 2026 (open, £18,472 due), Q2 submitted.",
        "  — HMRC MTD VAT API v1.0. 247 transactions to review.",
        "  — Use for: accountant demos, MTD readiness conversations, VAT submission flow.",
        "",
        "##Statements",
        "  — Transaction history with full filtering: date range, category, counterparty,",
        "    payment method (card / DD / SO / FP / BACS).",
        "  — Counterparty deep-dive with spend summary and transaction list.",
    ])

# ── SLIDE 10: Section C divider ────────────────────────────────────────────────
section_divider(prs,
    section_label = "Section C  ·  Slides 10–18",
    section_title = "Command Centre\nDeep-Dive",
    section_sub   = "The Home screen — 15 components, one screen that knows your business")

# ── SLIDE 11: Command Centre overview ─────────────────────────────────────────
content_slide(prs,
    section_label = "C  ·  Command Centre",
    title = "Command Centre — philosophy and components",
    bullets = [
        "The Home screen is called the 'Command Centre'. Its design principle:",
        "\"One screen that knows everything about your business — before you ask.\"",
        "",
        "##The 15 components (top to bottom)",
        "•  1.  Editorial greeting + date",
        "•  2.  Total Balance hero card (dark, with Santander red floating orb)",
        "•  3.  Proactive Forecast strip — 3 colour-coded nudge cards",
        "•  4.  Demo controls (entity switcher, cooling-off, simulate timeout)",
        "•  5.  Session anomaly alert (new sign-in detected — Edinburgh · new device)",
        "•  6.  Pre-approved lending offer — £45,000 at 6.2% APR (or active loan card)",
        "•  7.  Business Health Score card",
        "•  8.  Cooling-off cards — live progress bar, cancel or 'Do it now'",
        "•  9.  Stalled requests — 'Priya's on it' RM escalation",
        "• 10.  MTD VAT teaser — amount owed, days to deadline, transactions to review",
        "• 11.  Paperless banner — '3 paper forms retired · No more posting to Sunderland'",
        "• 12.  Pending approvals — items with expiry countdown",
        "• 13.  Paperless actions grid — 10 workflow tiles",
        "• 14.  Accounts list — balances, sort codes, mandate badges",
        "• 15.  Supplier Radar, 13-week forecast chart, Director Centre, Cards, Privacy controls",
    ])

# ── SLIDE 12: Greeting + hero card ────────────────────────────────────────────
content_slide(prs,
    section_label = "C  ·  Command Centre",
    title = "Editorial greeting + Total Balance hero card",
    bullets = [
        "##Editorial greeting",
        "  — Date line in small caps with a horizontal rule: 'Monday, 23 June'.",
        "  — Large serif headline (Fraunces typeface): 'Good morning, James.' (or Treasurer",
        "    when the entity type is a charity, club, or society).",
        "  — Context sentence: 'Everything's running normally today. 3 things need your",
        "    signature — they're below.'",
        "",
        "##Total Balance hero card",
        "  — Dark (stone-900 #1c1917) rounded card with a Santander red (#c8102e) glowing",
        "    orb in the top-right corner and a subtle grain texture overlay.",
        "  — Displays: 'Total balance' label, entity name, account count, sort code.",
        "  — Total across all accounts for the current entity. For a limited company:",
        "    Operating (£284,750) + Reserve (£1,240,500) + Payroll (£195,830) = £1,721,080.",
        "  — Three badge pills: entity type label, mandate rule, FSCS Protected.",
        "  — If credit ring-fencing is active: a fourth 'Credit ring-fenced' emerald badge.",
        "  — The balance figure uses tabular figures (num-tab class) for digit alignment.",
    ])

# ── SLIDE 13: Proactive Forecast strip ────────────────────────────────────────
content_slide(prs,
    section_label = "C  ·  Command Centre",
    title = "Proactive Forecast strip — Predict, Explain, Act",
    bullets = [
        "Three horizontally scrolling nudge cards. Design principle: Predict → Explain → Act.",
        "Each card names the event, tells the user why it matters, and links to the exact workflow.",
        "",
        "##Card 1 — Payroll run (amber border)",
        "  — 'In 3 days'  ·  'Payroll run · £42,180'",
        "  — 'Balance covers it — review if payees changed'",
        "  — Taps directly into the Bulk Payments (Wages) workflow.",
        "",
        "##Card 2 — VAT return due (red border)",
        "  — '22 days'  ·  'VAT return due 7 Aug'",
        "  — 'Q3 Making Tax Digital · HMRC submission'",
        "  — Taps directly into the MTD/Tax tab.",
        "",
        "##Card 3 — Companies House (blue border)",
        "  — '14 days'  ·  'Annual confirmation due'",
        "  — 'Companies House · Whitfield Holdings Ltd'",
        "  — Taps directly into the KYB / ID register workflow.",
        "",
        "##Why this matters",
        "  — No hunting through menus. The app surfaces what needs doing without being asked.",
        "  — Colour-coding (amber / red / blue) communicates urgency at a glance.",
        "  — Each card is fully interactive — tapping starts the right workflow immediately.",
    ])

# ── SLIDE 14: Pre-approved lending ────────────────────────────────────────────
content_slide(prs,
    section_label = "C  ·  Command Centre",
    title = "Pre-approved lending offer — contextual, not interruptive",
    bullets = [
        "##The offer card",
        "  — Dark hero card (same style as the Total Balance card) with subtle circular overlays.",
        "  — Badge: 'Pre-approved offer · valid until 30 Sep 2026'.",
        "  — Large figure: £45,000.",
        "  — 'From 6.2% APR · Based on your 6-month trading history'.",
        "  — 'View offer & draw down' — one tap opens the lending workflow.",
        "",
        "##Inside the lending workflow (4 steps)",
        "  — Step 0: Offer details, rate, representative example.",
        "  — Step 1: Choose repayment term (12, 24, or 36 months). Monthly cost updates live.",
        "  — Step 2: Confirm purpose of funds.",
        "  — Step 3: Sign & draw down. OTP confirmation.",
        "",
        "##After acceptance",
        "  — The offer card is replaced by a 'Business loan active · £45,000 drawn' confirmation",
        "    in emerald green. First repayment date shown. Repayment term badge.",
        "  — State persists for the session — switching entity type or tabs doesn't reset it.",
        "",
        "##Design principle",
        "  — The offer appears in context, not as a pop-up or interstitial.",
        "  — Based on demonstrated trading data — not a cold credit offer.",
        "  — One tap to accept, not a branch visit or paper form.",
    ])

# ── SLIDE 15: Health Score + Cooling-off ──────────────────────────────────────
content_slide(prs,
    section_label = "C  ·  Command Centre",
    title = "Business Health Score + Cooling-off cards",
    bullets = [
        "##Business Health Score",
        "  — A composite score card covering: payment history, signatory compliance,",
        "    MTD submission status, document freshness, and account activity.",
        "  — Visual gauge or score band. Tappable for a breakdown.",
        "  — Helps a director understand their overall compliance health at a glance.",
        "",
        "##Cooling-off cards",
        "  — Regulatory requirement: high-risk changes (closures, mandate changes) have a",
        "    mandatory 24-hour delay before execution.",
        "  — Pauses on weekends and bank holidays (working days only).",
        "  — Each card shows: request type, description, execution time in plain English",
        "    ('We'll do this tomorrow afternoon').",
        "  — Live amber progress bar — width updates every 30 seconds via a tick state.",
        "  — Two buttons: Cancel (with confirmation sheet) and 'Do it now' (bypasses delay).",
        "",
        "##Demo cooling-off",
        "  — Tap the 'Demo cooling-off' pill to inject a mock cooling-off card instantly.",
        "  — The progress bar starts at 0% and grows in real time.",
        "  — 'Do it now' fires a toast and dismisses the card immediately.",
        "  — Cancel triggers a bottom sheet with a final confirmation step.",
    ])

# ── SLIDE 16: Stalled requests ────────────────────────────────────────────────
content_slide(prs,
    section_label = "C  ·  Command Centre",
    title = "Stalled requests — 'Priya's on it'",
    bullets = [
        "##The scenario",
        "  — A dual-authorisation request was sent to a co-signer (e.g. Sarah Chen).",
        "  — The co-signer did not approve within the required window.",
        "  — Normally, the request would silently expire — causing frustration and re-work.",
        "",
        "##What the prototype does instead",
        "  — The app detects the missed window automatically.",
        "  — A 'Priya's on it' card appears on the Command Centre (blue, escalated styling).",
        "  — Message: 'A signature didn't arrive in time. Priya saw this within 2 minutes",
        "    and will reach out today. Nothing's been lost — just paused.'",
        "  — Two buttons: 'Got it' (dismisses the card) and 'Call Priya' (opens the RM sheet).",
        "",
        "##Why this matters",
        "  — Stalled requests are a real pain point in business banking.",
        "  — The current experience: the request expires, nobody knows why, the customer",
        "    restarts the entire process.",
        "  — This design: the bank takes ownership, the RM is already involved before",
        "    the customer has to call.",
        "",
        "##Demo",
        "  — Tap 'Simulate timeout' on the Home screen to inject a stalled card.",
    ])

# ── SLIDE 17: MTD teaser + Paperless banner ───────────────────────────────────
content_slide(prs,
    section_label = "C  ·  Command Centre",
    title = "MTD VAT teaser + Paperless banner",
    bullets = [
        "##MTD VAT teaser card",
        "  — Appears when there is an open VAT obligation (Q3 Jul–Sep 2026).",
        "  — Santander red left accent bar and subtle red gradient background.",
        "  — Shows: VAT amount owed (£18,472.40), days to deadline, transaction count to review.",
        "  — 'VAT due in 22 days · MTD' — taps directly into the Tax tab.",
        "  — VRN: GB 287 4513 92  ·  HMRC MTD VAT API v1.0  ·  247 transactions.",
        "",
        "##Paperless banner",
        "  — Emerald green card. Icon: envelope with an X (no more paper post).",
        "  — '3 paper forms retired'",
        "  — 'No more posting to Sunderland · 5 days → minutes'",
        "  — 'Sunderland' is a deliberate real-world reference — Santander's UK head office",
        "    used to require physical form ANB9 0370 to be posted there for account closures.",
        "  — Taps a modal showing the three retired forms and estimated time saved.",
        "",
        "##Why these two cards sit together",
        "  — They represent the same message: the bank now does in minutes what used to",
        "    take days of paperwork and postal delays.",
        "  — Together they make a compelling narrative for both customers and internal teams.",
    ])

# ── SLIDE 18: Pending approvals + Actions grid ────────────────────────────────
content_slide(prs,
    section_label = "C  ·  Command Centre",
    title = "Pending approvals + Paperless actions grid",
    bullets = [
        "##Pending approvals",
        "  — 'Action required' section with a red count badge.",
        "  — Shows the first 2 of 3 pending items (to avoid overwhelming the screen):",
        "    1. Bulk Payment — October Wages · 47 employees · £183,450.00 · expires 46h",
        "    2. Mandate Change — Add Mark Patel as signatory · expires 4d",
        "  — Each card shows: type, description, amount, initiator, expiry. Taps Approve tab.",
        "",
        "##Paperless actions grid (2 × 5 tiles)",
        "  — 'No more posting to Sunderland · Paperless actions'",
        "  — 10 tiles, each opens a full-screen workflow:",
        "  •  Bulk payments   — CSV · BACS, FP, CHAPS  (highlighted tile — primary CTA)",
        "  •  International   — FX · SWIFT · SEPA",
        "  •  Change mandate  — Add, remove, signing rule",
        "  •  Business details — Name, address, contact",
        "  •  Scan receipt    — Auto-categorise for MTD",
        "  •  ID register     — Lists 1, 2 & 3 (KYC/KYB)",
        "  •  Dormant accounts — Reactivate or close  (badge: 1)",
        "  •  Close account   — Form ANB9 0370",
        "  •  Voice memo      — Speak → expense auto-tagged",
        "  •  Optimise payments — 30-day payment sequencer",
    ])

# ── SLIDE 19: Accounts list ────────────────────────────────────────────────────
content_slide(prs,
    section_label = "C  ·  Command Centre",
    title = "Accounts list — balances and mandate badges",
    bullets = [
        "##The four accounts (Limited Company default)",
        "  1. Operating       · 09-01-29 · ····2841 · £284,750.32   · Any 1 signature",
        "  2. Reserve         · 09-01-29 · ····9012 · £1,240,500.00 · Any 2 of 3",
        "  3. Payroll         · 09-01-29 · ····6633 · £195,830.50   · Any 2 of 3",
        "  4. Legacy Trading  · 09-01-29 · ····4422 · £0.00         · All · Dormant",
        "",
        "##Mandate rule badges",
        "  — 'Any 1 signature' (green badge) — single signatory sufficient",
        "  — 'Any 2 of 3'      (stone badge) — dual authorisation required",
        "  — 'All 3'           (dark badge)  — all verified signatories must sign",
        "  — 'Dormant'         (amber badge) — account inactive, needs reactivation or closure",
        "",
        "##Entity switching changes the accounts",
        "  — Sole Trader: 2 accounts (Trading, Tax Reserve) — all 'Any 1 signature'",
        "  — Partnership: 3 accounts (Operating, Reserve, Drawing) — 'All' mandate on Operating",
        "  — Limited / LLP: 4 accounts as above",
        "  — Charities / Clubs / Societies: same 4-account structure as limited",
        "",
        "##Design detail",
        "  — Account numbers are masked (····XXXX). Sort codes are shown in full.",
        "  — All monetary amounts use tabular figures (font-variant-numeric: tabular-nums)",
        "    so digits align in columns — a detail that matters in financial UI.",
    ])

# ── SLIDE 20: Forecast chart + Director Centre ─────────────────────────────────
content_slide(prs,
    section_label = "C  ·  Command Centre",
    title = "13-week forecast chart + Director Centre",
    bullets = [
        "##13-week cash flow forecast chart",
        "  — Gradient-filled SVG area chart. Smooth Catmull-Rom curve (converted to",
        "    cubic Bézier control points) — not a jagged step chart.",
        "  — Y-axis zoomed to the actual balance band (min - 35% padding to max + 35%)",
        "    so weekly movements are visible even when the balance is large.",
        "  — Event markers on the curve: red dots for outflows, pink for receipts,",
        "    amber circle for the lowest balance week.",
        "  — Dashed amber 'floor warning' line at £80,000 working capital threshold.",
        "  — Below the chart: two summary tiles (Lowest week · Week 13 close).",
        "  — Named largest outflow: e.g. 'HMRC VAT £41,614 in week 3 (7 Jul)'.",
        "  — Amber warning if balance dips below the £80k floor.",
        "",
        "##Director Centre",
        "  — Below the forecast chart, a section labelled 'Director Centre'.",
        "  — Supplier Radar: lists 5 key suppliers with Companies House filing status.",
        "    Meridian Logistics (261 days overdue — red risk), CloudStack (95 days — amber),",
        "    Apex Print, TechLink, Office Base (all current — green).",
        "  — Scheduled payments: PAYE/NIC £29,650 (24 Jun), Meridian £8,750 (26 Jun),",
        "    CloudStack £2,240 (27 Jun), Office Lease £6,500 (1 Jul), HMRC VAT £41,614 (7 Jul).",
        "  — Cards section: Business Debit (····4821) and Payroll Card (····3927) — freeze / View PIN.",
        "  — Privacy controls: personal accounts linked / separated, open banking consents,",
        "    credit decisioning ring-fence.",
    ])

# ── SLIDE 21: Section D divider ────────────────────────────────────────────────
section_divider(prs,
    section_label = "Section D  ·  Slides 21–26",
    section_title = "Workflow\nWalkthroughs",
    section_sub   = "All 11 paperless workflows — step-by-step, full-screen overlays")

# ── SLIDE 22: Workflow overview ────────────────────────────────────────────────
content_slide(prs,
    section_label = "D  ·  Workflow Walkthroughs",
    title = "How workflows work — 11 paperless journeys",
    bullets = [
        "##Mechanics",
        "  — Tapping an action tile or forecast card opens a full-screen workflow overlay.",
        "  — Each workflow is step-based (0-indexed). Progress dots across the top.",
        "  — Back button: goes to previous step. At step 0, closes the workflow entirely.",
        "  — 'Replaces' footer: shows which paper form or manual process this replaces.",
        "",
        "##The 11 workflows",
        "  1.  Account closure        — Form ANB9 0370 · 4 steps · 24h cooling-off",
        "  2.  Business details (Biz) — Name, address, authorised contact · 4 steps",
        "  3.  Mandate change         — Add/remove signatories, change signing rule · 5 steps",
        "  4.  Bulk payments (Wages)  — Persistent payee book, CSV upload, BACS · 4 steps",
        "  5.  Lending draw-down      — Pre-approved offer, term selector · 4 steps",
        "  6.  FX / International     — SWIFT/SEPA, rate lock, counterparty · 4 steps",
        "  7.  Dormancy management    — Reactivate or close dormant accounts · 3 steps",
        "  8.  Account unlinking      — Separate personal from business data · 3 steps",
        "  9.  Credit ring-fence      — Exclude personal data from credit decisions · 3 steps",
        " 10.  ID register (KYB)      — KYC lists 1, 2 & 3 per signatory · 4 steps",
        " 11.  MTD VAT submission     — Transaction review, categorise, declare · 4 steps",
        "",
        "##Mandate logic",
        "  — getMandateFor() picks the strictest rule across selected accounts.",
        "  — Determines whether a second signature is required and what cooling-off applies.",
    ])

# ── SLIDE 23: Mandate change workflow ─────────────────────────────────────────
content_slide(prs,
    section_label = "D  ·  Workflow Walkthroughs",
    title = "Mandate change workflow",
    bullets = [
        "##What it covers",
        "  — Add a new signatory, remove an existing one, or change the signing rule",
        "    (Any 1 / Any 2 / All).",
        "  — Applies to one or more selected accounts.",
        "",
        "##Current signatories (Limited Company)",
        "  — James Whitfield — Director — verified (passport, council tax, residential match)",
        "  — Sarah Chen — Finance Director — verified (UK photocard, bank statement, lease)",
        "  — Tom Bridges — Operations Lead — under review (documents from Aug 2024 ⚠)",
        "  — Anita Roy — Director — verified (passport, HMRC letter, utility bill)",
        "",
        "##Steps",
        "  1.  Select accounts to modify",
        "  2.  Choose change type: add signatory / remove / change rule",
        "  3.  Configure the change (select person, new rule)",
        "  4.  Review mandate impact — what changes, which accounts are affected",
        "  5.  OTP confirmation → cooling-off (24h) for high-risk changes",
        "",
        "##Mandate rules by account",
        "  — Strictest rule across selected accounts determines overall requirement.",
        "  — Removing a signatory from a dual-auth account triggers a 24h cooling-off.",
        "  — Adding a signatory on a single-auth account: immediate with OTP only.",
        "",
        "##Replaces",
        "  — Manual paper form  ·  'Post to Sunderland'  ·  5-7 working days",
    ])

# ── SLIDE 24: Bulk payments + MTD ─────────────────────────────────────────────
two_col_slide(prs,
    section_label = "D  ·  Workflow Walkthroughs",
    title = "Bulk payments (wages) + MTD VAT submission",
    left_bullets = [
        "##Bulk Payments — Wages",
        "",
        "• Persistent payee book (10 employees pre-loaded):",
        "  — Aisha Khan · Sr Engineer · £3,200",
        "  — Ben Carter · Ops Lead · £2,850",
        "  — Carmen Diaz · Finance Mgr · £4,100",
        "  — David Olu · Designer · £3,950",
        "  — Elena Petrov · Product Lead · £5,200",
        "  — + 5 more employees",
        "",
        "• Steps:",
        "  1. Select payees (search, add new, deselect)",
        "  2. Review total and CoP (Confirmation of Payee) statuses",
        "  3. Schedule: now / future date",
        "  4. Mandate check + sign",
        "",
        "• Supports CSV upload for one-off bulk imports.",
        "• CoP verified badge per payee (green tick).",
        "• Replaces: paper BACS instruction forms.",
    ],
    right_bullets = [
        "##MTD / VAT Submission",
        "",
        "• Q3 Jul–Sep 2026 — open obligation:",
        "  — £18,472.40 VAT due to HMRC",
        "  — 247 transactions to review",
        "  — VRN: GB 287 4513 92",
        "  — Deadline: 7 Nov 2026",
        "",
        "• Steps:",
        "  1. Review transactions — categorise each",
        "     as business / personal / exempt",
        "  2. Auto-categorised at high/medium/low",
        "     confidence — low-confidence flagged",
        "  3. Confirm VAT totals — editable",
        "  4. Legal declaration + countdown timer",
        "     before submission to HMRC API",
        "",
        "• Receipt scan: camera → OCR → auto-tag",
        "• Voice memo: speak → expense auto-tagged",
        "  (merchant, amount, category, VAT rate)",
        "",
        "• Replaces: Bridging software, spreadsheets,",
        "  manual HMRC portal entry.",
    ])

# ── SLIDE 25: FX + Account closure ────────────────────────────────────────────
two_col_slide(prs,
    section_label = "D  ·  Workflow Walkthroughs",
    title = "International payments (FX) + Account closure",
    left_bullets = [
        "##FX / International Payments",
        "",
        "• Currencies: EUR, USD, CHF, AUD, CAD",
        "• Live indicative rates (prototype data):",
        "  — EUR: 1.1742  ·  fee £0.50",
        "  — USD: 1.2681  ·  fee £0.50",
        "  — CHF: 1.1234  ·  fee £0.75",
        "",
        "• Steps:",
        "  1. Amount + currency selection",
        "  2. Beneficiary IBAN + name",
        "  3. Rate confirmation (rate locked",
        "     for 30 seconds — countdown shown)",
        "  4. Reference + sign",
        "",
        "• SWIFT and SEPA payment rails.",
        "• Counterparty management — saved",
        "  beneficiaries with full audit trail.",
        "",
        "• Replaces: branch foreign payments form,",
        "  phone-based FX desk instructions.",
    ],
    right_bullets = [
        "##Account Closure",
        "",
        "• Replaces Form ANB9 0370.",
        "  'No more posting to Sunderland.'",
        "",
        "• Steps (standard path):",
        "  1. Select accounts to close",
        "  2. Destination for remaining funds",
        "  3. In-credit check (must be in credit)",
        "  4. Sign per mandate → 24h cooling-off",
        "",
        "• Escalation path (Partnerships):",
        "  — A partner is unreachable / deceased /",
        "    health issue / dispute.",
        "  — 4-step regulatory escalation:",
        "    reason → contact evidence →",
        "    restrict & refer to Specialist",
        "    Vulnerable Customer Team.",
        "  — FCA Consumer Duty PS22/9.",
        "  — BCOBS due diligence requirements.",
        "",
        "• 24h cooling-off: execution delayed,",
        "  pausable, cancellable, or skippable",
        "  ('Do it now' for urgent closures).",
    ])

# ── SLIDE 26: Approve screen ───────────────────────────────────────────────────
content_slide(prs,
    section_label = "D  ·  Workflow Walkthroughs",
    title = "Approve screen — pending signature queue",
    bullets = [
        "##Three demo items waiting for the director's signature",
        "  1. Bulk Payment — October Wages · 47 employees · £183,450.00",
        "     — Initiated by Sarah Chen  ·  Expires in 46 hours",
        "  2. Mandate Change — Add Mark Patel as signatory",
        "     — Initiated by Sarah Chen  ·  Expires in 4 days",
        "  3. Account Closure — Legacy Trading · ····4422",
        "     — Initiated by James Whitfield  ·  Expires in 11 days",
        "",
        "##Signing flow",
        "  — Tap any pending item to open its detail view.",
        "  — Review the full request: what, who initiated it, when it expires, what it costs.",
        "  — Biometric or OTP confirmation to approve or reject.",
        "  — Approved items move to the Audit trail immediately.",
        "  — Rejected items fire a toast and notify the initiator.",
        "",
        "##SCA (Strong Customer Authentication)",
        "  — OTP Sheet: 6-digit code, 30-second resend timer.",
        "  — Voice ID: enrol 3 phrases ('My voice is my passport' etc.)",
        "  — Device biometric: Face ID / Fingerprint / Face (Android) — switchable.",
        "  — SCA step-up matrix: 6 tiers, PSD2 RTS Art.97 compliant.",
    ])

# ── SLIDE 27: Audit + Statements ──────────────────────────────────────────────
content_slide(prs,
    section_label = "D  ·  Workflow Walkthroughs",
    title = "Audit trail + Statements",
    bullets = [
        "##Audit trail",
        "  — Immutable, chronological log of all events in the prototype.",
        "  — Every workflow completion, approval, rejection, and state change is recorded.",
        "  — Shows: timestamp, action type, actor (director name), affected entity,",
        "    result, and any associated reference numbers.",
        "  — Export options: CSV download, PDF statement, share link.",
        "  — Designed for: compliance teams, external auditors, HMRC enquiries.",
        "",
        "##Statements (Statements tab)",
        "  — Full transaction history with rich filtering:",
        "    Date range selector, category (Equipment / Travel / Utilities / Sales…),",
        "    counterparty name search, payment method (Card / DD / SO / FP / BACS).",
        "  — Two views: chronological and by category.",
        "  — Counterparty deep-dive: tap any merchant to see all transactions with them,",
        "    total spend, and average transaction value.",
        "  — Supplier Radar cross-reference: Meridian Logistics flagged as high-risk",
        "    (261 days overdue at Companies House, £34,700 spend).",
        "",
        "##Key transactions in the demo dataset (Sep 2026)",
        "  — Whitfield invoice #4827: +£12,500 (sales) · Office Tech Solutions: -£2,400",
        "  — BT Business Broadband: -£89.99 · M&S groceries: -£67.42 (personal, excluded)",
        "  — Premier Inn (client meeting): -£154.00 (low confidence — needs categorisation)",
    ])

# ── SLIDE 28: Section E divider ────────────────────────────────────────────────
section_divider(prs,
    section_label = "Section E  ·  Slides 28–30",
    section_title = "Presenter Notes",
    section_sub   = "Demo script, entity switching guide, back cover")

# ── SLIDE 29: Suggested demo script ───────────────────────────────────────────
content_slide(prs,
    section_label = "E  ·  Presenter Notes",
    title = "Suggested demo script — Limited Company walkthrough",
    bullets = [
        "##Opening (2 min)",
        "  1. Open https://imscots-ui.github.io/santander on a phone or DevTools emulation.",
        "  2. Entity should read 'Limited Company'. Confirm with the pill top-left.",
        "  3. Scroll the Command Centre slowly — name each section as you pass it.",
        "  4. Highlight: 'One screen. Everything Whitfield Holdings needs to know today.'",
        "",
        "##Forecast → Payroll (3 min)",
        "  5. Tap the amber 'Payroll run' forecast card. Walk through the wages workflow.",
        "  6. Show the payee book. Deselect one employee. Show the total update.",
        "  7. Proceed to schedule → mandate check → sign (OTP sheet).",
        "  8. Return to home. Show the 24h cooling-off card if mandate triggered it.",
        "",
        "##MTD VAT (3 min)",
        "  9. Tap the red 'VAT return' forecast card. Switch to the Tax tab.",
        " 10. Show the Q3 obligation card — £18,472 due, 22 days, 247 transactions.",
        " 11. Enter the submission workflow. Review 2–3 transactions, categorise one.",
        " 12. Show the countdown timer before HMRC submission fires.",
        "",
        "##Stalled co-signer (2 min)",
        " 13. Return to Home. Tap 'Simulate timeout' — show the 'Priya's on it' card.",
        " 14. Explain: 'In the current world this would silently expire.'",
        " 15. Tap 'Call Priya' to show the RM sheet.",
        "",
        "##Close (1 min)",
        " 16. Show the 'Paperless banner' — 3 forms retired, Sunderland reference.",
        " 17. Toggle to Desktop mode. Show same state, wider layout.",
    ])

# ── SLIDE 30: Entity switching guide ──────────────────────────────────────────
content_slide(prs,
    section_label = "E  ·  Presenter Notes",
    title = "Entity switching — live demo of all 7 types",
    bullets = [
        "The entity switcher is a key differentiator. It shows that the UI adapts",
        "dynamically to the legal structure of the business — not a one-size-fits-all form.",
        "",
        "##How to switch",
        "  — Tap the pill labelled with the current entity name (e.g. 'Limited Company ›').",
        "  — A bottom sheet opens with 7 options. Tap any to switch immediately.",
        "",
        "##What changes per entity",
        "  — Company name: 'Whitfield Holdings Ltd' → 'Whitfield & Co Partnership' etc.",
        "  — Greeting: 'James' → 'Treasurer' (for charities, clubs, societies)",
        "  — Signatory labels: 'director' → 'partner' / 'trustee' / 'member'",
        "  — Account list: sole trader has 2 accounts; limited has 4; partnership has 3",
        "  — Mandate rules: sole trader is always 'Any 1'; partnership defaults to 'All'",
        "  — Companies House number shown for: Limited (07429183), LLP (OC428192)",
        "  — Charity Commission number shown for: Registered Charity (1192847)",
        "  — Minutes required: Charity, Club, Society — extra document upload step",
        "",
        "##The 7 entity types",
        "  1. Limited Company  — the default view in this walkthrough",
        "  2. Sole Trader      — simplest: one person, single-auth, 2 accounts",
        "  3. Partnership      — all-party mandate by default, escalation for unreachable partner",
        "  4. LLP              — like limited but member labels, OC- Companies House number",
        "  5. Registered Charity — trustee labels, Charity Commission number, minutes required",
        "  6. Club             — committee member labels, no CH/CC number",
        "  7. Society          — same as club, different name prefix",
    ])

# ── SLIDE 31: Back cover ───────────────────────────────────────────────────────
back_cover(prs,
    contact   = "Created by Alan Davidson  ·  Alan.Davidson@santander.co.uk  ·  May 2026",
    disclaimer= (
        "Prototype only. No backend systems, no live data, no real customer information. "
        "All names, account numbers, sort codes, VAT registration numbers, and financial figures "
        "are fabricated for demonstration purposes. This document and the associated prototype "
        "are provided for internal discussion only and do not represent Santander's committed "
        "product roadmap or any regulatory commitment."
    ))


# ── Save ───────────────────────────────────────────────────────────────────────
output_path = "/home/user/santander/Santander_LimitedCompany_Walkthrough.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
