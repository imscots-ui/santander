"""
build_changes_slides.py
Appends slides 12-15 (Session 2 — 23 June 2026) to Santander_June2026_Changes.pptx.

Style mirrors the existing deck:
  - Red thin stripe across top (#DA291C)
  - Warm off-white content area (#FAF6EF)
  - Title in dark charcoal 1A1A1A, 26 pt bold
  - Section labels in #DA291C, 14 pt bold
  - Body bullets in 1A1A1A, 12 pt
  - Footer strip in #F5F0E8 with muted code/location note
  - Slide number + eyebrow in top-right / top-left

Run:
    python3 build_changes_slides.py
"""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
RED      = RGBColor(0xDA, 0x29, 0x1C)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xFA, 0xF6, 0xEF)
DARK     = RGBColor(0x1A, 0x1A, 0x1A)
MUTED    = RGBColor(0x78, 0x71, 0x6C)
FOOTERBG = RGBColor(0xF5, 0xF0, 0xE8)

# ── EMU constants (matching existing slide geometry) ─────────────────────────
W        = 12188952   # slide width
H        = 6858000    # slide height
STRIPE_H = 73152      # red top stripe height
MARGIN_L = 457200     # left margin
MARGIN_R = 457200     # right margin
BODY_W   = W - MARGIN_L - MARGIN_R   # 11274552
TEXT_W   = W - MARGIN_L - MARGIN_R - 182880  # ~10972800 (matching existing)
INDENT_W = W - MARGIN_L - MARGIN_R - 182880 - 182880  # ~10515600 (bullet indent)

FOOTER_TOP = 6217920
FOOTER_H   = 457200

# Font sizes in EMU (914400 EMU = 1 pt ... actually pptx Pt() is helper)
# pptx stores font size in 100ths of a point (EMU-like): Pt(26) = 330200
# Verify: Pt(1) = 12700 half-points? No — pptx uses 1/100 pt = 12700 EMU per pt
# python-pptx: font.size is in EMUs where 1pt = 12700


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_rgb):
    """Add a filled rectangle with no line."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1 → actually add_shape uses MSO_AUTO_SHAPE_TYPE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()   # no border
    return shape


def add_textbox(slide, left, top, width, height, text, font_size_pt,
                bold=False, color=DARK, align=PP_ALIGN.LEFT, wrap=True):
    """Add a textbox with a single paragraph/run."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    # Remove default padding on textbox
    tf.margin_left   = 0
    tf.margin_right  = 0
    tf.margin_top    = 0
    tf.margin_bottom = 0
    return txb


def add_bullets(slide, left, top, width, bullets, font_size_pt=12, color=DARK, line_gap_emu=292608):
    """
    Add multiple bullet textboxes stacked vertically.
    Returns the bottom Y position after the last bullet.
    """
    y = top
    for bullet in bullets:
        h = line_gap_emu
        add_textbox(slide, left, y, width, h, bullet, font_size_pt, color=color)
        y += h
    return y


def build_slide(prs, slide_num, total_slides, title, session_label,
                sections, footer_note):
    """
    Build a single content slide.

    sections: list of dicts:
        { 'heading': str, 'bullets': [str, ...] }
    """
    layout = prs.slide_layouts[6]   # Blank
    slide  = prs.slides.add_slide(layout)

    # ── Red top stripe ────────────────────────────────────────────────────────
    add_rect(slide, 0, 0, W, STRIPE_H, RED)

    # ── Main background ───────────────────────────────────────────────────────
    add_rect(slide, 0, STRIPE_H, W, H - STRIPE_H, OFFWHITE)

    # ── Session eyebrow (top-left, red) ───────────────────────────────────────
    add_textbox(slide,
                MARGIN_L, STRIPE_H + 18000, 4000000, 320000,
                session_label, 9, bold=True, color=RED)

    # ── Slide number (top-right, muted) ───────────────────────────────────────
    num_label = f"{slide_num} / {total_slides}"
    add_textbox(slide,
                W - MARGIN_R - 800000, STRIPE_H + 18000, 760000, 320000,
                num_label, 9, color=MUTED, align=PP_ALIGN.RIGHT)

    # ── Title ─────────────────────────────────────────────────────────────────
    add_textbox(slide,
                MARGIN_L, 228600, TEXT_W, 640080,
                title, 26, bold=True, color=DARK)

    # ── Red rule under title ───────────────────────────────────────────────────
    add_rect(slide, MARGIN_L, 868680, BODY_W, 36576, RED)

    # ── Sections ─────────────────────────────────────────────────────────────
    y = 1097280
    HEADING_H    = 365760
    BULLET_H     = 285750   # slightly tighter than existing (292608) to fit more
    SECTION_GAP  = 73152    # extra space before next heading

    for section in sections:
        # Heading label
        add_textbox(slide, MARGIN_L, y, TEXT_W, HEADING_H,
                    section['heading'], 14, bold=True, color=RED)
        y += HEADING_H

        # Bullets
        indent_left = MARGIN_L + 182880
        for bullet in section['bullets']:
            add_textbox(slide, indent_left, y, INDENT_W, BULLET_H,
                        bullet, 12, color=DARK)
            y += BULLET_H

        y += SECTION_GAP

    # ── Footer bar ────────────────────────────────────────────────────────────
    add_rect(slide, MARGIN_L, FOOTER_TOP, BODY_W, FOOTER_H, FOOTERBG)
    add_textbox(slide, MARGIN_L + 182880, FOOTER_TOP + 45720,
                6400000, 365760,
                footer_note, 10, color=MUTED)
    # Classification + credit (right side of footer bar)
    add_textbox(slide, MARGIN_L + BODY_W - 4500000, FOOTER_TOP + 45720,
                4318000, 365760,
                "INTERNAL · CONFIDENTIAL   ·   A. Davidson · June 2026",
                9, bold=True, color=RED, align=PP_ALIGN.RIGHT)

    return slide


# ── Slide definitions ─────────────────────────────────────────────────────────

EYEBROW = "SESSION 2  —  23 June 2026"
TOTAL   = 15

SLIDES = [
    {
        "num":   12,
        "title": "Forecast Chart — Smooth Gradient Area Rebuild",
        "sections": [
            {
                "heading": "Before",
                "bullets": [
                    "▪  Flat coloured bars with no visual hierarchy — movement looked identical week-to-week",
                    "▪  Chart axis started at zero, compressing weekly swings to near-invisible deltas",
                    "▪  No event markers; no distinction between inflow weeks and outflow weeks",
                    "▪  Amber floor warning (£80k) never triggered because the scale hid dips",
                ]
            },
            {
                "heading": "After",
                "bullets": [
                    "▪  Smooth area chart using Catmull-Rom → cubic Bézier conversion for a natural curve",
                    "▪  Y-axis zoomed to actual [min, max] balance band (+35% padding) — weekly movement is now legible",
                    "▪  Event markers: red dot = largest outflow, rose dot = inflow, amber ring = lowest week",
                    "▪  Event text block names the dominant scheduled outflow (e.g. 'VAT direct debit £41,614')",
                    "▪  Amber warning fires when balance dips below the £80k floor",
                ]
            },
            {
                "heading": "Why it matters",
                "bullets": [
                    "▪  A cash-flow chart that looks flat communicates nothing. The rebuild makes variance immediately visible",
                    "▪  and gives the SME owner one named event to act on — the core job of a forecast widget.",
                ]
            },
        ],
        "footer": "Session 2 · 23 Jun 2026  ·  forecast useMemo ~line 609, render block ~line 3207"
    },
    {
        "num":   13,
        "title": "Card Buttons — Visual Hierarchy Polish",
        "sections": [
            {
                "heading": "Before",
                "bullets": [
                    "▪  Two ghost buttons (\"View PIN\" and \"Freeze\") with identical weight and styling",
                    "▪  No visual signal to distinguish a read-only action from a destructive one",
                    "▪  \"Unfreeze\" used the same ghost style as the initial Freeze button",
                ]
            },
            {
                "heading": "After",
                "bullets": [
                    "▪  \"View PIN\" — visible ghost: stone-50 background, stone-300 border, stone-800 text",
                    "▪  \"Freeze\" — solid stone-900 fill (filled black), communicating a serious / committed action",
                    "▪  \"Unfreeze\" (frozen state) — solid blue-600 fill, clearly signals a reversible recovery action",
                ]
            },
            {
                "heading": "Why it matters",
                "bullets": [
                    "▪  Identical button weight is one of Refactoring UI's canonical anti-patterns. The new hierarchy",
                    "▪  signals action severity at a glance and reduces accidental freezes in user testing.",
                ]
            },
        ],
        "footer": "Session 2 · 23 Jun 2026  ·  Card buttons render block ~line 2600 (SavingsSheet / card section)"
    },
    {
        "num":   14,
        "title": "Biometric Picker — Segmented Control Pattern",
        "sections": [
            {
                "heading": "Before",
                "bullets": [
                    "▪  Flat colour tiles for Face ID / Touch ID / Voice ID with no clear selected-state affordance",
                    "▪  All three options looked identical whether selected or not",
                    "▪  No depth cue to indicate which option was active",
                ]
            },
            {
                "heading": "After",
                "bullets": [
                    "▪  Grey track (stone-100 background) containing three equal-width segments",
                    "▪  Active state: the selected segment lifts as a white card with drop shadow + stone-200 border",
                    "▪  Inactive segments sit flush in the track — classic iOS-style segmented control pattern",
                    "▪  Three options preserved: Face ID  /  Touch ID  /  Voice ID",
                ]
            },
            {
                "heading": "Why it matters",
                "bullets": [
                    "▪  The segmented control is a well-established mobile pattern for mutually exclusive choices.",
                    "▪  The lifted-card active state is immediately legible without colour dependency (accessibility win).",
                ]
            },
        ],
        "footer": "Session 2 · 23 Jun 2026  ·  Biometric picker ~line 2390 (renderIdCheck / VoiceIdSheet)"
    },
    {
        "num":   15,
        "title": "VAT Countdown Overlay — Calm Redesign",
        "sections": [
            {
                "heading": "Before",
                "bullets": [
                    "▪  Header read 'One last chance to stop' — alarming, adversarial framing",
                    "▪  Cancel button used brand red (#DA291C) — visually screamed danger for a routine cancel",
                    "▪  Countdown ring radius r=42 with no unit label beneath the seconds number",
                    "▪  HMRC warning was verbose and lacked visual containment",
                ]
            },
            {
                "heading": "After",
                "bullets": [
                    "▪  Header changed to 'Filing your VAT return' — confident, process-oriented framing",
                    "▪  Cancel button: calm stone-100 background / stone-700 text — clearly secondary, non-alarming",
                    "▪  Ring enlarged from r=42 to r=46; 'sec' label added beneath the countdown number",
                    "▪  Tax amount + label moved into a left-aligned stone-50 detail pill for clarity",
                    "▪  HMRC amber warning shortened and given a visible border for containment",
                    "▪  Footer line: 'Filing automatically in Xs' replaces the redundant descriptive text",
                ]
            },
            {
                "heading": "Why it matters",
                "bullets": [
                    "▪  Alarm framing on a confirmation screen increases task abandonment. Calm copy + a neutral cancel",
                    "▪  button reduce anxiety without hiding the escape hatch — a standard reassurance pattern.",
                ]
            },
        ],
        "footer": "Session 2 · 23 Jun 2026  ·  VAT countdown overlay ~line 4998 (renderMtdSubmit)"
    },
]


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    path  = "/home/user/santander/Santander_June2026_Changes.pptx"
    base  = "/tmp/changes_base11.pptx"   # clean 11-slide template from git
    import shutil, os
    # Always start from the clean 11-slide base so this script is idempotent
    shutil.copy2(base, path)
    prs = Presentation(path)

    before = len(prs.slides)
    print(f"Base slides: {before}")

    for slide_def in SLIDES:
        build_slide(
            prs,
            slide_num     = slide_def["num"],
            total_slides  = TOTAL,
            title         = slide_def["title"],
            session_label = EYEBROW,
            sections      = slide_def["sections"],
            footer_note   = slide_def["footer"],
        )
        print(f"  Built slide {slide_def['num']}: {slide_def['title']}")

    prs.save(path)
    after = len(prs.slides)
    print(f"\nSaved. Slide count: {before} → {after}")
    assert after == 15, f"Expected 15 slides, got {after}"
    print("✓ 15 slides confirmed.")


if __name__ == "__main__":
    main()
