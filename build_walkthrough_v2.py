"""
build_walkthrough_v2.py
Best-quality rebuild of Santander_LimitedCompany_Walkthrough.pptx

Design philosophy:
  · Every headline is an assertion, not a label
  · Lead with impact — big numbers, specific claims
  · The Command Centre is the centrepiece — 9 slides of full depth
  · Less text, more space, more signal
  · Dark / warm-white / split-panel rhythm throughout
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

# ── Palette ───────────────────────────────────────────────────────────────────
RED    = RGBColor(0xC8, 0x10, 0x2E)
DARK   = RGBColor(0x1C, 0x19, 0x17)
DARK2  = RGBColor(0x29, 0x24, 0x20)
WARM   = RGBColor(0xFA, 0xF6, 0xEF)
STONE1 = RGBColor(0xF5, 0xF0, 0xE8)
STONE2 = RGBColor(0xE7, 0xE5, 0xE4)
STONE5 = RGBColor(0x78, 0x71, 0x6C)
STONE7 = RGBColor(0x44, 0x40, 0x3C)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
AMBER  = RGBColor(0xD9, 0x77, 0x06)
AMBERLT= RGBColor(0xFE, 0xF3, 0xC7)
BLUE   = RGBColor(0x25, 0x63, 0xEB)
BLUELT = RGBColor(0xDB, 0xEA, 0xFE)
EMRLD  = RGBColor(0x05, 0x96, 0x69)
EMRLDLT= RGBColor(0xD1, 0xFA, 0xE5)
REDLT  = RGBColor(0xFE, 0xF2, 0xF2)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Primitives ────────────────────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill=None, line_color=None, line_pt=0.75):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_pt)
    else:
        s.line.fill.background()
    return s

def text(slide, txt, l, t, w, h, size=14, bold=False, color=DARK,
         align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = txt
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = "Calibri"
    return tb

def multiline(slide, lines, l, t, w, h, size=12, color=DARK,
              line_space_pt=4, indent=""):
    """lines = list of str — rendered as separate paragraphs."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = indent + line
        run.font.size  = Pt(size)
        run.font.color.rgb = color
        run.font.name  = "Calibri"
        p.space_after  = Pt(line_space_pt)
    return tb

def top_bar(slide, color=RED):
    rect(slide, 0, 0, 13.33, 0.09, fill=color)

def bottom_bar(slide, color=RED):
    rect(slide, 0, 7.41, 13.33, 0.09, fill=color)

def eyebrow(slide, txt, l=0.55, t=0.18, color=RED):
    text(slide, txt, l, t, 12, 0.32, size=9, bold=True, color=color)

def rule(slide, l=0.55, t=1.1, w=12.23, color=RED):
    rect(slide, l, t, w, 0.035, fill=color)

# ── Layout patterns ───────────────────────────────────────────────────────────

def cover(title_main, title_sub, tagline, url=""):
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
    top_bar(slide)
    bottom_bar(slide)
    # Left red accent column
    rect(slide, 0.55, 1.6, 0.18, 4.1, fill=RED)
    # Eyebrow
    text(slide, "SANTANDER BUSINESS BANKING", 0.9, 1.7, 10, 0.4,
         size=9, bold=True, color=RED)
    # Main title
    text(slide, title_main, 0.9, 2.15, 11, 1.6,
         size=52, bold=True, color=WHITE)
    # Subtitle
    text(slide, title_sub, 0.9, 3.85, 11, 0.55,
         size=18, color=RGBColor(0xC0, 0xBB, 0xB5))
    # Tagline
    text(slide, tagline, 0.9, 4.55, 11, 0.6,
         size=13, italic=True, color=RGBColor(0x78, 0x71, 0x6C))
    if url:
        text(slide, url, 0.9, 6.8, 10, 0.4,
             size=11, color=RGBColor(0x78, 0x71, 0x6C))
    return slide

def section_header(title, subtitle="", eyebrow_txt=""):
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
    top_bar(slide)
    bottom_bar(slide)
    # Thick red left bar
    rect(slide, 0.55, 2.5, 0.2, 2.5, fill=RED)
    if eyebrow_txt:
        text(slide, eyebrow_txt, 0.9, 2.55, 11, 0.4,
             size=10, bold=True, color=RED)
    text(slide, title, 0.9, 3.0, 11, 1.4,
         size=42, bold=True, color=WHITE)
    if subtitle:
        text(slide, subtitle, 0.9, 4.5, 10, 0.8,
             size=15, color=RGBColor(0xA8, 0xA2, 0x9A))
    return slide

def split_panel(eyebrow_txt, headline, left_stat, left_label, body_lines,
                right_accent_color=None):
    """Dark left panel (3.8") with big stat; warm white right panel with body."""
    slide = prs.slides.add_slide(BLANK)
    # Background
    rect(slide, 0, 0, 13.33, 7.5, fill=WARM)
    # Left dark panel
    rect(slide, 0, 0, 4.1, 7.5, fill=DARK)
    top_bar(slide)
    bottom_bar(slide)
    # Left content
    eyebrow(slide, eyebrow_txt, l=0.35, color=RED)
    text(slide, left_stat, 0.3, 1.5, 3.5, 2.0,
         size=72, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    text(slide, left_label, 0.3, 3.7, 3.5, 1.2,
         size=13, color=RGBColor(0x90, 0x8A, 0x84))
    # Right content
    text(slide, headline, 4.5, 0.8, 8.4, 1.2,
         size=28, bold=True, color=DARK)
    rect(slide, 4.5, 2.1, 8.4, 0.04, fill=RED)
    multiline(slide, body_lines, 4.5, 2.25, 8.4, 4.8, size=13, color=STONE7,
              line_space_pt=6)
    return slide

def stat_grid(eyebrow_txt, headline, stats):
    """Dark background, 2×2 or 2×3 grid of big stats."""
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
    top_bar(slide)
    bottom_bar(slide)
    eyebrow(slide, eyebrow_txt, color=RED)
    text(slide, headline, 0.55, 0.35, 12, 0.8,
         size=32, bold=True, color=WHITE)
    rect(slide, 0.55, 1.2, 12.23, 0.04, fill=RED)
    # stats = [(value, label, color), ...]
    cols = 4 if len(stats) == 4 else 3
    W = 12.23 / cols
    for i, (val, lbl, col) in enumerate(stats):
        c = i % cols
        r = i // cols
        x = 0.55 + c * W
        y = 1.4 + r * 2.6
        rect(slide, x + 0.08, y, W - 0.16, 2.35, fill=DARK2)
        rect(slide, x + 0.08, y, 0.09, 2.35, fill=col)
        text(slide, val, x + 0.32, y + 0.2, W - 0.5, 1.3,
             size=52, bold=True, color=WHITE)
        text(slide, lbl, x + 0.32, y + 1.6, W - 0.5, 0.6,
             size=11, color=RGBColor(0x90, 0x8A, 0x84))
    return slide

def content_slide(eyebrow_txt, headline, left_lines, right_lines=None,
                  left_title="", right_title="", note=""):
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, fill=WARM)
    top_bar(slide)
    bottom_bar(slide)
    eyebrow(slide, eyebrow_txt)
    text(slide, headline, 0.55, 0.35, 12.2, 0.85,
         size=30, bold=True, color=DARK)
    rule(slide)
    if right_lines:
        lw = 5.9
        if left_title:
            text(slide, left_title, 0.55, 1.25, lw, 0.32,
                 size=9, bold=True, color=RED)
        rect(slide, 0.55, 1.62, lw, 5.5, fill=WHITE, line_color=STONE2)
        multiline(slide, left_lines, 0.75, 1.75, lw - 0.25, 5.3,
                  size=10, color=DARK, line_space_pt=4)
        rx = 0.55 + lw + 0.42
        rw = 13.33 - rx - 0.42
        if right_title:
            text(slide, right_title, rx, 1.25, rw, 0.32,
                 size=9, bold=True, color=RED)
        rect(slide, rx, 1.62, rw, 5.5, fill=WHITE, line_color=STONE2)
        multiline(slide, right_lines, rx + 0.2, 1.75, rw - 0.25, 5.3,
                  size=10, color=DARK, line_space_pt=4)
    else:
        rect(slide, 0.55, 1.22, 12.23, 5.9, fill=WHITE, line_color=STONE2)
        multiline(slide, left_lines, 0.75, 1.38, 11.8, 5.7,
                  size=11, color=DARK, line_space_pt=5)
    if note:
        text(slide, note, 0.55, 7.1, 12.23, 0.28,
             size=9, italic=True, color=STONE5)
    return slide

def highlight_slide(eyebrow_txt, headline, subhead, callouts, footer=""):
    """Headline + 3-4 coloured callout boxes."""
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, fill=WARM)
    top_bar(slide)
    bottom_bar(slide)
    eyebrow(slide, eyebrow_txt)
    text(slide, headline, 0.55, 0.35, 12.2, 0.85,
         size=30, bold=True, color=DARK)
    rule(slide)
    if subhead:
        text(slide, subhead, 0.55, 1.2, 12.2, 0.5,
             size=13, italic=True, color=STONE5)
    # callouts = [(title, body, bg_color, text_color), ...]
    n = len(callouts)
    W = 12.23 / n
    for i, (ctitle, cbody, bg, tc) in enumerate(callouts):
        x = 0.55 + i * W
        y = 1.85 if subhead else 1.4
        rect(slide, x + 0.06, y, W - 0.12, 5.0, fill=bg)
        text(slide, ctitle, x + 0.22, y + 0.22, W - 0.36, 0.5,
             size=11, bold=True, color=tc)
        multiline(slide, cbody.split('\n'), x + 0.22, y + 0.75,
                  W - 0.36, 4.0, size=10, color=tc, line_space_pt=4)
    if footer:
        text(slide, footer, 0.55, 7.1, 12.23, 0.28,
             size=9, italic=True, color=STONE5)
    return slide

def three_col(eyebrow_txt, headline, cols):
    """Three equal columns with title + body each."""
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, fill=WARM)
    top_bar(slide)
    bottom_bar(slide)
    eyebrow(slide, eyebrow_txt)
    text(slide, headline, 0.55, 0.35, 12.2, 0.8,
         size=30, bold=True, color=DARK)
    rule(slide)
    W = 4.0
    for i, (ctitle, cdesc, lines, accent) in enumerate(cols):
        x = 0.55 + i * (W + 0.16)
        rect(slide, x, 1.25, W, 5.9, fill=WHITE, line_color=STONE2)
        rect(slide, x, 1.25, W, 0.08, fill=accent)
        text(slide, ctitle, x + 0.18, 1.42, W - 0.3, 0.45,
             size=13, bold=True, color=DARK)
        text(slide, cdesc, x + 0.18, 1.9, W - 0.3, 0.38,
             size=10, italic=True, color=STONE5)
        multiline(slide, lines, x + 0.18, 2.35, W - 0.3, 4.55,
                  size=10, color=STONE7, line_space_pt=4)
    return slide

def back_cover(name, email, url, note):
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
    top_bar(slide)
    bottom_bar(slide)
    rect(slide, 0.55, 2.8, 0.18, 2.0, fill=RED)
    text(slide, "Thank you.", 0.9, 2.9, 11, 1.0,
         size=52, bold=True, color=WHITE)
    text(slide, f"{name}  ·  {email}", 0.9, 4.1, 11, 0.45,
         size=14, color=RGBColor(0xC0, 0xBB, 0xB5))
    text(slide, url, 0.9, 4.65, 11, 0.38,
         size=12, color=RED)
    text(slide, note, 0.9, 6.7, 11.5, 0.5,
         size=9, italic=True, color=RGBColor(0x55, 0x50, 0x4C))
    return slide

# ═══════════════════════════════════════════════════════════════════════════════
# BUILD THE DECK
# ═══════════════════════════════════════════════════════════════════════════════

# ── 1. COVER ──────────────────────────────────────────────────────────────────
cover(
    title_main="Santander\nBusiness Banking",
    title_sub="A complete digital prototype — built alone — to show what's possible.",
    tagline="Limited Company walkthrough · June 2026",
    url="imscots-ui.github.io/santander  ·  Alan Davidson  ·  Alan.Davidson@santander.co.uk"
)

# ── 2. THE PROBLEM ────────────────────────────────────────────────────────────
content_slide(
    eyebrow_txt="THE PROBLEM  ·  WHY THIS EXISTS",
    headline="UK business banking is still mostly paper. In 2026.",
    left_lines=[
        "Four things every UK limited company has to do — that no banking app can do today:",
        "",
        "01  CHANGE A MANDATE",
        "    Add or remove a signatory. The form goes to Sunderland by post.",
        "    Average time: five working days. Often longer.",
        "",
        "02  FILE A VAT RETURN",
        "    Making Tax Digital requires software. The bank is not involved.",
        "    The connection between your transactions and your HMRC submission is manual.",
        "",
        "03  PAY YOUR TEAM",
        "    Bulk BACS payments require either a separate system or a phone call.",
        "    The bank holds the data. The bank doesn't offer the workflow.",
        "",
        "04  CONFIRM YOUR IDENTITY",
        "    KYC and KYB renewals. Another form. Another envelope. Another wait.",
        "",
        "─────────────────────────────────────────────────────────────",
        "These aren't edge cases. They are the most common tasks",
        "a business banking customer needs to do — every quarter.",
    ],
    right_lines=[
        "WHAT THE PROTOTYPE PROVES",
        "",
        "Every one of these tasks can be done:",
        "",
        "  ✓  On a mobile phone",
        "  ✓  In under ten minutes",
        "  ✓  Without a form, an envelope, or a phone call",
        "  ✓  With a full audit trail",
        "  ✓  Within FCA BCOBS 4A cooling-off rules",
        "",
        "The prototype is not a wireframe.",
        "It is not a Figma mock-up.",
        "It is not a proof of concept.",
        "",
        "It is a fully functional, interactive product",
        "that runs in any mobile browser — right now.",
        "",
        "imscots-ui.github.io/santander",
        "",
        "Built by one person.",
        "Self-initiated.",
        "In spare time.",
    ],
    left_title="THE FOUR BANKING PAIN POINTS",
    right_title="THE RESPONSE"
)

# ── 3. WHAT WAS BUILT ─────────────────────────────────────────────────────────
stat_grid(
    eyebrow_txt="WHAT WAS BUILT  ·  JUNE 2026",
    headline="One engineer. One file. Everything a UK business bank needs.",
    stats=[
        ("5,700",  "lines of React — single file, zero dependencies beyond Vite + Tailwind", RED),
        ("11",     "complete paperless workflows — each a multi-step process with audit trail", AMBER),
        ("7",      "business entity types — Limited, Sole Trader, LLP, Partnership, Charity, Club, Society", BLUE),
        ("0",      "paper forms required for any task the prototype covers", EMRLD),
    ]
)

# ── 4. HOW TO ACCESS ──────────────────────────────────────────────────────────
content_slide(
    eyebrow_txt="ACCESS  ·  IT'S LIVE RIGHT NOW",
    headline="Open it. Use it. No login, no install, no server.",
    left_lines=[
        "STEP 1 — OPEN THE URL",
        "    imscots-ui.github.io/santander",
        "    Works on any modern mobile browser.",
        "    Chrome on Android or Safari on iPhone recommended.",
        "",
        "STEP 2 — IF THE SCREEN LOOKS OUTDATED",
        "    The CDN caches aggressively. Long-press the reload button",
        "    on your phone and choose 'Hard refresh' — or open",
        "    a private / incognito tab.",
        "",
        "STEP 3 — NAVIGATE",
        "    Bottom bar: Home · Approve · Audit · Tax · Statements",
        "    Home screen → tap any action tile to start a workflow.",
        "    Workflows are full-screen overlays. Back button always cancels.",
        "",
        "STEP 4 — SWITCH ENTITY TYPES",
        "    Home screen → tap the entity pill (e.g. 'Limited Company')",
        "    to switch between all 7 business types live, mid-demo.",
        "",
        "DESKTOP MODE",
        "    Tap the Desktop toggle in the header for a sidebar layout.",
        "    All state is shared — switching preserves everything in progress.",
    ],
    right_lines=[
        "WHAT WORKS",
        "",
        "  ✓  All 5 screens — Home, Approve, Audit, Tax, Statements",
        "  ✓  All 11 workflow wizards",
        "  ✓  OTP confirmation (fake 6-digit code — any digits work)",
        "  ✓  24-hour cooling-off with live progress bar",
        "  ✓  VAT countdown overlay before HMRC submission",
        "  ✓  Voice memo (mock AI transcription)",
        "  ✓  Receipt scanning (mock OCR categorisation)",
        "  ✓  Relationship manager escalation (Priya's on it)",
        "  ✓  Payment sequencer optimiser",
        "  ✓  Desktop / mobile view toggle",
        "  ✓  Accessibility settings (dyslexia, motion, contrast, text size)",
        "  ✓  Voice ID set-up and session anomaly detection",
        "",
        "WHAT DOESN'T WORK (intentionally)",
        "",
        "  ✗  No real bank connection",
        "  ✗  No real HMRC submission",
        "  ✗  No real payment processing",
        "  ✗  No authentication or user accounts",
        "",
        "This is a prototype. Every number is realistic",
        "but fictional. No real data is stored.",
    ],
    left_title="HOW TO USE IT",
    right_title="WHAT'S REAL"
)

# ── SECTION: THE COMMAND CENTRE ───────────────────────────────────────────────
section_header(
    title="The Command\nCentre",
    subtitle="One screen that knows everything about your business — before you ask.",
    eyebrow_txt="THE HOME SCREEN"
)

# ── 5. COMMAND CENTRE OVERVIEW ────────────────────────────────────────────────
content_slide(
    eyebrow_txt="COMMAND CENTRE  ·  WHAT'S ON THE HOME SCREEN",
    headline="Fifteen components. Every signal a director needs. Zero navigation required.",
    left_lines=[
        "01  Editorial greeting  — personalised, date, context sentence",
        "02  Total Balance hero card  — dark, animated, all accounts in one number",
        "03  Proactive Forecast strip  — 3 nudge cards, colour-coded by urgency",
        "04  Pre-approved lending offer  — contextual, based on trading history",
        "05  Business Health Score  — live 0–100 weighted metric",
        "06  Cooling-off cards  — 24h delay requests with live progress bar",
        "07  Stalled requests  — 'Priya's on it' RM escalation card",
        "08  MTD VAT teaser  — live countdown, amount owed, transactions to review",
    ],
    right_lines=[
        "09  Paperless banner  — '3 paper forms retired · No more posting to Sunderland'",
        "10  Pending approvals  — items awaiting the director's signature",
        "11  Paperless Actions grid  — 10 tiles, every banking task in one tap",
        "12  Accounts list  — all accounts, sort codes, balances, mandate badges",
        "13  Supplier Radar  — counterparty intelligence",
        "14  13-week forecast chart  — smooth area chart, event markers, floor warning",
        "15  Director Centre  — company filings and compliance calendar",
        "",
        "─────────────────────────────────────────────────",
        "Design principle: surface the right signal at the right time.",
        "The user should never have to navigate to find out what needs doing today.",
        "The bank should already know — and tell them.",
    ],
    left_title="COMPONENTS 01–08",
    right_title="COMPONENTS 09–15"
)

# ── 6. GREETING + BALANCE HERO ────────────────────────────────────────────────
highlight_slide(
    eyebrow_txt="COMMAND CENTRE  ·  COMPONENT 01–02",
    headline="Good morning, James. Everything's running normally.",
    subhead="The most important number in your business — front and centre, the moment you open the app.",
    callouts=[
        (
            "The editorial greeting",
            "Personalised to the logged-in director.\n"
            "Shows today's date.\n"
            "A single plain-English sentence:\n"
            "'Everything's running normally today.'\n"
            "Or: '2 things need your signature — they're below.'\n\n"
            "No jargon. No codes.\n"
            "Written in the voice of a trusted banker.",
            STONE1, DARK
        ),
        (
            "Total Balance hero card",
            "Dark card — animated red orb floats in the background.\n"
            "One large number: total across all accounts.\n"
            "Entity name and sort code.\n"
            "Mandate rule badge (e.g. 'Any 1 signatory').\n"
            "FSCS Protected tag.\n"
            "Credit ring-fenced badge (if applicable).\n\n"
            "Why dark: the balance is the most important number.\n"
            "It deserves weight.",
            DARK, WHITE
        ),
        (
            "Why it matters",
            "Traditional banking apps open to a list of accounts.\n"
            "This opens to a greeting and a number.\n\n"
            "The difference is intent:\n"
            "An account list is a ledger.\n"
            "A greeting is a relationship.\n\n"
            "In user research, 'knowing what I have'\n"
            "is the #1 reason people open their banking app.\n"
            "We answer it immediately.",
            REDLT, RGBColor(0x99, 0x0D, 0x24)
        ),
    ]
)

# ── 7. PROACTIVE FORECAST STRIP ───────────────────────────────────────────────
highlight_slide(
    eyebrow_txt="COMMAND CENTRE  ·  COMPONENT 03  ·  THE FLAGSHIP",
    headline="Predict. Explain. Act.",
    subhead="The banking app that tells you what needs doing — before you think to ask.",
    callouts=[
        (
            "AMBER — Payroll run · in 3 days",
            "Calendar icon · Amber border\n\n"
            "Payroll run · £42,180\n"
            "'Balance covers it — review if payees changed'\n\n"
            "→  Review payroll\n\n"
            "One tap: opens the bulk payments workflow.\n"
            "No navigation. No menu.\n"
            "Predict → Explain → Act.",
            AMBERLT, RGBColor(0x78, 0x35, 0x00)
        ),
        (
            "RED — VAT return · due in 22 days",
            "Receipt icon · Red border\n\n"
            "VAT return due 7 Aug\n"
            "'Q3 Making Tax Digital · HMRC submission'\n\n"
            "→  Start now\n\n"
            "One tap: opens the MTD tax screen.\n"
            "The bank knows the deadline.\n"
            "The bank shows the amount.\n"
            "The bank offers to help file it.",
            REDLT, RGBColor(0x7F, 0x08, 0x1D)
        ),
        (
            "BLUE — Companies House · in 14 days",
            "Building icon · Blue border\n\n"
            "Annual confirmation due\n"
            "'Companies House · Northgate Systems Ltd'\n\n"
            "→  Run KYB check\n\n"
            "One tap: opens the ID register workflow.\n"
            "Cross-references the company's current\n"
            "directors and filings automatically.",
            BLUELT, RGBColor(0x1A, 0x3D, 0x7F)
        ),
    ],
    footer="These three cards update dynamically based on the entity type selected. Charity → Charity Commission. Sole Trader → no payroll card shown."
)

# ── 8. LENDING OFFER ─────────────────────────────────────────────────────────
split_panel(
    eyebrow_txt="COMMAND CENTRE  ·  COMPONENT 04",
    headline="You're pre-approved for £45,000 — based on your real trading history.",
    left_stat="£45k",
    left_label="pre-approved business loan\nfrom 6.9% APR\nvalid until 31 July 2026",
    body_lines=[
        "The lending card appears on the Home screen before the customer goes looking for a loan.",
        "",
        "WHY THIS WORKS",
        "  •  Not a generic offer — it references the customer's actual 6-month trading record",
        "  •  The APR and term are shown immediately — no 'call us' gatekeeping",
        "  •  One tap to draw down — a 4-step workflow with OTP confirmation",
        "  •  Once accepted, the card is replaced by a 'loan active' confirmation",
        "",
        "WHAT CHANGES AFTER DRAW-DOWN",
        "  The hero card updates. The Home screen shows:",
        "  'Business loan active · £45,000 drawn'",
        "  'First repayment 1 Aug 2026 · repaying over 36 months'",
        "",
        "DESIGN INTENT",
        "  The banking relationship already has the data to make a decision.",
        "  The offer surfaces that decision proactively.",
        "  The customer doesn't apply — they accept or decline.",
        "  This is banking as a service, not banking as a counter.",
    ]
)

# ── 9. COOLING-OFF + STALLED ──────────────────────────────────────────────────
content_slide(
    eyebrow_txt="COMMAND CENTRE  ·  COMPONENTS 06–07",
    headline="The 24-hour safety net — and what happens when it catches something.",
    left_lines=[
        "COOLING-OFF  (Component 06)",
        "",
        "Some requests — account closures, mandate changes to higher-risk rules —",
        "are subject to a mandatory 24-hour cooling-off period.",
        "FCA BCOBS 4A — consumer protection regulation.",
        "",
        "The Home screen shows each request as a card with:",
        "  •  A plain-English description of what was requested",
        "  •  When it will execute ('We'll do this on Monday morning')",
        "  •  A live amber progress bar — drains in real time",
        "  •  'Cancel' — triggers a confirmation before dropping the request",
        "  •  'Do it now' — removes the delay and executes immediately",
        "",
        "The progress bar updates every 30 seconds via a tick useEffect.",
        "No polling. No intervals without cleanup.",
        "Working days only — pauses on weekends and bank holidays.",
    ],
    right_lines=[
        "STALLED REQUESTS  (Component 07)  — 'Priya's on it'",
        "",
        "When a workflow requires two or more signatories (mandate rule 'any-2' or 'all'),",
        "the co-signer gets a push notification and has a window to respond.",
        "",
        "If the co-signer doesn't respond in time — the request stalls.",
        "",
        "WHAT HAPPENS",
        "  The Home screen shows a 'Priya's on it' card.",
        "  Priya is the customer's Relationship Manager.",
        "  The system auto-notifies Priya within 2 minutes.",
        "",
        "The card says:",
        "  'A signature didn't arrive in time.'",
        "  'Priya saw this within 2 minutes and will reach out today.'",
        "  'Nothing's been lost — just paused.'",
        "",
        "Two actions: 'Got it' (dismiss) and 'Call Priya' (opens RM sheet).",
        "",
        "DESIGN INTENT",
        "  A missed co-signer is a known pain point. Instead of showing",
        "  'Request expired — please try again', the system surfaces a human.",
        "  The bank doesn't fail silently. It escalates visibly.",
    ],
    left_title="24-HOUR COOLING-OFF",
    right_title="RELATIONSHIP MANAGER ESCALATION"
)

# ── 10. MTD VAT TEASER + ACTIONS ─────────────────────────────────────────────
content_slide(
    eyebrow_txt="COMMAND CENTRE  ·  COMPONENTS 08, 09, 10, 11",
    headline="The bank knows your VAT deadline. Your bank knows your payees. Your bank knows everything.",
    left_lines=[
        "MTD VAT TEASER  (Component 08)",
        "",
        "A persistent card on the Home screen, updated in real time:",
        "  '22 days · VAT due 7 Aug'",
        "  '£4,118.00 owed to HMRC'",
        "  '47 transactions to review · Q3 Jul–Sep 2026'",
        "",
        "One tap: opens the Making Tax Digital screen.",
        "The bank has already categorised the transactions.",
        "You review, declare, and submit.",
        "",
        "PAPERLESS BANNER  (Component 09)",
        "",
        "  '3 paper forms retired'",
        "  'No more posting to Sunderland · 5 days → minutes'",
        "",
        "This is deliberately specific — 'Sunderland' is where Santander's",
        "business banking paper forms are physically processed.",
        "Real customers know this. Seeing it named builds immediate trust.",
    ],
    right_lines=[
        "PAPERLESS ACTIONS GRID  (Components 10–11)",
        "",
        "Ten action tiles in a 2-column grid. Every banking task.",
        "",
        "  Bulk payments    — CSV · BACS · FP · CHAPS",
        "  International    — FX · SWIFT · SEPA",
        "  Change mandate   — Add, remove, signing rule",
        "  Business details — Name, address, contact",
        "  Scan receipt     — Auto-categorise for MTD",
        "  ID register      — KYC Lists 1, 2 & 3",
        "  Dormant accounts — Reactivate or close",
        "  Close account    — Form ANB9 0370",
        "  Voice memo       — Speak → expense auto-tagged",
        "  Optimise payments — 30-day payment sequencer",
        "",
        "Each tile links directly to a workflow.",
        "No sub-menus. No search. One tap.",
        "",
        "The tiles adapt to entity type:",
        "  Charity → 'Mandate & members' instead of 'Change mandate'",
        "  Sole Trader → Payroll tile hidden (no employees)",
    ],
    left_title="MTD + PAPERLESS BANNER",
    right_title="TEN ACTION TILES"
)

# ── 11. FORECAST CHART ────────────────────────────────────────────────────────
split_panel(
    eyebrow_txt="COMMAND CENTRE  ·  COMPONENT 14",
    headline="Thirteen weeks of cash flow — before you need to worry about it.",
    left_stat="13",
    left_label="weeks of projected\ncash flow — from\nlive account balances",
    body_lines=[
        "The 13-week forecast chart is a smooth gradient area chart — not a bar chart.",
        "",
        "WHY AREA, NOT BARS",
        "  A balance over time is a trend, not a set of independent values.",
        "  A smooth curve communicates continuity. Bars imply discrete events.",
        "  The chart is zoomed to the actual min–max balance band,",
        "  so week-on-week movement is legible even when the balance is large.",
        "",
        "WHAT'S ON THE CHART",
        "  •  Red dot = scheduled outflow (e.g. payroll run)",
        "  •  Rose dot = expected receipt (e.g. customer payment)",
        "  •  Amber ring = lowest forecast week",
        "  •  Dashed amber line = £80k working-capital floor (if crossed)",
        "",
        "ALWAYS EXPLAINS ITSELF",
        "  'Largest outflow: Payroll run £42,180 in week 3 (15 Jul)'",
        "  'Balance dips to £68,210 in week 3, below the £80k",
        "   working-capital floor — consider transferring from Reserve.'",
        "",
        "A chart that doesn't explain itself is decoration.",
        "This one tells the director what to do next.",
    ]
)

# ── SECTION: ELEVEN WORKFLOWS ─────────────────────────────────────────────────
section_header(
    title="11 Workflows.\nZero Paper.",
    subtitle="Every banking task a UK limited company needs — digital, audited, done.",
    eyebrow_txt="THE WORKFLOWS"
)

# ── 12. ALL 11 WORKFLOWS ─────────────────────────────────────────────────────
content_slide(
    eyebrow_txt="WORKFLOWS  ·  COMPLETE REFERENCE",
    headline="Every paper form, replaced. Every phone call, self-served.",
    left_lines=[
        "01  ACCOUNT CLOSURE  (renderClosure)",
        "    Form ANB9 0370 — 5 steps — 24h cooling-off — full audit",
        "",
        "02  BUSINESS DETAILS  (renderBiz)",
        "    Name, address, contact — 4 steps — OTP confirmation",
        "",
        "03  MANDATE CHANGE  (renderMandate)",
        "    Add / remove signatories, change signing rule",
        "    Triggered by strictest account mandate — co-signer OTP if required",
        "",
        "04  BULK PAYMENTS  (renderWages)",
        "    CSV upload → payee review → schedule",
        "    BACS · Faster Payments · CHAPS",
        "    10-second countdown before execution",
        "",
        "05  LENDING DRAW-DOWN  (renderLending)",
        "    Pre-approved offer → terms review → OTP → active",
        "",
        "06  FX / INTERNATIONAL  (renderFX)",
        "    SWIFT · SEPA · rate lock → counterparty → OTP",
    ],
    right_lines=[
        "07  DORMANT ACCOUNTS  (renderDormancy)",
        "    Reactivate or close — 3 steps",
        "",
        "08  ACCOUNT UNLINKING  (renderUnlink)",
        "    Personal / business separation — 4 steps — FCA BCOBS",
        "",
        "09  RING-FENCING  (renderRingfence)",
        "    Credit ring-fence opt-in — 5 steps — compliance declaration",
        "",
        "10  KYC / KYB ID CHECK  (renderIdCheck)",
        "    ID register: Lists 1, 2 & 3",
        "    Companies House annual confirmation",
        "    Charity Commission (for charity entity type)",
        "",
        "11  MTD VAT SUBMISSION  (renderMtdSubmit)",
        "    Transaction review → categorisation → declaration",
        "    5-second HMRC countdown overlay",
        "    HMRC MTD VAT API → receipt in audit log",
        "",
        "─────────────────────────────────────────────────",
        "Every workflow uses the same StepFrame shell:",
        "Back button · step indicator · sticky footer.",
        "Every workflow produces an audit trail entry.",
        "Every workflow resets cleanly on cancel.",
    ],
    left_title="WORKFLOWS 01–06",
    right_title="WORKFLOWS 07–11"
)

# ── 13. MANDATE CHANGE ────────────────────────────────────────────────────────
split_panel(
    eyebrow_txt="WORKFLOW 03  ·  MOST COMMON PAIN POINT",
    headline="Change a signatory — from 5 days by post to 6 minutes on your phone.",
    left_stat="5 days",
    left_label="→ 6 minutes\nmandate change\nwithout a form",
    body_lines=[
        "THE CURRENT STATE",
        "  Adding or removing a company signatory at Santander requires a paper form.",
        "  That form goes to Sunderland. Average turnaround: 5 working days.",
        "",
        "THE PROTOTYPE WORKFLOW  (6 steps)",
        "  1.  Select the account(s) to change the mandate on",
        "  2.  Choose new signing rule: Any 1 / Any 2 / All signatories",
        "  3.  Add new signatories (name, email, role)",
        "  4.  Review — the strictest mandate across selected accounts applies",
        "  5.  OTP confirmation via SMS",
        "  6.  If new rule requires cooling-off: 24h delay, shown on Home screen",
        "",
        "MANDATE LOGIC",
        "  getMandateFor(accountNos) picks the strictest rule across selected accounts.",
        "  'all' > 'any-2' > 'any-1'",
        "  A request covering two accounts — one 'any-1', one 'all' — requires all signatories.",
        "  This is the correct regulatory behaviour, enforced in the UI automatically.",
        "",
        "WHAT THE CUSTOMER SEES",
        "  The mandate rule badge on each account changes instantly on confirmation.",
        "  The cooling-off card appears on the Home screen if applicable.",
        "  The audit trail records the exact change, who made it, and when.",
    ]
)

# ── 14a. MTD VAT SUBMISSION — THE FIVE STEPS ─────────────────────────────────
content_slide(
    eyebrow_txt="WORKFLOW 11  ·  MTD HMRC VAT SUBMISSION  ·  PART 1 OF 2",
    headline="File your VAT return from your phone — in 5 steps.",
    left_lines=[
        "STEP 1  —  REVIEW TRANSACTIONS",
        "  All transactions in the VAT quarter shown by the bank.",
        "  Each one is pre-categorised. Confirm or correct each line.",
        "",
        "STEP 2  —  CATEGORISE",
        "  Receipts scanned earlier auto-populate this step.",
        "  Voice memos from the Home screen also feed in.",
        "  Manual override available per transaction.",
        "",
        "STEP 3  —  REVIEW FIGURES",
        "  VAT due: £4,118.00",
        "  Output tax (sales): £18,472.40",
        "  Input tax (purchases): £14,354.40",
        "  Period: Q3 Jul–Sep 2026 · Due: 7 Nov 2026",
    ],
    right_lines=[
        "STEP 4  —  DECLARATION",
        "  'I declare that the information given in this return",
        "   is true, complete and accurate.'",
        "  — statutory HMRC wording.",
        "  Checkbox. Timestamp. Logged in audit trail.",
        "  Director's name and role recorded.",
        "",
        "STEP 5  —  THE COUNTDOWN",
        "  'Filing your VAT return'",
        "  A 5-second countdown ring with 'sec' label.",
        "  Amber note: 'Filed returns cannot be corrected",
        "   through this app — contact HMRC to amend.'",
        "  Grey Cancel button — not red, not alarming.",
        "  At zero: submitted to HMRC.",
        "  Receipt number logged in the Audit tab.",
    ],
    left_title="STEPS 1–3",
    right_title="STEPS 4–5"
)

# ── 14b. MTD VAT SUBMISSION — DESIGN RATIONALE ────────────────────────────────
content_slide(
    eyebrow_txt="WORKFLOW 11  ·  MTD HMRC VAT SUBMISSION  ·  PART 2 OF 2",
    headline="Why the countdown. Why grey. Why the bank — not the spreadsheet.",
    left_lines=[
        "WHY THE COUNTDOWN",
        "",
        "HMRC VAT returns, once filed, cannot be amended through the same",
        "API call — a correction requires a separate amendment.",
        "",
        "The 5-second countdown gives the user one last chance to stop.",
        "Without panic. Without red buttons.",
        "",
        "DESIGN DECISIONS IN THE COUNTDOWN OVERLAY",
        "",
        "  Heading: 'Filing your VAT return'",
        "  — not 'One last chance to stop'",
        "    (previous version; changed — it was alarming)",
        "",
        "  Cancel button: stone-100 grey (secondary)",
        "  — not red (previous version confused users:",
        "    cancel should not look like a danger action)",
        "",
        "  Ring: r=46, larger than original, with 'sec' label",
        "  Details pill: left-aligned · label · £4,118.00",
        "  Amber warning: shortened to one clear sentence",
    ],
    right_lines=[
        "WHY IT MATTERS — THE BIGGER PICTURE",
        "",
        "Most UK businesses still file VAT via bridging software:",
        "  •  Accountant exports bank data",
        "  •  Accountant pastes into spreadsheet",
        "  •  Spreadsheet feeds into bridging software",
        "  •  Bridging software submits to HMRC MTD API",
        "",
        "The bank already holds all the transaction data.",
        "",
        "This prototype shows that the bank could:",
        "  •  Auto-categorise every transaction",
        "  •  Calculate the VAT figures",
        "  •  Present the declaration",
        "  •  Submit directly to HMRC",
        "",
        "Without the customer ever opening a spreadsheet.",
        "Without an accountant involved in the mechanics.",
        "Without bridging software.",
        "",
        "The bank becomes the tax agent —",
        "not just the ledger.",
    ],
    left_title="COUNTDOWN — DESIGN RATIONALE",
    right_title="THE BIGGER PICTURE"
)

# ── 15. OTHER WORKFLOWS ───────────────────────────────────────────────────────
three_col(
    eyebrow_txt="WORKFLOWS  ·  FURTHER DETAIL",
    headline="Bulk payments. International FX. Account closure. All paperless.",
    cols=[
        (
            "Bulk Payments — Wages",
            "renderWages  ·  5 steps",
            [
                "CSV file → payee review",
                "Each payee: name, sort code,",
                "account, amount, reference",
                "",
                "Mandate check: if total > threshold,",
                "co-signer approval required",
                "",
                "Payment rail: BACS / Faster",
                "Payments / CHAPS",
                "",
                "10-second countdown before",
                "execution — cancel anytime",
                "",
                "Audit entry: each payee,",
                "amount, time, authorising directors",
            ],
            AMBER
        ),
        (
            "FX / International",
            "renderFX  ·  4 steps",
            [
                "Currency selection",
                "Amount + rate lock (live mock rate)",
                "SWIFT / SEPA / correspondent",
                "",
                "Counterparty management:",
                "saved payees with full bank",
                "details + compliance flags",
                "",
                "OTP confirmation",
                "",
                "Supports: EUR, USD, CHF,",
                "JPY, CAD, AUD — plus",
                "SEPA zone countries",
                "",
                "FX rate displayed with",
                "spread and settlement date",
            ],
            BLUE
        ),
        (
            "Account Closure",
            "renderClosure  ·  Form ANB9 0370",
            [
                "Select account to close",
                "State reason (regulatory)",
                "Transfer remaining balance",
                "",
                "Mandatory 24h cooling-off:",
                "FCA consumer protection",
                "",
                "Cannot close if:",
                "  · Outstanding direct debits",
                "  · Pending transactions",
                "  · Only business account",
                "",
                "Audit entry: full closure",
                "record, timestamp, director",
                "who authorised",
            ],
            EMRLD
        ),
    ]
)

# ── SECTION: HOW TO DEMO ──────────────────────────────────────────────────────
section_header(
    title="The Demo",
    subtitle="A complete Limited Company walkthrough — from open to submitted VAT return.",
    eyebrow_txt="PRESENTER GUIDE"
)

# ── 16a. DEMO SCRIPT — PART 1 ─────────────────────────────────────────────────
content_slide(
    eyebrow_txt="PRESENTER GUIDE  ·  LIMITED COMPANY WALKTHROUGH  ·  PART 1 OF 2",
    headline="Suggested script — 10 minutes, every headline feature. Opening + Command Centre + VAT.",
    left_lines=[
        "OPENING  (1 min)",
        "  Open imscots-ui.github.io/santander on your phone.",
        "  'This is a fully working banking app — built as a prototype",
        "   to show what Business Banking could look like.'",
    ],
    right_lines=[
        "VAT SUBMISSION  (3 min)",
        "  Tap the Tax tab (bottom bar).",
        "  'Making Tax Digital. The bank has already categorised",
        "   your transactions. You review, declare, and submit.'",
        "  Walk through steps 1–3.",
        "  On step 3: 'Figures come directly from the bank's",
        "   transaction data — no spreadsheet, no bridging software.'",
        "  Tap Submit → show the countdown:",
        "  '5 seconds. One last chance to cancel.",
        "   Then it goes to HMRC. Receipt in the audit log.'",
    ],
    left_title="OPENING",
    right_title="VAT SUBMISSION"
)

# command centre gets its own full-width block
content_slide(
    eyebrow_txt="PRESENTER GUIDE  ·  COMMAND CENTRE SECTION",
    headline="The Command Centre — 3 minutes of your 10-minute demo.",
    left_lines=[
        "THE COMMAND CENTRE  (3 min)",
        "",
        "  Point to the balance hero card:",
        "  'Total balance, one number, immediately.'",
        "",
        "  Scroll to the Forecast strip:",
        "  'The bank tells you what needs doing.",
        "   Payroll in 3 days. VAT in 22 days.",
        "   Companies House in 14 days.",
        "   Each one taps directly into the workflow.'",
        "",
        "  Tap 'Review payroll'",
        "  → show the wages workflow briefly → back.",
        "",
        "  Scroll to the Actions grid:",
        "  '10 tasks. All the ones that currently require",
        "   a phone call or a letter to Sunderland.'",
        "",
        "  Scroll to the forecast chart:",
        "  '13 weeks of cash flow. The bank knows payroll",
        "   lands in week 3. It flags that the balance",
        "   dips below £80k — consider a Reserve transfer.'",
    ],
)

# ── 16b. DEMO SCRIPT — PART 2 ─────────────────────────────────────────────────
content_slide(
    eyebrow_txt="PRESENTER GUIDE  ·  LIMITED COMPANY WALKTHROUGH  ·  PART 2 OF 2",
    headline="Mandate change. Entity switching. Closing line. Follow-up questions.",
    left_lines=[
        "MANDATE CHANGE  (2 min)",
        "  Tap Home → Paperless Actions → 'Change mandate'",
        "  'Adding a signatory currently takes 5 days by post.",
        "   Here it takes 6 minutes. OTP confirmation.",
        "   Cooling-off if required by the mandate rule.",
        "   Audit trail automatically created.'",
        "  Show step 1: select account.",
        "  Show step 2: choose signing rule.",
        "  Cancel — no need to complete the workflow.",
        "",
        "CLOSING  (1 min)",
        "  Tap the entity pill at the top of the Home screen.",
        "  'It supports 7 business types — Limited Company,",
        "   Sole Trader, LLP, Partnership, Charity, Club, Society.",
        "   Switch live in the demo and every screen updates.",
        "   Different accounts. Different mandate rules.",
        "   Different KYC requirements.'",
        "  Switch to Charity → show entity name changes.",
        "  Switch back to Limited Company.",
    ],
    right_lines=[
        "THE CLOSING LINE",
        "",
        "  'This is one prototype.",
        "   Built alone. In spare time.",
        "   It's live right now — you can open it on",
        "   your phone in this room.',",
        "   'The question is: what could a team build?'",
        "",
        "─────────────────────────────────────────────",
        "AFTER THE DEMO",
        "Suggested follow-up questions:",
        "",
        "  •  What workflows are most painful for your",
        "     business customers right now?",
        "",
        "  •  Would you like to see any specific entity",
        "     type — Charity, LLP, Sole Trader?",
        "",
        "  •  What would you want to see in a",
        "     production roadmap?",
        "",
        "  •  Who else in your team should see this?",
    ],
    left_title="MANDATE CHANGE + CLOSING",
    right_title="CLOSING LINE + FOLLOW-UP"
)

# ── 17. ENTITY SWITCHING ──────────────────────────────────────────────────────
highlight_slide(
    eyebrow_txt="PRESENTER GUIDE  ·  ENTITY SWITCHING",
    headline="Switch between 7 business types live — mid-demo. Everything updates.",
    subhead="Tap the entity pill on the Home screen (e.g. 'Limited Company') to open the Entity sheet.",
    callouts=[
        (
            "Limited Company",
            "4 accounts\nAny-1 · Any-2 · All\nmandates\n\nDirectors\nCompanies House #\nForm ANB9 0370\nMTD VAT enabled",
            STONE1, DARK
        ),
        (
            "Sole Trader",
            "2 accounts\nAny-1 mandate only\n\nNo directors —\njust the owner\nNo payroll tile\nSelf-assessment MTD",
            STONE1, DARK
        ),
        (
            "Charity",
            "3 accounts\nAll-signatory mandate\n\nTrustees not directors\nCharity Commission #\nNot Companies House\nGift Aid tile added",
            STONE1, DARK
        ),
        (
            "LLP / Partnership",
            "3–4 accounts\nAny-2 or All mandate\n\nDesignated members\nor partners\nLLP agreement refs\nNo share capital",
            STONE1, DARK
        ),
    ],
    footer="Club and Society entity types are also available — they use a Treasurer role rather than Director, and have simplified account structures."
)

# ── 18. BACK COVER ────────────────────────────────────────────────────────────
back_cover(
    name="Alan Davidson",
    email="Alan.Davidson@santander.co.uk",
    url="imscots-ui.github.io/santander",
    note="Self-initiated concept piece. Prototype only — no backend, no real data. Built to demonstrate what's possible."
)

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/home/user/santander/Santander_LimitedCompany_Walkthrough.pptx"
prs.save(out)
from pptx import Presentation as _P
p = _P(out)
import os
print(f"Saved: {out}")
print(f"Slides: {len(p.slides)}")
print(f"Size: {os.path.getsize(out)//1024} KB")
