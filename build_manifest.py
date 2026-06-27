"""
HMS 1701 — Project Record / Manifest Deck
Produces: Santander_Project_Record.pptx
A comprehensive document-style deck capturing everything built, written, and delivered.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

# ── Brand colours ─────────────────────────────────────────────────────────────
RED      = RGBColor(0xC8, 0x10, 0x2E)
DARK     = RGBColor(0x1C, 0x19, 0x17)
WARM     = RGBColor(0xFA, 0xF6, 0xEF)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
STONE1   = RGBColor(0xF5, 0xF5, 0xF4)
STONE2   = RGBColor(0xE7, 0xE5, 0xE4)
STONE3   = RGBColor(0xD6, 0xD3, 0xD1)
STONE5   = RGBColor(0x78, 0x71, 0x6C)
STONE7   = RGBColor(0x44, 0x40, 0x3C)
LIGHTRED = RGBColor(0xFE, 0xF2, 0xF2)
EMERALD  = RGBColor(0x05, 0x96, 0x69)
AMBER    = RGBColor(0xD9, 0x77, 0x06)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────

def R(slide, l, t, w, h, fill=None, line=None):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.line.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    return sh

def T(slide, text, l, t, w, h, size=12, bold=False, color=DARK,
      align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb

def ML(slide, lines, l, t, w, h, size=11, color=DARK, leading=5):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(leading)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"

def red_bar(s):
    R(s, 0, 0, 13.33, 0.06, fill=RED)

def section_header(s, title, subtitle=None):
    red_bar(s)
    T(s, title, 0.5, 0.15, 12, 0.52, size=22, bold=True, color=DARK)
    if subtitle:
        T(s, subtitle, 0.5, 0.68, 12, 0.3, size=11, color=STONE5)
    R(s, 0.5, 0.98, 12.33, 0.018, fill=STONE3)

def advisor_stamp(s):
    R(s, 0, 7.22, 13.33, 0.28, fill=DARK)
    T(s, "Alan Davidson · Business Banking Advisor · Self-initiated · Completed out of hours in own time · June 2026",
      0.5, 7.24, 12.33, 0.22, size=8, color=STONE3, align=PP_ALIGN.CENTER)

def slide_num(s, n):
    T(s, str(n), 12.9, 7.2, 0.35, 0.25, size=8, color=STONE5, align=PP_ALIGN.RIGHT)

def tag(s, text, x, y, w=1.8, fill=STONE1, color=STONE7):
    R(s, x, y, w, 0.27, fill=fill, line=STONE3)
    T(s, text, x+0.1, y+0.03, w-0.15, 0.22, size=8, bold=True, color=color)

def row_block(s, items, x, y, col_w, row_h=0.36, alt=True):
    for i, (label, value) in enumerate(items):
        bg = STONE1 if (alt and i % 2 == 0) else WHITE
        R(s, x, y + i*row_h, col_w[0]+col_w[1], row_h, fill=bg, line=STONE2)
        T(s, label, x+0.12, y+i*row_h+0.06, col_w[0]-0.15, row_h-0.08,
          size=9, bold=True, color=STONE7)
        T(s, value, x+col_w[0]+0.1, y+i*row_h+0.06, col_w[1]-0.15, row_h-0.08,
          size=9, color=DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, fill=DARK)
R(s, 0, 0, 13.33, 0.08, fill=RED)
R(s, 0, 7.42, 13.33, 0.08, fill=RED)

# Santander flamingo logo mark
R(s, 0.6, 1.6, 0.55, 0.55, fill=RED)
R(s, 0.6, 1.6, 0.55, 0.55, fill=RED)

T(s, "Project Record", 0.6, 2.4, 10, 0.55, size=13, color=STONE5, bold=False)
T(s, "Santander Business Banking", 0.6, 2.9, 11, 1.0, size=44, bold=True, color=WHITE)
T(s, "Digital Prototype", 0.6, 3.85, 11, 0.7, size=36, bold=False, color=LIGHTRED, italic=True)

R(s, 0.6, 4.72, 4.5, 0.02, fill=RED)

ML(s, [
    "Author: Alan Davidson · Business Banking Advisor",
    "Capacity: Self-initiated · Completed entirely out of hours · Own time",
    "Prototype: https://imscots-ui.github.io/santander/",
    "Date: June 2026",
], 0.6, 4.85, 7, 1.0, size=11, color=STONE3, leading=6)

# Stat boxes
stats = [("21", "Features"), ("12", "Workflows"), ("7", "Entity types"), ("8", "Frameworks")]
for i, (num, lbl) in enumerate(stats):
    bx = 8.2 + i * 1.28
    R(s, bx, 4.6, 1.15, 1.1, fill=RGBColor(0x2C, 0x27, 0x24))
    R(s, bx, 4.6, 1.15, 0.04, fill=RED)
    T(s, num, bx, 4.68, 1.15, 0.55, size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    T(s, lbl, bx, 5.2, 1.15, 0.3, size=8, color=STONE5, align=PP_ALIGN.CENTER)

slide_num(s, 1)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROJECT OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, fill=WARM)
section_header(s, "Project Overview", "A single-file React prototype — fully functional, browser-only, no backend")

# Left column — what it is
T(s, "WHAT WAS BUILT", 0.5, 1.15, 5.5, 0.28, size=9, bold=True, color=RED)
ML(s, [
    "App.jsx — ~5,300 lines · ~300KB source · Single React component",
    "5 main screens accessible from the bottom nav / sidebar",
    "12 step-based workflow wizards covering every key banking operation",
    "5 sheet overlays for quick-access tasks (notifications, receipts, PIN, open banking, voice)",
    "7 entity types — every compliance path forks by entity",
    "21 features across dashboard intelligence, paperless workflows, security, and compliance",
    "8 regulatory frameworks implemented end-to-end",
    "WCAG 2.1 AA accessibility — focus-visible on every interactive element",
    "No backend · No API calls · Runs from any browser as a single HTML file",
], 0.5, 1.45, 5.8, 3.8, size=10.5, color=DARK, leading=7)

# Right column — tech stack
T(s, "TECHNOLOGY", 7.0, 1.15, 5.8, 0.28, size=9, bold=True, color=RED)
tech = [
    ("Framework",        "React 18 · single-component architecture"),
    ("Styling",          "Tailwind CSS + inline brand CSS template literal"),
    ("Icons",            "Lucide React"),
    ("Build",            "Vite 5 + vite-plugin-singlefile → 832KB HTML"),
    ("Fonts",            "Fraunces (display) · Geist (body) · Geist Mono"),
    ("Deploy",           "GitHub Pages · gh-pages branch"),
    ("State",            "React useState / useMemo / useEffect — all at App() top"),
    ("Source control",   "Git · branch claude/add-claude-documentation-eaPc5"),
]
row_block(s, tech, 7.0, 1.45, [2.2, 3.9], row_h=0.38)

# Bottom strip
R(s, 0.5, 5.9, 12.33, 0.8, fill=DARK)
T(s, "Live prototype", 0.75, 5.98, 2.5, 0.28, size=9, bold=True, color=STONE3)
T(s, "https://imscots-ui.github.io/santander/", 3.1, 5.98, 6.5, 0.28, size=10, color=WHITE)
T(s, "832 KB · self-contained", 10.2, 5.98, 2.4, 0.28, size=9, color=STONE5, align=PP_ALIGN.RIGHT)
T(s, "All work self-initiated · out of hours · own time", 0.75, 6.36, 12, 0.25,
   size=8, color=STONE5)

advisor_stamp(s)
slide_num(s, 2)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — FIVE SCREENS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, fill=WARM)
section_header(s, "The Five Screens", "All accessible from the bottom navigation bar (mobile) or sidebar (desktop)")

screens = [
    ("Home Dashboard",
     "tab: 'home'",
     [
         "Balance hero card — accounts, available balance, mandate badge",
         "Business health score — live 0-100 circular SVG gauge across 5 factors",
         "Session anomaly alert — new device/location security warning",
         "13-week cash flow forecast — SVG bar chart with risk threshold",
         "Supplier risk radar — Companies House RAG status for top 5 counterparties",
         "Director command centre — all signatories with KYC status and pending counts",
         "10 action tiles — quick launch for every workflow and feature",
         "Cooling-off progress bars — live countdowns with cancel option",
         "Voice ID security card — enrolment status, cross-channel, SCA tiers",
         "Relationship Manager card — Priya Desai with escalation trigger",
     ]),
    ("Signature Queue",
     "tab: 'approve'",
     [
         "Pending mandates, closures, and bulk payments awaiting your signature",
         "Dual-authorisation: Face ID biometric sign or explicit reject",
         "Each item shows initiator, amount, expiry window, and rule type",
         "Badge count on nav tab — updates live as items are signed",
         "Stalled requests shown with RM escalation card when co-signer misses window",
     ]),
    ("Financial Statements",
     "tab: 'statements'",
     [
         "6 months of transactions — chronological and by category views",
         "Full-text counterparty search with instant results",
         "Counterparty detail sheet — spend breakdown, transaction history, method split",
         "PDF / CSV / Excel export with simulated download feedback",
         "Transaction method filter: all, bank transfer, direct debit, standing order, card",
         "Outgoing / incoming colour-coded with tabular-num amounts",
     ]),
    ("Making Tax Digital",
     "tab: 'mtd'",
     [
         "HMRC VAT and ITSA obligations with deadline countdown badges",
         "Quarterly submission wizard — 4 steps: categorise, review, declare, submit",
         "Auto-categorisation with confidence scoring for each transaction",
         "VAT100 form view with all 9 boxes calculated",
         "YTD income/expense totals and projected tax liability",
         "ITSA readiness indicator with quarters ahead of schedule",
         "Voice memo integration — tap-to-record expense auto-fills to ledger",
         "Receipt scan — OCR categorise and match to open transactions",
     ]),
    ("Audit Trail",
     "tab: 'audit'",
     [
         "FCA SYSC 9 compliant — 7-year retention, append-only",
         "Every action: workflow steps, sign events, rejections, system alerts",
         "Each entry: timestamp, actor name, action type, account/reference",
         "Session anomaly events logged with device/location metadata",
         "Filter by event type or actor",
         "Export for compliance reporting",
     ]),
]

# 3+2 layout: 3 panels in row 0, 2 centred in row 1
col_w  = 4.04
gap_x  = 0.19
step_x = col_w + gap_x
ph     = 2.68
step_y = ph + 0.22

# Row 0 — first 3 screens
for ci, (name, state, pts) in enumerate(screens[:3]):
    x = 0.42 + ci * step_x
    y = 1.15
    R(s, x, y, col_w, ph, fill=WHITE, line=STONE3)
    R(s, x, y, col_w, 0.06, fill=RED)
    T(s, name,  x+0.15, y+0.10, col_w-0.30, 0.30, size=11, bold=True,  color=DARK)
    T(s, state, x+0.15, y+0.40, col_w-0.30, 0.20, size=8,  italic=True, color=STONE5)
    ML(s, ["· " + p for p in pts], x+0.15, y+0.62, col_w-0.28, ph-0.70, size=8.5, color=STONE7, leading=3)

# Row 1 — last 2 screens centred
x_start = (13.33 - 2 * col_w - gap_x) / 2
for ci, (name, state, pts) in enumerate(screens[3:]):
    x = x_start + ci * step_x
    y = 1.15 + step_y
    R(s, x, y, col_w, ph, fill=WHITE, line=STONE3)
    R(s, x, y, col_w, 0.06, fill=RED)
    T(s, name,  x+0.15, y+0.10, col_w-0.30, 0.30, size=11, bold=True,  color=DARK)
    T(s, state, x+0.15, y+0.40, col_w-0.30, 0.20, size=8,  italic=True, color=STONE5)
    ML(s, ["· " + p for p in pts], x+0.15, y+0.62, col_w-0.28, ph-0.70, size=8.5, color=STONE7, leading=3)

advisor_stamp(s)
slide_num(s, 3)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ELEVEN WORKFLOWS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, fill=WARM)
section_header(s, "Twelve Workflow Wizards", "Step-based overlays — each is a full multi-step process rendered on top of the main app")

workflows = [
    ("Account Closure",
     "4 steps",
     "Select accounts → destination CoP check → credit check → sign. 24h cooling-off on Any-1 accounts. Mandate enforcement throughout."),
    ("Partner Unreachable Escalation",
     "4 steps",
     "4 regulatory branches: health/welfare, no contact, deceased/probate, dispute. Evidence upload, specialist team referral, account restriction. FCA Consumer Duty PS22/9. Case ref SVC-2026-4471."),
    ("Mandate Changes",
     "6 steps",
     "Add / remove / update signatories or change signing rule. Full KYC/KYB: GOV.UK One Login List 1, address List 2, trading List 3. Board minutes for clubs/charities. 7 entity type variations."),
    ("Bulk Payments / Wages",
     "4 steps",
     "Payee book, CoP verification per payee, CSV import, one-off or monthly scheduling. Mandate enforcement: Any-1/2/All per source account."),
    ("Business Details Update",
     "3 steps",
     "Address, trading name, phone, email. Proof of address upload. Companies House sync notification. RM escalation for partnership renames."),
    ("Dormant Reactivation",
     "1 step",
     "Select dormant account and request reactivation. FSCS notice displayed. 12-month inactivity threshold enforced."),
    ("MTD VAT Submission",
     "4 steps",
     "Auto-categorise transactions (confidence score per item) → review VAT100 form → declare accuracy → HMRC API submit. Returns receipt number."),
    ("Personal/Business Unlink",
     "4 steps",
     "Scope confirmation → data separation notice → all-channels toggle (CRM ref REL-2026-0291, 2 working days) + postal toggle → signed declaration. Removes personal data from JSX render tree — not CSS hide."),
    ("Credit Ring-Fence",
     "2 steps",
     "Risk overview with GDPR Art.5(1)(c) legal basis → signed formal instruction. Personal account data permanently excluded from business credit decisioning. State persists across sessions."),
    ("Pre-approved Business Lending",
     "3 steps",
     "£45,000 offer with live monthly repayment calculator (12/24/36 months) → CCA 1974 terms and 14-day cooling-off rights → biometric confirm and draw-down. lendingCompleted state persists."),
    ("International FX Payment",
     "3 steps",
     "Amount + currency + IBAN + beneficiary → live rate + FCA fee disclosure (Consumer Rights Act 2015) → biometric confirm. 5 currencies: EUR/USD/CHF/AUD/CAD. MLR 2017 screening flag ≥£50k. SWIFT ref to audit trail."),
    ("Complaint Handling",
     "4 steps",
     "Intake → escalation triage → denial → case outcome. FCA DISP-compliant. Eligible complainant check (DISP 2.7), escalation flag selection, denial reason, optional goodwill gesture. Triggers from Home screen Log complaint tile."),
]

pw = 3.9
ph = 1.38
for i, (name, steps, desc) in enumerate(workflows):
    col = i % 3
    row = i // 3
    x = 0.42 + col * (pw + 0.19)
    y = 1.12 + row * (ph + 0.16)
    R(s, x, y, pw, ph, fill=WHITE, line=STONE2)
    R(s, x, y, 0.055, ph, fill=RED)
    T(s, name, x+0.18, y+0.1, pw-0.28, 0.28, size=9.5, bold=True, color=DARK)
    T(s, steps, x+0.18, y+0.36, 1.2, 0.2, size=7.5, bold=True, color=RED)
    ML(s, [desc], x+0.18, y+0.55, pw-0.3, 0.78, size=8, color=STONE5, leading=3)

advisor_stamp(s)
slide_num(s, 4)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ENTITY TYPES & MANDATE RULES
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, fill=WARM)
section_header(s, "Seven Entity Types & Mandate Rules",
               "Every compliance path diverges by entity — KYC docs, principals, authorisation rules, board minute requirements")

# Table header
cols = [("Entity Type", 2.3), ("Mandate Default", 1.7), ("Principal Title", 1.8),
        ("Companies House", 1.85), ("Board Minutes", 2.25), ("Key Divergence", 2.6)]
hx = 0.42
hy = 1.15
R(s, hx, hy, sum(c[1] for c in cols) + 0.05, 0.36, fill=DARK)
cx = hx + 0.1
for label, cw in cols:
    T(s, label, cx, hy+0.06, cw-0.1, 0.26, size=8, bold=True, color=STONE3)
    cx += cw

# Rows
entities = [
    ("Sole Trader",       "Any-1",   "Sole trader",      "No",                   "No",                    "Single principal — no co-signer ever required"),
    ("Partnership",       "Any-1/2", "Partner",          "No",                   "No (RM on rename)",     "Partner-unreachable path active; Any-2 if >2 partners"),
    ("Limited Company",   "Any-2",   "Director",         "Yes — Companies No.",  "No",                    "CH reg verified; director register linked to mandate"),
    ("LLP",               "Any-2",   "Member",           "Yes — LLP No.",        "No",                    "LLP number shown; member terminology throughout"),
    ("Registered Charity","All",     "Trustee",          "No — Charity Comm.",   "Yes — all changes",     "All trustees must sign; Charity Commission ref shown"),
    ("Club",              "All",     "Committee member", "No",                   "Yes — all changes",     "Full committee approval for every mandate/detail change"),
    ("Society",           "All",     "Committee member", "No",                   "Yes — all changes",     "Same as club; society terminology throughout"),
]
col_widths = [c[1] for c in cols]

for ri, row in enumerate(entities):
    ry = hy + 0.36 + ri * 0.44
    bg = STONE1 if ri % 2 == 0 else WHITE
    R(s, hx, ry, sum(col_widths) + 0.05, 0.44, fill=bg, line=STONE2)
    cx2 = hx + 0.1
    rule_color = RED if row[1] == "All" else (AMBER if "2" in row[1] else EMERALD)
    for ci, (val, cw) in enumerate(zip(row, col_widths)):
        c = rule_color if ci == 1 else DARK
        T(s, val, cx2, ry+0.09, cw-0.12, 0.3, size=8.5, color=c, bold=(ci == 1))
        cx2 += cw

# Legend
legend_y = hy + 0.36 + 7*0.44 + 0.1
T(s, "Any-1 = any single signatory authorises", 0.42, legend_y, 3.5, 0.25, size=8, color=EMERALD)
T(s, "Any-2 = any two signatories must co-sign", 3.9, legend_y, 3.5, 0.25, size=8, color=AMBER)
T(s, "All = every signatory must sign", 7.35, legend_y, 3.0, 0.25, size=8, color=RED)
T(s, "getMandateFor() picks the strictest rule across selected accounts", 0.42, legend_y+0.25, 8, 0.22, size=8, color=STONE5, italic=True)

advisor_stamp(s)
slide_num(s, 5)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SECURITY & COMPLIANCE (PART 1)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, fill=WARM)
section_header(s, "Security & Compliance — Part 1 of 2",
               "Every feature tied to a specific regulation or standard — no vague compliance, only implemented rules")

compliance_1 = [
    ("24h cooling-off",
     "FCA BCOBS 4A",
     "Working-hours aware (bank holidays skipped). Live progress bar refreshes every 30s. Cancellable. Triggered automatically on Any-1 accounts."),
    ("Dual-authorisation",
     "PSR 2017",
     "Any-1 / Any-2 / All per account. getMandateFor() picks strictest rule. Co-signer notified. 48-hour sign window. Stalled escalation to RM."),
    ("KYC / KYB",
     "MLR 2017 / JMLSG",
     "GOV.UK One Login List 1 photo ID. Address List 2. Trading address List 3. Sanctions, PEP, adverse media, visa status checks. Document expiry tracking."),
    ("Audit trail",
     "FCA SYSC 9",
     "7-year retention. Append-only — no edits or deletes. Every workflow step, sign event, rejection, and system alert. Actor name + timestamp on every entry."),
    ("MTD compliance",
     "HMRC MTD VAT / ITSA",
     "Digital links maintained. 6-year records. HMRC MTD VAT API v1.0 recognised software. VAT100 direct submission. Receipt scan closes manual-entry gap."),
    ("Confirmation of Payee",
     "PSR SI 2019/1215",
     "Pre-payment account name verification. Mismatch warning shown before proceeding. Explicit user override required — cannot be skipped silently."),
    ("Accessibility",
     "WCAG 2.1 AA",
     "focus-visible:ring on every interactive element. Contrast ≥4.5:1 throughout. Tabular-nums on all monetary amounts. ARIA labels where needed. No bare focus:outline-none."),
    ("Partner unreachable",
     "FCA Consumer Duty PS22/9 / BCOBS",
     "Documented contact attempts. Evidence upload. 4-path regulatory branching: health/welfare, no contact, deceased/probate, dispute. Specialist team referral. Account restriction."),
    ("Personal/business separation (app)",
     "GDPR Art.5(1)(c)",
     "personalLinked state removes personal section from JSX render tree — not CSS visibility, not display:none. Actual data removal, not visual hiding."),
    ("Personal/business separation (call centre)",
     "GDPR Art.5(1)(c)",
     "unlinkAllChannels toggle triggers CRM back-office update ref REL-2026-0291. CLI lookup excludes personal accounts. Explicit 2-working-day SLA disclosed."),
]

pw2 = 5.9
ph2 = 1.16
for i, (name, reg, desc) in enumerate(compliance_1):
    col = i % 2
    row = i // 2
    x = 0.42 + col * (pw2 + 0.25)
    y = 1.12 + row * (ph2 + 0.1)
    R(s, x, y, pw2, ph2, fill=WHITE, line=STONE2)
    R(s, x, y, pw2, 0.055, fill=RED)
    T(s, name, x+0.15, y+0.1, pw2*0.55, 0.3, size=9.5, bold=True, color=DARK)
    T(s, reg, x+pw2*0.58, y+0.12, pw2*0.4, 0.25, size=7.5, color=RED, bold=True)
    ML(s, [desc], x+0.15, y+0.42, pw2-0.25, 0.68, size=8, color=STONE5, leading=3)

advisor_stamp(s)
slide_num(s, 6)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — SECURITY & COMPLIANCE (PART 2)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, fill=WARM)
section_header(s, "Security & Compliance — Part 2 of 2",
               "Advanced features: biometric authentication, intelligent session monitoring, and real-time notification")

compliance_2 = [
    ("Postal statement separation",
     "GDPR Art.5(1)(c)",
     "unlinkPostal toggle — separate delivery preference per channel. Personal statements excluded from business delivery address. Explicit per-channel granularity."),
    ("Credit ring-fence",
     "GDPR Art.5(1)(c) / ICO LIA",
     "Formal written instruction persisted in state. Personal account history permanently excluded from business loan and overdraft underwriting. State survives session."),
    ("PSD2 consent audit",
     "PSD2 RTS Art.29",
     "OBSheet shows all third-party consents with scope, purpose, and expiry date. Flags personal data exposure (Funding Circle). Revocation within 90 seconds per RTS requirement."),
    ("Receipt scan → MTD",
     "HMRC MTD digital links",
     "Simulated OCR: scan receipt → extract merchant, amount, VAT rate, category, reference → match to open transaction. Closes the manual-entry gap in the digital links chain."),
    ("Voice ID biometric",
     "PSD2 RTS Art.4(30) · GDPR Art.9 · NIST SP 800-63B",
     "3-phrase voice enrolment with GDPR Art.9 special-category consent gate. Active across app, phone banking, and video call channels. Anti-spoofing: liveness + voice clone detection. voiceIdEnrolled persists."),
    ("SCA step-up matrix",
     "PSD2 RTS Art.97 · FCA PS19/4",
     "6-tier dynamic authentication: device passkey → biometric fingerprint → Voice ID → Voice ID + co-sig → Voice ID + cooling-off → all three combined. Interactive table in VoiceIdSheet."),
    ("Session anomaly detection",
     "FCA SYSC 10A · PSD2 RTS Art.2",
     "Dismissible amber banner on HomeScreen. New device/location login alert. Two CTAs: 'That was me' (dismisses and confirms) or 'Review activity' (navigates to Audit tab). Logged to audit trail."),
    ("Supplier risk radar",
     "Companies House filing obligations",
     "5 key suppliers shown with CH registration, last filed date, days overdue, annual spend, and RAG risk badge. Red = >180 days overdue, Amber = 90–180 days, Green = current. Links to CH API."),
    ("Live notification bell",
     "FCA SYSC 10A · PSD2 RTS Art.2",
     "Bell icon (desktop header + mobile header) opens NotificationsSheet. Shows: session anomaly alert with CTAs, unsigned pending approvals with Sign-tab navigation, cooling-off progress bars, stalled co-signer escalations. Badge dot while anomaly or approvals exist."),
]

pw3 = 5.9
ph3 = 1.12
for i, (name, reg, desc) in enumerate(compliance_2[:8]):   # first 8 in 2-col grid
    col = i % 2
    row = i // 2
    x = 0.42 + col * (pw3 + 0.25)
    y = 1.12 + row * (ph3 + 0.1)
    R(s, x, y, pw3, ph3, fill=WHITE, line=STONE2)
    R(s, x, y, pw3, 0.055, fill=RED)
    T(s, name, x+0.15, y+0.1, pw3*0.55, 0.3, size=9.5, bold=True, color=DARK)
    T(s, reg, x+pw3*0.58, y+0.12, pw3*0.4, 0.25, size=7, color=RED, bold=True)
    ML(s, [desc], x+0.15, y+0.42, pw3-0.25, 0.65, size=8, color=STONE5, leading=3)

# 9th item — centred below the 4-row grid
name9, reg9, desc9 = compliance_2[8]
cx9 = (13.33 - pw3) / 2          # truly centred on the slide
cy9 = 1.12 + 4 * (ph3 + 0.1)
R(s, cx9, cy9, pw3, ph3, fill=WHITE, line=STONE2)
R(s, cx9, cy9, pw3, 0.055, fill=RED)
T(s, name9, cx9+0.15, cy9+0.1, pw3*0.55, 0.3, size=9.5, bold=True, color=DARK)
T(s, reg9, cx9+pw3*0.58, cy9+0.12, pw3*0.4, 0.25, size=7, color=RED, bold=True)
ML(s, [desc9], cx9+0.15, cy9+0.42, pw3-0.25, 0.65, size=8, color=STONE5, leading=3)

advisor_stamp(s)
slide_num(s, 7)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — HOME SCREEN INTELLIGENCE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, fill=WARM)
section_header(s, "Home Screen Intelligence",
               "Eight features that go beyond banking — proactive intelligence surfaced directly on the dashboard")

intel = [
    ("Pre-approved Lending Offer",
     "Hero card — £45,000 offer surfaced proactively on home screen. Based on 6-month transaction history, reserve balance, and payment record. 3-step draw-down in-app. Replaced by active loan summary after draw-down. lendingCompleted persists."),
    ("13-Week Cash Flow Forecast",
     "SVG bar chart (forecastWeeks useMemo) — 13 weekly balance projections computed from known direct debits, payroll cycles, VAT commitments, and income patterns. Amber warning line at £80k. Actionable nudge when threshold breached."),
    ("Business Health Score",
     "Circular SVG gauge (0–100) computed live from 5 weighted factors: liquidity 0–20, tax compliance 0–20, cash flow headroom 0–20, payroll ratio 0–20, mandate health 0–20. Each factor shown as a colour-coded progress bar. Grade A–D. healthScore useMemo refreshes on accounts/payees/forecast changes."),
    ("Supplier Risk Radar",
     "5 key counterparties shown with Companies House registration number, last filing date, days overdue, and annual spend. Red (>180d overdue), Amber (90–180d), Green (current). Critical badge if any supplier in red. SUPPLIER_RISK static data. Links to CH API."),
    ("Director Command Centre",
     "2×2 governance grid showing all 4 signatories. Per-director: initials avatar, role title, KYC verification status, last-active timestamp, pending action count (red badge). One tap navigates to full Signature Queue."),
    ("Voice Memo → MTD",
     "Tap-to-record expense entry — no typing required. 1.8s simulated transcription extracts merchant, amount, VAT rate, category, and reference. One-tap confirm posts to MTD ledger and audit trail. voiceMemoAdded Set tracks confirmed entries per session."),
    ("Smart Payment Sequencer",
     "30-day balance outlook SVG bar chart with daily projections and £80k risk threshold line. 7 scheduled payments with type colour-coding (tax/fixed/supplier). Optimise button reschedules discretionary payments to maximise minimum balance. Locked payments (tax, rent) never moved."),
    ("Voice ID Biometric Setup",
     "3-phrase voice enrolment with GDPR Art.9 special-category consent gate displayed on first entry. Cross-channel status: app, phone banking, video call. Anti-spoofing badge. SCA step-up matrix (6 tiers, PSD2 RTS Art.97). voiceIdEnrolled persists."),
]

pw4 = 5.9
ph4 = 1.28
for i, (name, desc) in enumerate(intel):
    col = i % 2
    row = i // 2
    x = 0.42 + col * (pw4 + 0.25)
    y = 1.12 + row * (ph4 + 0.12)
    R(s, x, y, pw4, ph4, fill=WHITE, line=STONE2)
    R(s, x, y, 0.055, ph4, fill=RED)
    T(s, name, x+0.18, y+0.1, pw4-0.3, 0.3, size=10, bold=True, color=DARK)
    ML(s, [desc], x+0.18, y+0.44, pw4-0.3, 0.8, size=8.5, color=STONE5, leading=3)

advisor_stamp(s)
slide_num(s, 8)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — SHIPS COMPANY — AGENT ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, fill=WARM)
section_header(s, "Ships Company — Agent Architecture",
               "Six slash-command officers built into Claude Code — each specialist enforces a discipline before every commit")

officers = [
    ("/xo",         "Commander (XO)",           RED,    "Architecture quality · Hook discipline · Single-file mandate · Pattern enforcement · closeWorkflow() invariant"),
    ("/security",   "Lt Cdr Security",           RGBColor(0xDC, 0x26, 0x26),
                                                         "XSS scan · Secrets scan · focus:outline-none audit · PII check · dangerouslySetInnerHTML check"),
    ("/engineer",   "Lt Cdr Engineer",           RGBColor(0x1D, 0x4E, 0xD8),
                                                         "Bundle analysis · React performance · key props check · tick interval guard · Vite config · CSS delivery"),
    ("/bosun",      "Lieutenant Bosun",           AMBER,  "Spacing scale violations · Colour token violations (stone not gray) · CTA hierarchy · Shadow system · Empty states"),
    ("/signals",    "Lieutenant Signals",         EMERALD,"State signal integrity · useMemo dependency audit · Direct mutation safety · Cooling system · API contracts"),
    ("/ship-ready", "All stations",               DARK,   "7-station pre-commit checklist — GREEN/AMBER/RED per station. A single RED means the ship stands down and fixes first."),
]

for i, (cmd, rank, col, duties) in enumerate(officers):
    row = i // 2
    ci2 = i % 2
    x = 0.42 + ci2 * 6.47
    y = 1.15 + row * 1.68
    R(s, x, y, 6.25, 1.52, fill=WHITE, line=STONE2)
    R(s, x, y, 0.06, 1.52, fill=col)
    R(s, x+0.06, y, 1.5, 1.52, fill=STONE1)
    T(s, cmd, x+0.22, y+0.22, 1.28, 0.45, size=18, bold=True, color=col)
    T(s, rank, x+0.22, y+0.7, 1.3, 0.45, size=8, color=STONE5)
    T(s, "Duties:", x+1.72, y+0.12, 0.6, 0.25, size=8, bold=True, color=STONE5)
    ML(s, ["· " + d for d in duties.split(" · ")], x+1.72, y+0.38, 4.4, 1.1, size=8.5, color=DARK, leading=4)

# Settings strip
sy = 1.15 + 3 * 1.68 + 0.1
R(s, 0.42, sy, 12.49, 0.54, fill=DARK)
T(s, "Settings — .claude/settings.json", 0.62, sy+0.08, 3.5, 0.22, size=9, bold=True, color=STONE3)
T(s, "Allow: npm run dev/build/preview · git push/add/commit/status/diff/log/branch · python3 · node · npx vite",
   0.62, sy+0.28, 7.0, 0.2, size=8, color=STONE5)
T(s, "Deny: rm -rf · git push --force · chmod 777 · curl|bash · eval · exec",
   8.0, sy+0.18, 5.0, 0.28, size=8, color=RGBColor(0xFC, 0xA5, 0xA5))

advisor_stamp(s)
slide_num(s, 9)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — REFERENCE LIBRARY
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, fill=WARM)
section_header(s, "Reference Library — 1701-uniform/REFERENCE.md",
               "23,112 lines · 93 sections · 102 books & regulatory documents · §91–93 system design, PowerPoint, materials")

# ── Stats banner ──────────────────────────────────────────────────────────────
for i4, (val, lbl) in enumerate([
    ("23,112", "lines"),
    ("93",     "sections"),
    ("102",    "books & docs"),
    ("70+",    "quick answers"),
]):
    bx4 = 0.42 + i4 * 3.22
    R(s, bx4, 1.12, 3.05, 0.58, fill=RED)
    T(s, val, bx4, 1.14, 3.05, 0.36, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    T(s, lbl, bx4, 1.50, 3.05, 0.20, size=8, color=LIGHTRED, align=PP_ALIGN.CENTER)

# ── Domain blocks — two columns of 5 ─────────────────────────────────────────
domains = [
    [  # Left column — x=0.42
        (RED,     "Core Development",                   "§1–22, §91",
         "Python · SQL · FastAPI · Docker · Git · React/JS · JWT · HTTP · Linux · Security · Pydantic · async · pytest · System Design Interview (King)"),
        (RED,     "AI & Prompt Engineering",            "§20, §23–25, §58–65, §72–74",
         "Claude · AI Agents · Copilot · ChatGPT · AI Content Creation · Surveillance AI · Dynamics 365 AI · 600+ prompts"),
        (RED,     "Banking · Payments · Regulation",    "§35–41, §66–70, §84, §86–88",
         "PSD2 · FPS/BACS/CHAPS/SWIFT · PSR · FCA Consumer Duty · AML · Ring-Fencing · SCA · BCOBS · FCA Compliance · UK Resolution · Post-Trade"),
        (RED,     "Prudential & Market Regulation",     "§36–40, §87–88",
         "PRA/PRC · Basel 3.1 · SS19/13 Resolution Planning · MREL · CNRF · Deposit Aggregators · LEI/SSI · KYC Passporting · Post-Trade Task Force"),
        (RED,     "Data Protection",                    "§75, §77",
         "GDPR/DPDI/Marketing AI (Scheuing, Kogan Page 2024) · GDPR for Startups (Martin, Edward Elgar 2023)"),
    ],
    [  # Right column — x=6.68
        (DARK,    "UI · Frontend · Design",             "§11, §26, §42, §50–56, §78–81",
         "React/JS patterns · Refactoring UI · Tailwind CSS v3 (×3 books) · CSS Animation · Vite/TS · React 19 · React Hooks (Larsen) · UX · Photoshop · KiCad"),
        (DARK,    "Microsoft Ecosystem",                "§27–33, §64–65, §76, §92",
         "Power Teams · Power BI (×2) · PowerPoint (×3 — incl. 365 Pro/Lemmings) · SharePoint · Dynamics 365 BC AL · Azure AI Engineer AI-102"),
        (DARK,    "Employment & Equality Law",          "§34",
         "Equality Act 2010 · Nine protected characteristics · WCAG accessibility legal duty · Reasonable adjustments"),
        (DARK,    "Electronics · Hardware · Engineering Design", "§43–49, §54, §57, §89, §93",
         "MakerSpace · Electronics · Mechatronics · Digital Logic · Arduino/IoT · Embedded Linux · ISO GD&T (Green) · CAMSS Materials Selection (NMAB-467)"),
        (EMERALD, "Digital Banking Transformation",      "§82, §83, §85",
         "DBS World's Best Bank (Speculand) · The Autonomous Bank AI/ML (Dhillon) · TMRW Digital Bank Build (Khoo/UOB) · Quick Lookup Decision Index"),
    ],
]

col_xs = [0.42, 6.68]
for ci5, col_domains in enumerate(domains):
    cx5 = col_xs[ci5]
    cw5 = 6.05
    for di5, (hdr_col, domain, refs, desc) in enumerate(col_domains):
        dy5 = 1.82 + di5 * 0.98
        R(s, cx5, dy5, cw5, 0.26, fill=hdr_col)
        T(s, domain, cx5+0.10, dy5+0.03, cw5-0.72, 0.22, size=8.5, bold=True, color=WHITE)
        T(s, refs, cx5+cw5-0.60, dy5+0.04, 0.58, 0.20, size=7, bold=True,
          color=LIGHTRED, align=PP_ALIGN.RIGHT)
        bg5 = STONE1 if di5 % 2 == 0 else WHITE
        R(s, cx5, dy5+0.26, cw5, 0.68, fill=bg5, line=STONE2)
        T(s, desc, cx5+0.10, dy5+0.30, cw5-0.20, 0.60, size=7.5, color=STONE7)

advisor_stamp(s)
slide_num(s, 10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — DEPLOYMENT & BUILD
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, fill=WARM)
section_header(s, "Deployment & Build Configuration",
               "Single-file HTML distributed via GitHub Pages — no server, no CDN dependency, runs from a USB stick")

# GitHub Pages box — height 3.1 to contain 6 data rows without overflow
R(s, 0.42, 1.12, 7.8, 3.1, fill=WHITE, line=STONE3)
R(s, 0.42, 1.12, 7.8, 0.06, fill=RED)
T(s, "GitHub Pages", 0.62, 1.22, 4, 0.35, size=13, bold=True, color=DARK)
T(s, "Live · imscots-ui.github.io/santander/", 0.62, 1.58, 7.4, 0.28, size=11, color=RED)
pages_info = [
    ("Branch",    "gh-pages"),
    ("Contents",  "index.html (832KB) + .nojekyll + vercel.json"),
    ("Build",     "Vite 5 + vite-plugin-singlefile — all JS/CSS/fonts inlined"),
    ("Gzip",      "181.99 KB — fast first load even on 3G"),
    ("Key fix",   "DOMContentLoaded listener — singlefile inlines JS into <head> as plain <script> (no defer)"),
    ("Mirror",    "Vercel — santander-g671k7hgb-imscots-uis-projects.vercel.app"),
]
row_block(s, pages_info, 0.62, 1.95, [1.7, 5.8], row_h=0.33)

# Vite config box — height 3.1 to match GitHub Pages
R(s, 8.42, 1.12, 4.5, 3.1, fill=WHITE, line=STONE3)
R(s, 8.42, 1.12, 4.5, 0.06, fill=RED)
T(s, "vite.config.js", 8.62, 1.22, 4, 0.3, size=12, bold=True, color=DARK)
vite_info = [
    ("Plugin",   "vite-plugin-singlefile"),
    ("Target",   "esnext"),
    ("Output",   "dist/index.html"),
    ("Size",     "832 KB · 182 KB gzip"),
    ("Removed",  "removeViteModuleLoader: true"),
    ("Why",      "Strips module type → script exec order bug"),
]
row_block(s, vite_info, 8.62, 1.57, [1.4, 2.83], row_h=0.33)

# npm commands box — starts at y=4.42 (after row 1 boxes end at 4.22 + 0.2 gap)
R(s, 0.42, 4.42, 5.8, 2.58, fill=WHITE, line=STONE3)
R(s, 0.42, 4.42, 5.8, 0.06, fill=RED)
T(s, "npm Commands", 0.62, 4.52, 5.4, 0.3, size=12, bold=True, color=DARK)
cmds = [
    ("npm install",              "Install dependencies — first time only"),
    ("npm run dev",              "Dev server at localhost:5173 with hot reload"),
    ("npm run build",            "Production build to dist/ — 832KB HTML"),
    ("npm run preview",          "Serve dist/ locally for final testing"),
    ("python3 build_deck.py",           "Regenerate pitch deck (16 slides)"),
    ("python3 build_architecture.py",   "Regenerate architecture deck (10 slides)"),
    ("python3 build_manifest.py",       "Regenerate this project record (14 slides)"),
    ("python3 build_complaint_guide.py","Regenerate Santander_Complaint_Guide.docx"),
    ("python3 complaint_letter_writer.py", "Interactive FCA DISP letter writer (Evelyn/Angus)"),
]
for i2, (cmd, desc) in enumerate(cmds):
    ry2 = 4.84 + i2 * 0.30
    bg2 = STONE1 if i2 % 2 == 0 else WHITE
    R(s, 0.62, ry2, 5.4, 0.30, fill=bg2, line=STONE2)
    T(s, cmd, 0.72, ry2+0.04, 2.1, 0.23, size=8, bold=True, color=DARK)
    T(s, desc, 2.85, ry2+0.04, 3.1, 0.23, size=8, color=STONE5)

# CLAUDE.md box — same y as npm Commands
R(s, 6.42, 4.42, 6.5, 2.58, fill=WHITE, line=STONE3)
R(s, 6.42, 4.42, 6.5, 0.06, fill=RED)
T(s, "CLAUDE.md — Captain's Standing Orders", 6.62, 4.52, 6.1, 0.3, size=11, bold=True, color=DARK)
claude_items = [
    "Ships company command table — 6 officers with duties",
    "5 Standing Security Orders: no bare outline-none, no dangerouslySetInnerHTML on dynamic content, no secrets in source, no git push --force, run /ship-ready before every push",
    "Bosun's Law — spacing scale (4·8·12·16·24·32·48·64·96·128px), stone-* only, one primary CTA per view, num-tab on all monetary amounts",
    "closeWorkflow() invariant — all new state must be reset here",
    "Architecture constraints — no hooks inside closures, keep everything in App.jsx, single-file is intentional",
]
ML(s, ["· " + item for item in claude_items], 6.62, 4.88, 6.1, 1.95, size=8.5, color=DARK, leading=5)

advisor_stamp(s)
slide_num(s, 11)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — PRESENTATION DELIVERABLES
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, fill=WARM)
section_header(s, "Presentation Deliverables",
               "Three PowerPoint decks covering every audience — executives, technologists, and project stakeholders")

decks = [
    ("Santander_Digital_Banking_Future.pptx",
     "Pitch Deck · 16 slides",
     RED,
     [
         ("Slide 1",  "Title + Business Banking Advisor stamp"),
         ("Slide 2",  "Executive summary + 8 headline statistics"),
         ("Slide 3",  "The problem — 4 pain points"),
         ("Slide 4",  "What we built — 21 features (4-col grid)"),
         ("Slide 5",  "Privacy controls — 3-layer data separation"),
         ("Slide 6",  "Advanced features — lending, FX, forecast, receipt scan"),
         ("Slide 7",  "Intelligence & security — 7 panels (health score, radar, director, Voice ID, sequencer, voice memo, bell)"),
         ("Slide 8",  "Paperless workflows deep-dive"),
         ("Slide 9",  "Making Tax Digital deep-dive"),
         ("Slide 10", "Security & compliance (4 frameworks)"),
         ("Slide 11", "Financial intelligence features"),
         ("Slide 12", "Technical architecture (current + production)"),
         ("Slide 13", "Business case — £137M annualised opportunity"),
         ("Slide 14", "Roadmap — 4 phases Q3 2026 → Q2 2027"),
         ("Slide 15", "Design system & accessibility"),
         ("Slide 16", "Next steps — 3 decisions for senior management"),
     ]),
    ("Santander_Architecture.pptx",
     "Architecture Deck · 10 slides",
     RGBColor(0x1D, 0x4E, 0xD8),
     [
         ("Slide 1",  "Cover"),
         ("Slide 2",  "Component map — current prototype"),
         ("Slide 3",  "Data & state architecture"),
         ("Slide 4",  "Production system architecture (4 tiers)"),
         ("Slide 5",  "Security architecture (4 domains)"),
         ("Slide 6",  "Regulatory & compliance framework"),
         ("Slide 7",  "Privacy & account separation architecture"),
         ("Slide 8",  "Entity & mandate logic (7 types × 3 rules)"),
         ("Slide 9",  "API integration contracts (4 APIs)"),
         ("Slide 10", "Design system tokens"),
     ]),
    ("Santander_Project_Record.pptx",
     "Project Record · 14 slides (this deck)",
     STONE7,
     [
         ("Slide 1",  "Cover"),
         ("Slide 2",  "Project overview — stats, tech stack, deployment"),
         ("Slide 3",  "Five screens"),
         ("Slide 4",  "Twelve workflows"),
         ("Slide 5",  "Entity types & mandate rules"),
         ("Slide 6",  "Security & compliance — part 1"),
         ("Slide 7",  "Security & compliance — part 2"),
         ("Slide 8",  "Home screen intelligence"),
         ("Slide 9",  "Ships company agent architecture"),
         ("Slide 10", "Reference library (93 sections · 102 books)"),
         ("Slide 11", "Deployment & build"),
         ("Slide 12", "Presentation deliverables"),
         ("Slide 13", "Complaint handling tools (Evelyn/Angus)"),
         ("Slide 14", "Business case summary"),
     ]),
]

for ci4, (filename, subtitle, acol, slides) in enumerate(decks):
    x4 = 0.42 + ci4 * 4.32
    R(s, x4, 1.12, 4.1, 5.58, fill=WHITE, line=STONE3)
    R(s, x4, 1.12, 4.1, 0.07, fill=acol)
    T(s, subtitle, x4+0.18, 1.22, 3.8, 0.3, size=10, bold=True, color=acol)
    T(s, filename, x4+0.18, 1.52, 3.8, 0.3, size=8.5, color=STONE5, italic=True)
    R(s, x4+0.18, 1.86, 3.7, 0.016, fill=STONE3)
    for si, (num, desc) in enumerate(slides):
        sy2 = 1.92 + si * 0.29
        if sy2 > 6.5:
            break
        R(s, x4+0.18, sy2, 0.52, 0.26, fill=STONE1)
        T(s, num, x4+0.2, sy2+0.03, 0.5, 0.22, size=7.5, bold=True, color=acol)
        T(s, desc, x4+0.75, sy2+0.03, 3.3, 0.22, size=8, color=DARK)

T(s, "All decks rebuilt via python3 build_*.py — Santander brand colours throughout · Advisor credit on every slide",
  0.42, 6.82, 12.49, 0.22, size=8, color=STONE5, italic=True)

advisor_stamp(s)
slide_num(s, 12)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — COMPLAINT HANDLING TOOLS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, fill=WARM)
section_header(s, "Complaint Handling Tools — Evelyn Workflow (Angus)",
               "FCA DISP-compliant complaint record, letter writing, and Copilot agent — built alongside the in-app workflow")

# ── Four tool panels ──────────────────────────────────────────────────────────
tools = [
    ("complaint_letter_writer.py",
     "Interactive Letter Writer",
     RED,
     [
         "Trigger-phrase driven — Evelyn's exact workflow",
         "Stage 1 record: 5 verbatim questions",
         "Escalation record: 4 verbatim questions",
         "Uphold letter: Sorry → Cause → Impact → Fix → Payment",
         "Declined letter: Sorry → Not our fault → Reason → Explanation",
         "Escalation: resolution changed / no change",
         "Dual mode: complaint record + letter in one doc",
         "Output: [REF]_[SURNAME]_[TRIGGER]_[DATE].docx",
     ]),
    ("Santander_Complaint_Guide.docx",
     "Complaint Reference Guide",
     RGBColor(0x1D, 0x4E, 0xD8),
     [
         "7 complaint types × Stage 1 + escalation scenarios",
         "ISA Transfer Delay — 3 scenarios (upheld/declined/escalation)",
         "Incorrect Charge / Fee",
         "Payment Error",
         "APP Fraud — PSR PS23/3 reimbursement",
         "Mandate Change Error",
         "Poor Service / Communication",
         "Account Closure Delay",
         "Full letter templates, interest formula, DISP table",
     ]),
    ("Evelyn_Complaint_Agent_Angus.md",
     "Microsoft Copilot Agent",
     EMERALD,
     [
         "Complete system prompt for Copilot / Copilot Studio",
         "All 9 trigger phrases with output definitions",
         "Verbatim Stage 1 (5Q) and escalation (4Q) headings",
         "Upheld and declined letter templates with exact language",
         "FOS referral block — copy-exact text",
         "Standard sign-off block",
         "System behaviour rules (vault-ready, consistency)",
         "Paste into Copilot custom instructions or Copilot Studio",
     ]),
    ("In-app: renderComplaint()",
     "App.jsx Workflow",
     STONE7,
     [
         "4-step workflow wizard built into the prototype",
         "Step 1: Intake — category, date, eligible complainant check",
         "Step 2: Escalation triage — flag selection",
         "Step 3: Denial reason (if not upheld)",
         "Step 4: Case outcome — goodwill gesture option",
         "Triggered from Home screen 'Log complaint' action tile",
         "FCA DISP-compliant flow throughout",
         "State reset in closeWorkflow() — no ghost state",
     ]),
]

pw5 = 3.08
ph5 = 5.38
for ci6, (filename, subtitle, acol6, pts6) in enumerate(tools):
    x6 = 0.3 + ci6 * (pw5 + 0.19)
    y6 = 1.12
    R(s, x6, y6, pw5, ph5, fill=WHITE, line=STONE2)
    R(s, x6, y6, pw5, 0.07, fill=acol6)
    T(s, subtitle, x6+0.12, y6+0.10, pw5-0.2, 0.30, size=10, bold=True, color=acol6)
    T(s, filename, x6+0.12, y6+0.42, pw5-0.2, 0.22, size=7.5, color=STONE5, italic=True)
    R(s, x6+0.12, y6+0.66, pw5-0.22, 0.016, fill=STONE3)
    ML(s, ["· " + p for p in pts6], x6+0.12, y6+0.72, pw5-0.22, ph5-0.86, size=8, color=STONE7, leading=4)

# Bottom note
T(s, "Trigger phrases: 'evidence to support complaint stage 1' · 'Uphold - stage 1 letter' · 'Declined - stage 1 letter' · 'Escalation - resolution changed' and 5 more",
  0.3, 6.68, 12.7, 0.25, size=7.5, color=STONE5, italic=True)

advisor_stamp(s)
slide_num(s, 13)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — BUSINESS CASE SUMMARY
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
R(s, 0, 0, 13.33, 7.5, fill=DARK)
R(s, 0, 0, 13.33, 0.07, fill=RED)

T(s, "Business Case Summary", 0.6, 0.22, 10, 0.45, size=22, bold=True, color=WHITE)
T(s, "Phase 1 complete. Phase 2 awaiting stakeholder approval. Zero days to readiness.",
   0.6, 0.7, 12, 0.3, size=12, color=STONE5, italic=True)
R(s, 0.6, 1.05, 12.1, 0.02, fill=RED)

# Big metrics
metrics_top = [
    ("£137M", "Annualised opportunity\n(280,000 customers)"),
    ("£4.9M", "Cost saving per 10k\ncustomers · Year 1"),
    ("21",    "Features built\nacross 4 areas"),
    ("12",    "Paper workflows\ndigitised"),
    ("5",     "Paper forms\nretired"),
    ("8",     "Compliance\nframeworks"),
]
for i3, (num, lbl) in enumerate(metrics_top):
    bx3 = 0.5 + i3 * 2.05
    R(s, bx3, 1.22, 1.9, 1.45, fill=RGBColor(0x2C, 0x27, 0x24))
    R(s, bx3, 1.22, 1.9, 0.05, fill=RED)
    T(s, num, bx3, 1.3, 1.9, 0.7, size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    T(s, lbl, bx3, 1.98, 1.9, 0.65, size=8.5, color=STONE5, align=PP_ALIGN.CENTER)

# Forms retired
R(s, 0.5, 2.88, 5.8, 0.7, fill=RGBColor(0x2C, 0x27, 0x24))
T(s, "Paper forms retired:", 0.7, 2.96, 1.8, 0.28, size=9, bold=True, color=STONE3)
for fi, form in enumerate(["CA04 — Mandate change", "CA07 — Account closure",
                            "CA11 — Dormancy", "P17 — Bulk payment auth", "D18 — KYC update"]):
    R(s, 0.7 + fi*1.1, 3.06, 1.05, 0.3, fill=RGBColor(0x44, 0x40, 0x3C))
    T(s, form, 0.78 + fi*1.1, 3.1, 1.0, 0.22, size=7, color=LIGHTRED)

# Differentiators
R(s, 6.55, 2.88, 6.35, 2.3, fill=RGBColor(0x2C, 0x27, 0x24))
T(s, "Competitive Differentiators", 6.75, 2.96, 5.9, 0.28, size=10, bold=True, color=WHITE)
diffs = [
    "Only UK business bank with Voice ID biometric across app + phone + video",
    "Only prototype with live Companies House supplier risk scoring on dashboard",
    "Full personal/business data separation — app, call centre, post, credit — simultaneously",
    "MTD submission from receipt scan to HMRC in under 60 seconds",
    "Business health score — 5-factor live gauge no UK business bank currently offers",
]
ML(s, ["· " + d for d in diffs], 6.75, 3.3, 5.9, 1.8, size=9, color=STONE3, leading=6)

# Status
R(s, 0.5, 3.72, 5.8, 0.46, fill=EMERALD)
T(s, "Phase 1 — COMPLETE", 0.7, 3.8, 3.0, 0.3, size=11, bold=True, color=WHITE)
T(s, "Prototype live · All features functional", 3.6, 3.82, 2.6, 0.28, size=8.5, color=RGBColor(0xA7, 0xF3, 0xD0))

R(s, 0.5, 4.26, 5.8, 0.46, fill=RGBColor(0x1D, 0x4E, 0xD8))
T(s, "Phase 2 — 0 DAYS TO READY", 0.7, 4.34, 3.8, 0.3, size=11, bold=True, color=WHITE)
T(s, "Awaiting stakeholder approval", 4.4, 4.36, 1.8, 0.28, size=8.5, color=RGBColor(0xBF, 0xDB, 0xFE))

# Closing statement
R(s, 0, 5.35, 13.33, 1.78, fill=RGBColor(0x12, 0x0F, 0x0E))
T(s, "All work completed self-initiated, out of hours, in own time.", 0.6, 5.55, 9, 0.35, size=13, bold=True, color=WHITE)
T(s, "Alan Davidson · Business Banking Advisor · Santander · June 2026",
   0.6, 5.92, 9, 0.3, size=11, color=STONE5)
T(s, "£137M", 10.1, 5.48, 2.8, 0.75, size=48, bold=True, color=RED, align=PP_ALIGN.RIGHT)
T(s, "annualised opportunity", 10.1, 6.18, 2.8, 0.28, size=9, color=STONE5, align=PP_ALIGN.RIGHT)

slide_num(s, 14)

prs.save("Santander_Project_Record.pptx")
print("Saved → /home/user/santander/Santander_Project_Record.pptx")
