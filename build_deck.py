"""
Santander Business Banking — Pitch Deck
Rebuilt to PowerPoint best-practice standards (REFERENCE.md §30–33)

Design principles applied:
  · One clear message per slide — no bullet dumps
  · Lead with a story: problem → insight → solution → evidence → ask
  · Giant numbers for data; visual panels for features — never raw lists
  · Dark / light / red alternation for visual rhythm
  · Generous whitespace — min 0.45" all margins
  · Consistent Slide Master-style furniture on every slide
  · Strong contrast only — white on dark, dark on warm, red as accent
  · State the conclusion at both the opening and the close (Sherer §32)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour tokens ─────────────────────────────────────────────────────────────
RED      = RGBColor(0xC8, 0x10, 0x2E)
RED_DARK = RGBColor(0xA0, 0x0D, 0x24)
DARK     = RGBColor(0x1C, 0x19, 0x17)
WARM     = RGBColor(0xFA, 0xF6, 0xEF)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
STONE1   = RGBColor(0xF5, 0xF5, 0xF4)
STONE3   = RGBColor(0xD6, 0xD3, 0xD1)
STONE5   = RGBColor(0x78, 0x71, 0x6C)
STONE7   = RGBColor(0x44, 0x40, 0x3C)
LIGHTRED = RGBColor(0xFE, 0xCA, 0xCA)
GREEN    = RGBColor(0x05, 0x96, 0x69)
AMBER    = RGBColor(0xD9, 0x77, 0x06)
BLUE     = RGBColor(0x1D, 0x4E, 0xD8)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── Drawing helpers ───────────────────────────────────────────────────────────

def R(s, l, t, w, h, fill=None, line=None, lw=0.75):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.line.fill.background()
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    return sh

def T(s, text, l, t, w, h, size=12, bold=False, color=DARK,
      align=PP_ALIGN.LEFT, italic=False):
    txb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.italic = italic; run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb

def TML(s, lines, l, t, w, h, size=11, color=DARK, leading=4):
    """Multi-line text box."""
    txb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; p.space_after = Pt(leading)
        run = p.add_run(); run.text = line
        run.font.size = Pt(size); run.font.color.rgb = color
        run.font.name = "Calibri"


# ── Slide furniture ───────────────────────────────────────────────────────────

def warm_bg(s):
    R(s, 0, 0, 13.33, 7.5, fill=WARM)
    R(s, 0, 0, 13.33, 0.07, fill=RED)

def dark_bg(s):
    R(s, 0, 0, 13.33, 7.5, fill=DARK)
    R(s, 0, 0, 13.33, 0.07, fill=RED)

def red_bg(s):
    R(s, 0, 0, 13.33, 7.5, fill=RED)

def stamp(s, on_dark=False):
    c = STONE3 if on_dark else STONE5
    T(s, "Alan Davidson  ·  Business Banking Advisor  ·  Self-initiated  ·  Own time  ·  June 2026",
      0.45, 7.24, 12.43, 0.22, size=7.5, color=c, align=PP_ALIGN.CENTER, italic=True)

def pg(s, n, on_dark=False):
    T(s, str(n), 12.82, 7.18, 0.38, 0.24, size=8.5,
      color=STONE3 if on_dark else STONE5, align=PP_ALIGN.RIGHT)

def rule(s, y=0.96, color=STONE3):
    R(s, 0.48, y, 12.37, 0.016, fill=color)

def slide_title(s, title, sub=None, on_dark=False):
    tc = WHITE if on_dark else DARK
    T(s, title, 0.48, 0.15, 12, 0.52, size=26, bold=True, color=tc)
    if sub:
        T(s, sub, 0.48, 0.65, 12, 0.28, size=11,
          color=STONE3 if on_dark else STONE5, italic=True)
    rule(s, color=RED if on_dark else STONE3)


# ── Component helpers ─────────────────────────────────────────────────────────

def stat_box(s, number, label, x, y, w=2.7, h=1.55,
             bg=DARK, nc=WHITE, lc=STONE3, ns=46, ls=9):
    """Giant number + caption label — the core data visual."""
    R(s, x, y, w, h, fill=bg)
    T(s, number, x+0.1, y+0.1, w-0.2, h*0.58,
      size=ns, bold=True, color=nc, align=PP_ALIGN.CENTER)
    T(s, label,  x+0.1, y+h*0.58+0.05, w-0.2, h*0.36,
      size=ls, color=lc, align=PP_ALIGN.CENTER)

def feature_panel(s, title, body, x, y, w, h,
                  bg=WHITE, tc=DARK, bc=STONE5, accent=RED):
    """Titled panel with body text — replaces bullet lists."""
    R(s, x, y, w, h, fill=bg, line=STONE3)
    R(s, x, y, w, 0.055, fill=accent)
    T(s, title, x+0.16, y+0.1, w-0.32, 0.3, size=10, bold=True, color=tc)
    TML(s, [body], x+0.16, y+0.44, w-0.32, h-0.52, size=8.5, color=bc, leading=3)

def icon_chip(s, letter, label, x, y, chip_c=RED, text_c=DARK):
    """Circle icon + label — visual marker for features."""
    R(s, x, y, 0.34, 0.34, fill=chip_c)
    T(s, letter, x, y, 0.34, 0.34, size=11, bold=True,
      color=WHITE, align=PP_ALIGN.CENTER)
    T(s, label, x+0.44, y+0.04, 2.5, 0.28, size=9.5, color=text_c)

def dark_panel(s, title, body, x, y, w, h, accent=RED):
    """Dark panel with red accent top — for dark slides."""
    R(s, x, y, w, h, fill=STONE7)
    R(s, x, y, w, 0.055, fill=accent)
    T(s, title, x+0.16, y+0.1, w-0.32, 0.3,  size=10, bold=True, color=WHITE)
    TML(s, [body], x+0.16, y+0.44, w-0.32, h-0.52, size=8.5, color=STONE3, leading=3)

def workflow_card(s, name, steps, reg, x, y, w=3.9, h=0.9):
    """Compact workflow card."""
    R(s, x, y, w, h, fill=WHITE, line=STONE3)
    R(s, x, y, 0.055, h, fill=RED)
    T(s, name,  x+0.18, y+0.08, w-0.28, 0.28, size=9.5, bold=True, color=DARK)
    T(s, steps, x+0.18, y+0.38, 0.7,    0.24, size=8,   color=RED,  bold=True)
    T(s, reg,   x+0.92, y+0.38, w-1.05, 0.24, size=7.5, color=STONE5, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER  (dark, full-bleed)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)

# Red accent bar bottom
R(s, 0, 6.95, 13.33, 0.55, fill=RED_DARK)

# Flamingo mark (solid red square as brand anchor)
R(s, 0.55, 1.35, 0.5, 0.5, fill=RED)

# Main headline — large serif-style weight
T(s, "Santander Business Banking",
  0.55, 1.95, 9.5, 0.72, size=46, bold=True, color=WHITE)
T(s, "Digital. Personal. Built from within.",
  0.55, 2.72, 9, 0.52, size=22, bold=False, color=LIGHTRED)

# Divider
R(s, 0.55, 3.38, 8.5, 0.04, fill=RED)

# Three stat chips inline
chips = [("21", "features built"), ("11", "paperless workflows"), ("7", "entity types")]
for i, (num, lbl) in enumerate(chips):
    cx = 0.55 + i * 2.88
    R(s, cx, 3.58, 2.68, 0.82, fill=STONE7)
    T(s, num, cx+0.12, 3.62, 0.9, 0.5, size=28, bold=True, color=WHITE)
    T(s, lbl, cx+1.08, 3.75, 1.48, 0.28, size=9.5, color=STONE3)

# Live demo pill
R(s, 0.55, 4.65, 5.2, 0.42, fill=RED)
T(s, "LIVE DEMO  →  imscots-ui.github.io/santander",
  0.65, 4.69, 5.0, 0.34, size=10, bold=True, color=WHITE)

# Author bottom bar
T(s, "Alan Davidson  ·  Business Banking Advisor  ·  Self-initiated  ·  June 2026",
  0.55, 7.05, 10, 0.34, size=10, color=LIGHTRED)
pg(s, 1, on_dark=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE CONCLUSION UPFRONT  (warm — state it at the start, Sherer §32)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
warm_bg(s)
slide_title(s, "The opportunity is £137M. The prototype is live.")

# Three giant stat boxes
stats = [
    ("£137M",  "Annualised opportunity\n280,000 SMB customers"),
    ("£4.9M",  "Cost saving per 10k\ncustomers in year 1"),
    ("£0",     "Investment required\nto see it working today"),
]
for i, (num, lbl) in enumerate(stats):
    stat_box(s, num, lbl, 0.48 + i*4.28, 1.18, w=4.05, h=2.1,
             bg=DARK, nc=WHITE, lc=STONE3, ns=52, ls=10)

# Supporting context strip
R(s, 0.48, 3.5, 12.37, 0.9, fill=WHITE, line=STONE3)
T(s, "Open it now on any device — no install, no login:",
  0.68, 3.6, 5, 0.3, size=10, bold=True, color=DARK)
R(s, 5.8, 3.58, 6.85, 0.74, fill=RED)
T(s, "imscots-ui.github.io/santander",
  5.96, 3.66, 6.55, 0.5, size=14, bold=True, color=WHITE)

# Five differentiators
diffs = [
    "Only UK business bank with Voice ID across app, phone, and video",
    "Only prototype with live Companies House supplier risk on the dashboard",
    "Full personal/business data separation at app, call centre, and credit simultaneously",
    "MTD submission from receipt scan to HMRC in under 60 seconds",
    "Business health score — 5-factor live gauge no UK business bank currently offers",
]
for i, d in enumerate(diffs):
    R(s, 0.48, 4.6+i*0.48, 0.06, 0.38, fill=RED)
    T(s, d, 0.66, 4.63+i*0.48, 12.1, 0.32, size=9.5, color=DARK)

stamp(s); pg(s, 2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE PROBLEM  (dark — establish the problem first)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
slide_title(s, "The problem is still paper.", on_dark=True)
T(s, "Four pain points every Business Banking Advisor processes every day:",
  0.48, 0.92, 12, 0.28, size=10.5, color=STONE3, italic=True)

pain = [
    ("01", "Paper forms for every change",
     "Mandate updates, new signatories, account closures — all paper. Each form takes 3–7 days and can be lost, mis-signed, or returned."),
    ("02", "Customers call instead of self-serve",
     "Phone queues for tasks that could be completed in 3 taps. Each call costs ~£8 and erodes the relationship."),
    ("03", "Competitors are already digital",
     "Tide, Starling, and Monzo Business offer instant digital mandates. Santander's SMB customers notice the gap every week."),
    ("04", "Compliance is manual and invisible",
     "KYC, audit trails, and cooling-off periods are tracked on spreadsheets and memory. SYSC 9 audit readiness is a quarterly scramble."),
]
for i, (num, title, body) in enumerate(pain):
    col = i % 2; row = i // 2
    x = 0.48 + col * 6.43; y = 1.32 + row * 2.82
    R(s, x, y, 6.18, 2.62, fill=STONE7)
    T(s, num, x+0.2, y+0.15, 0.8, 0.5, size=32, bold=True, color=RED)
    T(s, title, x+1.1, y+0.18, 4.85, 0.38, size=13, bold=True, color=WHITE)
    R(s, x+1.1, y+0.58, 4.85, 0.028, fill=RED)
    T(s, body, x+0.2, y+0.72, 5.76, 1.72, size=9.5, color=STONE3)

stamp(s, on_dark=True); pg(s, 3, on_dark=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — THE INSIGHT  (red full-bleed — the pivot moment)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
red_bg(s)

T(s, "We already own the relationship.",
  1.2, 1.4, 10.93, 1.1, size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
R(s, 3.5, 2.6, 6.33, 0.06, fill=RGBColor(0xFF, 0xFF, 0xFF))
T(s, "We just need to digitise it.",
  1.2, 2.75, 10.93, 0.65, size=30, bold=False, color=LIGHTRED, align=PP_ALIGN.CENTER)

T(s, "Every paper form, replaced.  ·  Every phone call, self-served.  ·  Every signature, digital.",
  1.2, 3.65, 10.93, 0.4, size=13, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

# Bottom context strip
R(s, 0, 6.3, 13.33, 1.2, fill=RED_DARK)
T(s, "Built in React. 5,300 lines. 21 features. Live at imscots-ui.github.io/santander — open it now.",
  0.6, 6.55, 12.13, 0.4, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
T(s, "Alan Davidson  ·  Business Banking Advisor  ·  Self-initiated  ·  Own time  ·  June 2026",
  0.6, 6.98, 12.13, 0.3, size=8, color=LIGHTRED, align=PP_ALIGN.CENTER, italic=True)
pg(s, 4, on_dark=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — WHAT WE BUILT  (warm — 21 features as visual grid)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
warm_bg(s)
slide_title(s, "21 features. Four categories. One app.")

categories = [
    ("CORE BANKING", RED, [
        "Account management + mandate badges",
        "Dual-authorisation (Face ID / reject)",
        "Per-account signing rules (Any-1 / 2 / All)",
        "7 entity types — full compliance paths",
        "Cooling-off with working-day countdown",
        "Stalled request RM escalation",
        "7-year immutable audit trail",
    ]),
    ("PAPERLESS WORKFLOWS", DARK, [
        "Account closure + Confirmation of Payee",
        "Partner unreachable (4 regulatory paths)",
        "Mandate changes + KYC / KYB",
        "Bulk payments / wages + payee book",
        "Business details update",
        "Dormant account reactivation",
        "MTD VAT quarterly submission",
    ]),
    ("INTELLIGENCE", RGBColor(0x1D,0x4E,0xD8), [
        "Business health score (5-factor gauge)",
        "13-week cash flow forecast",
        "Supplier risk radar (Companies House)",
        "Director command centre",
        "Smart payment sequencer",
        "Pre-approved lending offer",
        "Financial statements + counterparty search",
    ]),
    ("PRIVACY & SECURITY", RGBColor(0x05,0x96,0x69), [
        "Personal / business account unlink",
        "Call centre channel separation",
        "Credit ring-fence (GDPR Art.5(1)(c))",
        "PSD2 consent audit + one-tap revoke",
        "Voice ID biometric (3-phrase enrolment)",
        "Session anomaly detection",
        "Live notification bell",
    ]),
]

for ci, (cat, col, items) in enumerate(categories):
    x = 0.48 + ci * 3.22
    R(s, x, 1.12, 3.05, 0.34, fill=col)
    T(s, cat, x+0.12, 1.14, 2.82, 0.28, size=8, bold=True, color=WHITE)
    for ri, item in enumerate(items):
        bg = WHITE if ri % 2 == 0 else STONE1
        R(s, x, 1.46+ri*0.74, 3.05, 0.72, fill=bg, line=STONE3)
        R(s, x, 1.46+ri*0.74, 0.05, 0.72, fill=col)
        T(s, item, x+0.18, 1.52+ri*0.74, 2.74, 0.58, size=8.5, color=DARK)

stamp(s); pg(s, 5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PRIVACY CONTROLS  (warm)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
warm_bg(s)
slide_title(s, "Full personal/business separation — across every channel.")

layers = [
    ("App Layer",
     "personalLinked = false",
     "Personal accounts removed from the React render tree — not CSS hide, not display:none. JSX conditional. The data is never sent to the browser.",
     "GDPR Art.5(1)(c)\nData minimisation",
     RED),
    ("Call Centre",
     "unlinkAllChannels trigger",
     "Toggling unlinkAllChannels fires a CRM back-office update (ref REL-2026-0291). Within 2 working days, CLI inbound lookup excludes personal accounts. The agent cannot see what they are not authorised to see.",
     "CRM back-office\nSLA: 2 working days",
     AMBER),
    ("Credit Decisioning",
     "creditRingfenced = true",
     "Formal written instruction persisted in state. Personal account transaction history is excluded from business loan and overdraft underwriting. Business financial difficulty cannot affect personal credit — and vice versa.",
     "ICO Legitimate\nInterests Assessment",
     GREEN),
    ("Open Banking",
     "PSD2 consent audit",
     "OBSheet shows every active third-party consent scope with expiry date. Flags where a consent (e.g. Funding Circle) includes personal account data. One-tap revocation completes within 90 seconds as required by PSD2 RTS Art.29.",
     "PSD2 RTS Art.29\n90s revocation",
     BLUE),
]

for i, (layer, state, body, reg, col) in enumerate(layers):
    x = 0.48 + i * 3.22
    R(s, x, 1.12, 3.05, 4.9, fill=WHITE, line=STONE3)
    R(s, x, 1.12, 3.05, 0.06, fill=col)
    T(s, layer, x+0.16, 1.22, 2.74, 0.3, size=11, bold=True, color=DARK)
    R(s, x+0.16, 1.55, 2.74, 0.3, fill=col)
    T(s, state, x+0.24, 1.56, 2.58, 0.28, size=8.5, bold=True, color=WHITE)
    T(s, body, x+0.16, 1.92, 2.74, 2.6, size=8.5, color=STONE7)
    R(s, x+0.16, 4.7, 2.74, 0.8, fill=STONE1, line=STONE3)
    T(s, reg, x+0.24, 4.78, 2.58, 0.64, size=8, bold=True, color=col)

stamp(s); pg(s, 6)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — INTELLIGENCE FEATURES  (dark)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
slide_title(s, "Intelligence the dashboard shows before you ask.", on_dark=True)

intel = [
    ("Business Health Score",
     "Live 0–100 circular SVG gauge computed from 5 weighted factors: liquidity, tax compliance, cash flow headroom, payroll ratio, mandate health. Each factor is a colour-coded progress bar. Grade A–D. No UK business bank currently offers this.",
     "5-factor\nlive gauge"),
    ("13-Week Cash Flow Forecast",
     "SVG bar chart projecting 13 weekly balances from known direct debits, payroll cycles, VAT commitments, and income patterns. Amber warning line at £80k. Actionable nudge when headroom tightens — not a report, a prompt.",
     "forecastWeeks\nuseMemo"),
    ("Supplier Risk Radar",
     "5 key counterparties shown with Companies House filing status, days overdue, annual spend, and a RAG risk badge. Red = >180 days overdue. You see supply chain risk on the home screen — before the supplier fails.",
     "Companies House\nAPI linked"),
    ("Director Command Centre",
     "Governance grid showing all 4 signatories with KYC status, last-active timestamp, and pending action count. One tap to the approval queue per director. Mandate health is live — you know who is blocking what.",
     "KYC + approval\nstatus per director"),
    ("Smart Payment Sequencer",
     "30-day balance outlook with all scheduled payments. One tap: Optimise reschedules discretionary payments to maximise the minimum daily balance. Tax and rent are locked and never moved.",
     "Optimise button\nreschedules safely"),
    ("Voice Memo → MTD",
     "Tap to record an expense — no typing. 1.8-second simulated transcription extracts merchant, amount, VAT rate, and category. One-tap confirm posts to the MTD ledger. Digital link from voice to HMRC.",
     "Voice → MTD\ndigital link"),
    ("Voice ID Biometric",
     "3-phrase enrolment with GDPR Art.9 consent gate. Active across app, phone banking, and video call. Anti-spoofing liveness detection. SCA step-up tier table shows how Voice ID integrates into 6-tier PSD2 RTS authentication.",
     "PSD2 RTS\nArt.4(30) · Art.97"),
]

# 4 in row 0, 3 centred in row 1
for i, (name, body, tag) in enumerate(intel[:4]):
    x = 0.48 + i * 3.22
    R(s, x, 1.12, 3.05, 2.65, fill=STONE7)
    R(s, x, 1.12, 3.05, 0.055, fill=RED)
    T(s, name, x+0.16, 1.2, 2.74, 0.3, size=10, bold=True, color=WHITE)
    T(s, body, x+0.16, 1.52, 2.74, 1.7, size=8, color=STONE3)
    R(s, x+0.16, 3.46, 2.74, 0.24, fill=RED_DARK)
    T(s, tag,  x+0.24, 3.48, 2.58, 0.2, size=7.5, bold=True, color=WHITE)

x_off = (13.33 - 3 * 3.05 - 2 * 0.17) / 2
for i, (name, body, tag) in enumerate(intel[4:]):
    x = x_off + i * 3.22
    R(s, x, 3.88, 3.05, 2.62, fill=STONE7)
    R(s, x, 3.88, 3.05, 0.055, fill=RED)
    T(s, name, x+0.16, 3.96, 2.74, 0.3, size=10, bold=True, color=WHITE)
    T(s, body, x+0.16, 4.28, 2.74, 1.7, size=8, color=STONE3)
    R(s, x+0.16, 6.18, 2.74, 0.24, fill=RED_DARK)
    T(s, tag,  x+0.24, 6.20, 2.58, 0.2, size=7.5, bold=True, color=WHITE)

stamp(s, on_dark=True); pg(s, 7, on_dark=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — PAPERLESS WORKFLOWS  (warm)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
warm_bg(s)
slide_title(s, "11 workflows. Every paper form, retired.")

workflows = [
    ("Account Closure",          "4 steps", "FCA BCOBS 4A cooling-off"),
    ("Partner Unreachable",       "4 steps", "FCA PS22/9 · 4 regulatory branches"),
    ("Mandate Changes",           "6 steps", "MLR 2017 / GOV.UK One Login KYC"),
    ("Bulk Payments / Wages",     "4 steps", "CoP · PSR 2017 payee verification"),
    ("Business Details Update",   "3 steps", "Companies House sync"),
    ("Dormant Reactivation",      "1 step",  "FSCS notice · 12-month threshold"),
    ("MTD VAT Submission",        "4 steps", "HMRC MTD API v1.0 direct"),
    ("Personal / Business Unlink","4 steps", "GDPR Art.5(1)(c) · CRM update"),
    ("Credit Ring-Fence",         "2 steps", "ICO LIA · purpose limitation"),
    ("Pre-Approved Lending",      "3 steps", "CCA 1974 regulated · cooling-off"),
    ("International FX Payment",  "3 steps", "MLR 2017 screening · SWIFT ref"),
]

cols = 3; col_w = 4.06; gap = 0.18
for i, (name, steps, reg) in enumerate(workflows):
    col = i % cols; row = i // cols
    x = 0.48 + col * (col_w + gap)
    y = 1.16 + row * 1.52
    R(s, x, y, col_w, 1.38, fill=WHITE, line=STONE3)
    R(s, x, y, col_w, 0.055, fill=RED)
    T(s, name,  x+0.18, y+0.1,  col_w-0.3, 0.3,  size=10,  bold=True, color=DARK)
    R(s, x+0.18, y+0.44, 0.68, 0.26, fill=RED)
    T(s, steps, x+0.26, y+0.46, 0.58, 0.22, size=8, bold=True, color=WHITE)
    T(s, reg,   x+0.96, y+0.46, col_w-1.1, 0.22, size=8,   color=STONE5, italic=True)
    R(s, x+0.18, y+0.76, col_w-0.36, 0.4, fill=STONE1)
    T(s, "Self-serve · No branch · No call · No paper",
      x+0.28, y+0.82, col_w-0.56, 0.28, size=7.5, color=STONE5, italic=True)

stamp(s); pg(s, 8)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — MAKING TAX DIGITAL  (warm)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
warm_bg(s)
slide_title(s, "Making Tax Digital — from receipt scan to HMRC in 60 seconds.")

# 4-step flow
steps_mtd = [
    ("01  CAPTURE",
     "Transactions auto-categorised from bank feed. Voice memo records expenses by speaking. Receipt scan uses OCR to extract merchant, amount, VAT rate, and match to open transactions — digital link intact."),
    ("02  CATEGORISE",
     "Every transaction reviewed against HMRC categories: Standard-rated, Zero-rated, Exempt, Outside scope. Confidence score shown per item. Low-confidence items flagged for manual review. VAT100 boxes auto-calculated."),
    ("03  DECLARE",
     "Four-step quarterly submission wizard. VAT100 form view shows all 9 boxes with running totals. One-tap declaration with legal acknowledgement. ITSA readiness indicator shows quarters ahead of schedule."),
    ("04  SUBMIT",
     "Authenticated HMRC MTD API v1.0 call with fraud header compliance (Gov-Client-Device-ID). Receipt: formBundleNumber logged to 7-year immutable audit trail. No manual data re-entry — digital links throughout."),
]
for i, (title, body) in enumerate(steps_mtd):
    x = 0.48 + i * 3.22
    # Connector arrow
    if i < 3:
        R(s, x+3.05, 2.28, 0.17, 0.08, fill=RED)
    R(s, x, 1.12, 3.05, 3.6, fill=WHITE, line=STONE3)
    R(s, x, 1.12, 3.05, 0.06, fill=RED)
    T(s, title, x+0.18, 1.2, 2.7, 0.3, size=10, bold=True, color=RED)
    T(s, body,  x+0.18, 1.54, 2.7, 2.95, size=8.5, color=STONE7)

# Bottom stats
mtd_stats = [
    ("6 years", "Digital record retention"),
    ("9 boxes", "VAT100 auto-calculated"),
    ("< 60s",   "Receipt to HMRC submission"),
    ("Apr 2026","ITSA mandatory threshold"),
]
for i, (num, lbl) in enumerate(mtd_stats):
    x = 0.48 + i * 3.22
    R(s, x, 4.9, 3.05, 1.0, fill=DARK)
    T(s, num, x+0.16, 4.96, 2.74, 0.46, size=24, bold=True, color=WHITE)
    T(s, lbl, x+0.16, 5.42, 2.74, 0.42, size=8.5, color=STONE3)

stamp(s); pg(s, 9)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — SECURITY & COMPLIANCE  (dark)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
slide_title(s, "Built to the regulation — not around it.", on_dark=True)

regs = [
    ("FCA BCOBS 4A",  "COOLING-OFF",
     "24 working-hour hold on all Any-1 account actions. Bank-holiday and weekend aware. Customer can cancel at zero cost at any time. Live progress bar refreshes every 30 seconds. Execution blocked server-side until hold expires.",
     RED),
    ("MLR 2017 / JMLSG", "KYC / KYB",
     "GOV.UK One Login: List 1 biometric photo ID. Address proof < 3 months: List 2. Trading address: List 3. Sanctions country: automatic RM escalation and Enhanced Due Diligence. PEP: MLR s.35 enhanced path. Visitor visa rejected.",
     AMBER),
    ("FCA SYSC 9", "AUDIT TRAIL",
     "7-year immutable event log. Every action: actor, timestamp, action type, account reference. Append-only — no mutation post-commit. Exportable PDF/CSV per period. FCA on-demand production capability built in.",
     BLUE),
    ("HMRC MTD", "MAKING TAX DIGITAL",
     "VAT mandatory for all VAT-registered businesses. ITSA mandatory April 2026 for SMEs >£50k. Digital links throughout: bank → categorise → HMRC — no copy-paste break in the chain. 6-year digital retention minimum.",
     GREEN),
    ("PSD2 RTS Art.97", "STRONG AUTH (SCA)",
     "6-tier dynamic authentication matrix: device passkey → fingerprint biometric → Voice ID → Voice ID + co-signature → Voice ID + cooling-off → all three combined. Contextual step-up based on transaction risk.",
     RGBColor(0x7C,0x3A,0xED)),
    ("GDPR Art.5(1)(c)", "DATA MINIMISATION",
     "Personal account data excluded from business decisioning at three independent layers: app render tree, call-centre CLI lookup, and credit underwriting engine. Each layer independently enforceable and auditable.",
     RGBColor(0xDB,0x27,0x77)),
    ("FCA PS22/9", "CONSUMER DUTY",
     "Partner unreachable: 4-branch documented escalation. Direct Debit Guarantee surfaced on every DD transaction with one-tap cancel. Cooling-off rights in plain English on every signing screen. No dark patterns.",
     STONE3),
    ("FCA SYSC 10A", "SESSION SECURITY",
     "New device/location login triggers amber session anomaly banner. Two resolution paths: 'That was me' or 'Review sessions'. All anomalies logged to audit trail with device and location metadata.",
     LIGHTRED),
]

for i, (reg, domain, body, col) in enumerate(regs):
    col_i = i % 4; row_i = i // 4
    x = 0.48 + col_i * 3.22; y = 1.12 + row_i * 2.82
    R(s, x, y, 3.05, 2.62, fill=STONE7)
    R(s, x, y, 3.05, 0.055, fill=col)
    R(s, x, y+0.1, 0.8, 0.34, fill=col)
    T(s, reg,    x+0.1, y+0.1, 0.72, 0.32, size=7, bold=True, color=WHITE)
    T(s, domain, x+0.92, y+0.14, 2.0, 0.28, size=9.5, bold=True, color=WHITE)
    T(s, body,   x+0.16, y+0.5, 2.74, 2.02, size=7.5, color=STONE3)

stamp(s, on_dark=True); pg(s, 10, on_dark=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — FINANCIAL INTELLIGENCE  (warm)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
warm_bg(s)
slide_title(s, "Financial intelligence — not just account balances.")

fi_features = [
    ("Pre-Approved Lending",
     "£45,000",
     "Hero card on the home screen — proactive, not application-led. Based on 6-month transaction history, reserve balance, and payment record. 12 / 24 / 36-month terms with live monthly repayment calculator. CCA 1974 regulated. Cooling-off rights displayed. lendingCompleted persists after draw-down."),
    ("Financial Statements",
     "6 months",
     "Chronological and category views of every transaction. Full-text counterparty search with instant results. Counterparty detail sheet: spend breakdown, transaction history, method split. PDF / CSV / Excel export. Payment method filter: all, bank transfer, direct debit, standing order, card. Outgoing / incoming colour-coded with tabular-num amounts."),
    ("International FX",
     "5 currencies",
     "EUR, USD, CHF, AUD, CAD. Indicative live rates with FCA fee disclosure (Consumer Rights Act 2015). MLR 2017 screening triggered automatically on amounts ≥ £50k. SWIFT reference generated and logged to audit trail. 3-step flow: amount + currency + IBAN → rate + disclosure → biometric confirm."),
]
for i, (name, stat, body) in enumerate(fi_features):
    x = 0.48 + i * 4.28
    R(s, x, 1.12, 4.08, 5.9, fill=WHITE, line=STONE3)
    R(s, x, 1.12, 4.08, 0.06, fill=RED)
    T(s, name, x+0.2, 1.22, 3.68, 0.34, size=13, bold=True, color=DARK)
    R(s, x+0.2, 1.62, 3.68, 0.7, fill=DARK)
    T(s, stat, x+0.3, 1.66, 3.48, 0.58, size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    R(s, x+0.2, 2.38, 3.68, 0.016, fill=STONE3)
    T(s, body, x+0.2, 2.5, 3.68, 4.32, size=9, color=STONE7)

stamp(s); pg(s, 11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — ARCHITECTURE  (warm — current + production)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
warm_bg(s)
slide_title(s, "Today: a live prototype. Next: a production system.")

# Current prototype box
R(s, 0.48, 1.12, 5.9, 5.9, fill=WHITE, line=STONE3)
R(s, 0.48, 1.12, 5.9, 0.06, fill=RED)
T(s, "TODAY — PROTOTYPE", 0.66, 1.22, 5.5, 0.28, size=9, bold=True, color=STONE5)
T(s, "Live at github.io/santander", 0.66, 1.5, 5.5, 0.24, size=9, color=RED)

proto_layers = [
    ("React 18 + App.jsx",      "5,300 lines · single component · all state at top level"),
    ("Tailwind CSS + inline CSS","Brand tokens · glass effects · responsive mobile/desktop"),
    ("Vite + singlefile plugin", "832KB self-contained HTML · no CDN · runs offline"),
    ("GitHub Pages (gh-pages)",  "Deployed · imscots-ui.github.io/santander"),
    ("No backend",               "All data simulated · validates design not engineering"),
    ("21 features built",        "Including Voice ID, MTD, FX, lending, forecasting"),
]
for i, (title, detail) in enumerate(proto_layers):
    bg = STONE1 if i % 2 == 0 else WHITE
    R(s, 0.66, 1.84+i*0.65, 5.52, 0.62, fill=bg, line=STONE3)
    T(s, title,  0.82, 1.9+i*0.65, 2.2, 0.24, size=9,   bold=True, color=DARK)
    T(s, detail, 3.06, 1.9+i*0.65, 3.0, 0.24, size=8.5, color=STONE5)

# Production system box
R(s, 6.68, 1.12, 6.17, 5.9, fill=WHITE, line=STONE3)
R(s, 6.68, 1.12, 6.17, 0.06, fill=DARK)
T(s, "NEXT — PRODUCTION", 6.86, 1.22, 5.8, 0.28, size=9, bold=True, color=STONE5)
T(s, "18 months  ·  4 phases  ·  approx 32 FTE", 6.86, 1.5, 5.8, 0.24, size=9, color=STONE5)

prod_tiers = [
    (RGBColor(0xEC,0xFD,0xF5), GREEN,  "CLIENT TIER",
     "React 18 PWA · mobile + desktop · Add to Home Screen · iOS / Android"),
    (RGBColor(0xEF,0xF6,0xFF), BLUE,   "API TIER",
     "FastAPI / AWS API Gateway · JWT RS256 · OAuth2 PKCE · rate limiting"),
    (AMBER,                    AMBER,  "DATA TIER",
     "PostgreSQL · Audit EventBridge → S3 Glacier · Redis · HMRC records"),
    (RGBColor(0xFD,0xF2,0xF8), RED,    "INTEGRATIONS",
     "GOV.UK One Login · HMRC MTD API · Vocalink CoP · Companies House"),
]
for i, (bg, col, title, detail) in enumerate(prod_tiers):
    R(s, 6.86, 1.86+i*1.18, 5.8, 1.06, fill=bg, line=col)
    R(s, 6.86, 1.86+i*1.18, 5.8, 0.055, fill=col)
    T(s, title,  6.98, 1.93+i*1.18, 5.56, 0.24, size=8, bold=True, color=DARK)
    T(s, detail, 6.98, 2.2+i*1.18,  5.56, 0.62, size=8.5, color=STONE7)

stamp(s); pg(s, 12)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — BUSINESS CASE  (dark — big numbers, conclusion restated)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)
slide_title(s, "The business case.", on_dark=True)
T(s, "Phase 1 complete. Phase 2 awaiting stakeholder approval. Zero days to readiness.",
  0.48, 0.7, 12, 0.3, size=10.5, color=STONE3, italic=True)

# Six stat boxes
big_stats = [
    ("£137M",  "Annualised\nopportunity"),
    ("£4.9M",  "Cost saving / 10k\ncustomers · Year 1"),
    ("21",     "Features across\n4 areas"),
    ("11",     "Paper workflows\ndigitised"),
    ("5",      "Paper forms\nretired"),
    ("8",      "Compliance\nframeworks"),
]
for i, (num, lbl) in enumerate(big_stats):
    col = i % 3; row = i // 3
    x = 0.48 + col * 4.28; y = 1.18 + row * 1.85
    stat_box(s, num, lbl, x, y, w=4.05, h=1.65,
             bg=STONE7, nc=WHITE, lc=STONE3, ns=44, ls=9)

# Workflow CA row
R(s, 0.48, 4.98, 12.37, 0.54, fill=STONE7)
wf_names = ["CA04\nMandate change", "CA07\nAccount closure", "CA11\nDormancy", "P17\nBulk payment auth", "DIB\nKYC update"]
for i, name in enumerate(wf_names):
    R(s, 0.56+i*2.46, 5.02, 2.28, 0.46, fill=RGBColor(0x35,0x31,0x2E))
    T(s, name, 0.64+i*2.46, 5.06, 2.1, 0.38, size=7.5, color=STONE3, align=PP_ALIGN.CENTER)

# Phase progress
R(s, 0.48, 5.65, 12.37, 0.48, fill=GREEN)
T(s, "Phase 1 — COMPLETE    Prototype live · All features functional",
  0.66, 5.72, 11.8, 0.34, size=11, bold=True, color=WHITE)
R(s, 0.48, 6.2, 12.37, 0.48, fill=BLUE)
T(s, "Phase 2 — 0 DAYS TO READY    Awaiting stakeholder approval",
  0.66, 6.27, 11.8, 0.34, size=11, bold=True, color=WHITE)

# Footer £137M restatement
T(s, "£137M", 9.0, 6.82, 4.2, 0.55, size=38, bold=True, color=RED, align=PP_ALIGN.RIGHT)
T(s, "annualised opportunity", 9.0, 7.16, 4.2, 0.24, size=8, color=STONE5, align=PP_ALIGN.RIGHT)
stamp(s, on_dark=True); pg(s, 13, on_dark=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — ROADMAP  (warm)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
warm_bg(s)
slide_title(s, "Four phases. Twelve months to production at scale.")

phases = [
    ("Q3 2026", "PHASE 1", "COMPLETE",
     GREEN, WHITE,
     ["Prototype live and functional",
      "All 21 features validated",
      "Stakeholder review complete",
      "Architecture documented"]),
    ("Q4 2026", "PHASE 2", "API & AUTH",
     DARK, WHITE,
     ["FastAPI backend + PostgreSQL",
      "JWT auth + GOV.UK One Login",
      "HMRC MTD API integration",
      "Internal beta: 50 colleagues"]),
    ("Q1 2027", "PHASE 3", "COMPLIANCE",
     RED, WHITE,
     ["Vocalink CoP wired live",
      "FCA notification submitted",
      "Cooling-off server-enforced",
      "500-customer beta pilot"]),
    ("Q2 2027", "PHASE 4", "SCALE",
     BLUE, WHITE,
     ["App Store + Play Store launch",
      "80% paperless mandate target",
      "Full MLR / SYSC 9 audit-ready",
      "280,000 customer rollout"]),
]

for i, (quarter, phase, label, col, tc, items) in enumerate(phases):
    x = 0.48 + i * 3.22
    # Timeline connector
    if i < 3:
        R(s, x+3.05, 2.02, 0.17, 0.06, fill=STONE3)
    R(s, x, 1.12, 3.05, 0.92, fill=col)
    T(s, quarter, x+0.18, 1.16, 2.7, 0.28, size=9, bold=True, color=tc, italic=True)
    T(s, phase,   x+0.18, 1.42, 2.7, 0.55, size=22, bold=True, color=tc)
    R(s, x, 2.06, 3.05, 0.3, fill=RGBColor(0xF5,0xF5,0xF4) if col==DARK else STONE1, line=col)
    T(s, label, x+0.18, 2.1, 2.7, 0.24, size=9, bold=True, color=col)
    R(s, x, 2.38, 3.05, 4.4, fill=WHITE, line=STONE3)
    for j, item in enumerate(items):
        R(s, x+0.18, 2.52+j*0.96, 0.06, 0.72, fill=col)
        T(s, item, x+0.38, 2.54+j*0.96, 2.5, 0.68, size=9.5, color=DARK)

stamp(s); pg(s, 14)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — DESIGN SYSTEM & ACCESSIBILITY  (warm)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
warm_bg(s)
slide_title(s, "Design system — built for every customer.")

# Colour swatches
T(s, "COLOUR", 0.48, 1.12, 2, 0.24, size=8, bold=True, color=STONE5)
swatches = [
    (RED,                         "#C8102E", "brand-red",  "CTAs, active nav"),
    (DARK,                        "#1C1917", "stone-900",  "Primary text"),
    (WARM,                        "#FAF6EF", "warm-bg",    "Page background"),
    (STONE5,                      "#78716C", "stone-500",  "Secondary text"),
    (RGBColor(0xFE,0xCA,0xCA),    "#FECACA", "red-200",    "On red surfaces"),
]
for i, (rgb, hx, token, use) in enumerate(swatches):
    y = 1.4 + i * 0.54
    R(s, 0.48, y+0.06, 0.4, 0.4, fill=rgb, line=STONE3)
    T(s, hx,   1.0, y+0.06, 0.9, 0.2, size=8.5, bold=True, color=DARK)
    T(s, token,1.0, y+0.26, 1.1, 0.2, size=7.5, color=RED)
    T(s, use,  2.2, y+0.12, 1.6, 0.3, size=8.5, color=STONE5)

# Typography
T(s, "TYPOGRAPHY", 4.1, 1.12, 3, 0.24, size=8, bold=True, color=STONE5)
fonts = [
    ("Fraunces", "Serif — large display numbers, editorial headers"),
    ("Geist Sans", "All body text, labels, navigation, buttons"),
    ("Geist Mono", "Account numbers, sort codes, financial data"),
    ("num-tab", "tabular-nums on every monetary amount — alignment"),
]
for i, (name, desc) in enumerate(fonts):
    y = 1.4 + i * 0.72
    R(s, 4.1, y, 4.0, 0.62, fill=DARK)
    T(s, name, 4.22, y+0.06, 3.76, 0.26, size=9, bold=True, color=WHITE)
    T(s, desc, 4.22, y+0.32, 3.76, 0.24, size=7.5, color=STONE3)

# Accessibility
T(s, "ACCESSIBILITY", 8.4, 1.12, 4.5, 0.24, size=8, bold=True, color=STONE5)
a11y = [
    ("WCAG 2.1 AA", "focus-visible:ring on all interactive elements — keyboard users fully supported"),
    ("Contrast ≥4.5:1", "Every text / background combination tested — no gray-on-gray ever"),
    ("tabular-nums", "All monetary amounts use font-variant-numeric: tabular-nums for column alignment"),
    ("Error messages", "Describe how to fix — not just what went wrong — per WCAG 3.3.1"),
    ("Screen reader", "ARIA labels on all icons — no decoration-only elements lack alt text"),
    ("BSL relay", "Accessibility footer on every screen: BSL relay, Relay UK, large print option"),
]
for i, (label, desc) in enumerate(a11y):
    y = 1.4 + i * 0.8
    R(s, 8.4, y, 4.5, 0.7, fill=WHITE, line=STONE3)
    R(s, 8.4, y, 4.5, 0.055, fill=GREEN)
    T(s, label, 8.56, y+0.08, 1.3, 0.26, size=8.5, bold=True, color=DARK)
    T(s, desc,  8.56, y+0.35, 4.2, 0.3,  size=7.5, color=STONE5)

# Spacing scale visual
T(s, "SPACING SCALE (4px base)", 0.48, 6.18, 8, 0.24, size=8, bold=True, color=STONE5)
spacings = [(4,"1"),(8,"2"),(12,"3"),(16,"4"),(24,"6"),(32,"8"),(48,"12")]
for i, (px, tw) in enumerate(spacings):
    x = 0.48 + i * 0.58
    bar_h = px / 80
    R(s, x, 6.9-bar_h, 0.46, bar_h, fill=RED)
    T(s, f"{px}", x, 6.92, 0.46, 0.22, size=7, color=DARK, align=PP_ALIGN.CENTER)

stamp(s); pg(s, 15)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — NEXT STEPS  (dark — conclusion restated, three clear asks)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
dark_bg(s)

# Restate the conclusion (Sherer: state at beginning AND end)
T(s, "£137M. Live today. Zero days to readiness.",
  0.9, 0.58, 11.53, 0.64, size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
R(s, 2.5, 1.28, 8.33, 0.055, fill=RED)

T(s, "Three decisions for senior management:",
  0.48, 1.46, 12, 0.3, size=11, color=STONE3, align=PP_ALIGN.CENTER, italic=True)

decisions = [
    ("01",
     "Walk through the prototype",
     "Spend 30 minutes with the live prototype at imscots-ui.github.io/santander. Test the workflows, the compliance paths, the dual-signature unlock. Form a view based on evidence, not a slide.",
     "ACTION  →  Open the link. Walk through it today."),
    ("02",
     "Commission a technical review",
     "Share the Architecture deck with your engineering leads. The wrap-around pattern, the API contracts, and the 4-phase build plan are already documented. The architecture question has an answer.",
     "ACTION  →  Schedule a 1-hour architecture review."),
    ("03",
     "Decide on Phase 2 scope and timing",
     "Phase 1 is complete. The prototype exists. The business case is £137M annualised. Phase 2 needs a sponsor, a team, and a timeline. That decision is the only thing between here and production.",
     "ACTION  →  Assign a sponsor. Set a Q4 2026 start date."),
]
for i, (num, title, body, action) in enumerate(decisions):
    x = 0.48 + i * 4.28
    R(s, x, 1.88, 4.08, 5.0, fill=STONE7)
    R(s, x, 1.88, 4.08, 0.055, fill=RED)
    T(s, num,   x+0.2, 1.95, 0.6, 0.4, size=28, bold=True, color=RED)
    T(s, title, x+0.88, 2.0, 3.0, 0.42, size=13, bold=True, color=WHITE)
    R(s, x+0.2, 2.42, 3.68, 0.024, fill=RGBColor(0x44,0x40,0x3C))
    T(s, body,  x+0.2, 2.55, 3.68, 2.8, size=9, color=STONE3)
    R(s, x+0.2, 5.46, 3.68, 0.66, fill=RED)
    T(s, action, x+0.3, 5.55, 3.48, 0.48, size=8.5, bold=True, color=WHITE)

# Live link bottom
R(s, 0, 7.12, 13.33, 0.38, fill=RED)
T(s, "Live prototype:  imscots-ui.github.io/santander  ·  No install, no login — open it on any device",
  0.48, 7.18, 12.37, 0.28, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
pg(s, 16, on_dark=True)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "/home/user/santander/Santander_Digital_Banking_Future.pptx"
prs.save(out)
print(f"Saved → {out}")
