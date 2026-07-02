#!/usr/bin/env python3
"""
Flag text boxes that visually overlap each other across all decks.

Complements scripts/deck-qa.py (which checks shapes vs the slide *bounds*);
this checks shapes vs *each other* — the "heading sitting on top of the tile
row" class of bug. Heuristic, so it reports candidates to eyeball with
scripts/deck_to_html.py, not a hard gate.

Known false-positive shapes (opposite-aligned footers, eyebrows above a wide
title, page-number over a full-width-but-short title box) are tagged so real
issues stand out. Usage:  python3 scripts/deck_overlap.py
"""
from pptx import Presentation
from pptx.util import Emu
import glob

# lowercased fragments that are habitually flagged but don't visually collide
FP = ('internal', 'confidential', 'alan davidson', 'session 2', ' / 15', 'live prototype',
      'since the june 2026 review', ' / 14')

def boxes(slide):
    out = []
    for sh in slide.shapes:
        try:
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text.strip()
            if not t:
                continue
            l, tp = Emu(sh.left).inches, Emu(sh.top).inches
            w, h = Emu(sh.width).inches, Emu(sh.height).inches
            out.append((sh.name, l, tp, l + w, tp + h, t[:26]))
        except Exception:
            pass
    return out

real = 0
for deck in sorted(glob.glob('Santander_*.pptx')):
    prs = Presentation(deck)
    for i, s in enumerate(prs.slides):
        bs = boxes(s)
        for x in range(len(bs)):
            for y in range(x + 1, len(bs)):
                ox = max(0, min(bs[x][3], bs[y][3]) - max(bs[x][1], bs[y][1]))
                oy = max(0, min(bs[x][4], bs[y][4]) - max(bs[x][2], bs[y][2]))
                if ox > 0.12 and oy > 0.10:
                    a, b = bs[x][5].lower(), bs[y][5].lower()
                    fp = any(k in a or k in b for k in FP)
                    if not fp:
                        real += 1
                        print(f"REAL  {deck} s{i}: {ox:.1f}x{oy:.1f}in  [{bs[x][5]!r}] x [{bs[y][5]!r}]")
print(f"\nlikely-real overlaps: {real}  (footer/eyebrow/page-number false-positives suppressed)")
