#!/usr/bin/env python3
"""Flag text boxes whose text likely OVERFLOWS the box (text-too-big-for-tile).
Heuristic — confirm flags by rendering with scripts/deck_to_html.py (clipping off).
Complements deck-qa.py (bounds) and deck_overlap.py (box-vs-box)."""
import glob, math
from pptx import Presentation
from pptx.util import Emu
flags=[]
for deck in sorted(glob.glob('Santander_*.pptx')):
    prs=Presentation(deck)
    for i,s in enumerate(prs.slides):
        for sh in s.shapes:
            try:
                if not sh.has_text_frame: continue
                txt=sh.text_frame.text.strip()
                if len(txt)<4: continue
                w=Emu(sh.width).inches; h=Emu(sh.height).inches
                if w<=0 or h<=0: continue
                # smallest explicit font size in the box (else assume 11pt body)
                szs=[r.font.size.pt for p in sh.text_frame.paragraphs for r in p.runs if r.font.size]
                fpt=min(szs) if szs else 11
                nowrap = sh.text_frame.word_wrap is False
                inner_w_pt=max(6, w*72-8)
                cpl=max(1, inner_w_pt/(fpt*0.48))       # chars per line
                # longest single paragraph drives wrapping
                longest=max((len(p.text) for p in sh.text_frame.paragraphs), default=len(txt))
                lines = 1 if nowrap else math.ceil(longest/cpl)
                # account for explicit newlines
                lines=max(lines, len(sh.text_frame.paragraphs))
                need_h=lines*fpt*1.25/72
                ratio=need_h/h
                # nowrap: check horizontal fit instead
                if nowrap:
                    need_w=longest*fpt*0.48/72
                    ratio=need_w/w
                if ratio>1.25:
                    flags.append((round(ratio,2),deck.split('_',1)[1][:12],i,sh.name,txt[:26],round(w,2),round(h,2),fpt,'NOWRAP' if nowrap else ''))
            except Exception: pass
flags.sort(reverse=True)
for f in flags[:40]:
    print(f"ratio {f[0]}  {f[1]} s{f[2]}  {f[3]} {f[4]!r}  box {f[5]}x{f[6]} {f[7]}pt {f[8]}")
print(f"\ntotal likely-overflow flags: {len(flags)}")
