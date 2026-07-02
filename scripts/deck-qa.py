from pptx import Presentation
from pptx.util import Emu
import glob
decks = sorted(glob.glob('Santander_*.pptx'))
total_issues = 0
for path in decks:
    prs = Presentation(path)
    W, H = prs.slide_width, prs.slide_height
    issues = []
    conf_slides = 0
    for si, slide in enumerate(prs.slides, 1):
        has_conf = False
        for sh in slide.shapes:
            # overflow check (only shapes with a position)
            try:
                l, t, w, h = sh.left, sh.top, sh.width, sh.height
            except Exception:
                continue
            if None in (l, t, w, h):
                continue
            tol = 9525  # ~0.01in tolerance
            if l < -tol or t < -tol or (l + w) > W + tol or (t + h) > H + tol:
                over = []
                if l < -tol: over.append(f"left {Emu(l).inches:.2f}in")
                if t < -tol: over.append(f"top {Emu(t).inches:.2f}in")
                if (l+w) > W+tol: over.append(f"right +{Emu(l+w-W).inches:.2f}in")
                if (t+h) > H+tol: over.append(f"bottom +{Emu(t+h-H).inches:.2f}in")
                issues.append(f"  slide {si}: '{(sh.name or '')[:24]}' overflow ({', '.join(over)})")
            if sh.has_text_frame and 'CONFIDENTIAL' in (sh.text_frame.text or '').upper():
                has_conf = True
        if has_conf: conf_slides += 1
    total_issues += len(issues)
    print(f"{path}: {len(list(prs.slides))} slides · {Emu(W).inches:.2f}x{Emu(H).inches:.2f}in · confidential-marked slides: {conf_slides}")
    for i in issues: print(i)
    if not issues: print("  ✓ no overflow")
print(f"\nTOTAL overflow issues across decks: {total_issues}")
