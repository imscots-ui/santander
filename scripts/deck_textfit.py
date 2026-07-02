#!/usr/bin/env python3
"""
Flag text boxes whose *rendered text* runs past the slide's safe area.

deck-qa.py checks shape bounds vs the slide; deck_overlap.py checks boxes vs
each other. Neither catches the case that bit us: a text box set to
auto-size SHAPE_TO_FIT_TEXT grows DOWNWARD as content is added, so the shape's
stored height looks fine but the actual text spills below the footer / off the
slide (the last rows vanish). This estimates rendered height (line ≈ 1.2×font ×
line-spacing, plus paragraph space before/after) and flags boxes whose content
bottom crosses into the footer zone.

Heuristic — eyeball flags with scripts/deck_shot.py. Usage: python3 scripts/deck_textfit.py
"""
from pptx import Presentation
from pptx.util import Emu
import glob, sys

FOOTER = 7.10   # inches — content should end above this (slide is 7.5in tall)

def para_h(p):
    sizes = [r.font.size.pt for r in p.runs if r.font.size]
    sz = max(sizes) if sizes else 12
    sb = p.space_before.pt if p.space_before else 0
    sa = p.space_after.pt if p.space_after else 0
    ls = p.line_spacing if isinstance(p.line_spacing, float) else 1.0
    return (sz * 1.2 * ls + sb + sa) / 72.0

flagged = 0
for deck in sorted(glob.glob('Santander_*.pptx')):
    prs = Presentation(deck)
    for i, s in enumerate(prs.slides):
        for sh in s.shapes:
            if not sh.has_text_frame or not sh.text_frame.text.strip():
                continue
            if len(sh.text_frame.paragraphs) < 3:
                continue  # single lines/labels aren't the growing-narration case
            top = Emu(sh.top).inches
            grows = str(sh.text_frame.auto_size) == 'SHAPE_TO_FIT_TEXT (1)'
            content_bottom = top + sum(para_h(p) for p in sh.text_frame.paragraphs)
            fixed_bottom = top + Emu(sh.height).inches
            bottom = content_bottom if grows else max(fixed_bottom, content_bottom)
            if bottom > FOOTER:
                flagged += 1
                print(f"OVERFLOW  {deck} s{i} [{sh.name}] "
                      f"top={top:.2f} estBottom={bottom:.2f}in "
                      f"({'auto-grows' if grows else 'clips'}) :: {sh.text_frame.text.strip()[:40]!r}")

print(f"\ntext boxes running into the footer zone (>{FOOTER}in): {flagged}")
sys.exit(1 if flagged else 0)
