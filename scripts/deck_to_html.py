#!/usr/bin/env python3
"""
Render a .pptx slide to standalone HTML for visual verification — no LibreOffice.

LibreOffice headless conversion is blocked in this sandbox (soffice/soffice.bin
IPC handshake fails), so this reads the deck with python-pptx and lays shapes out
as absolutely-positioned HTML that Chromium can screenshot. It is a faithful
proxy for the checks that matter — brand colour, text content, images/logo, table
and gross layout/overflow — not a pixel-perfect PowerPoint renderer (freeform
autoshapes render as their filled bounding box; gradients use their first stop).

Usage:  python3 scripts/deck_to_html.py <deck.pptx> <slide_index> > out.html
"""
import sys, base64, html
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

deck, idx = sys.argv[1], int(sys.argv[2]) if len(sys.argv) > 2 else 0
prs = Presentation(deck)
SW, SH = prs.slide_width, prs.slide_height
W = 1280
SCALE = W / SW
H = round(SH * SCALE)
def px(emu): return round((emu or 0) * SCALE)

def rgb_of(color):
    try:
        if color and color.type is not None and color.rgb is not None:
            return '#' + str(color.rgb)
    except Exception:
        pass
    return None

def fill_css(shape):
    try:
        f = shape.fill
        t = f.type
        if t == 1:  # solid
            c = rgb_of(f.fore_color)
            return f'background:{c};' if c else ''
        if t == 3:  # gradient — approximate with first stop
            try:
                stops = f._xPr.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}gs')
                cols = []
                for gs in stops:
                    srgb = gs.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                    if srgb is not None: cols.append('#' + srgb.get('val'))
                if len(cols) >= 2: return f'background:linear-gradient(135deg,{cols[0]},{cols[-1]});'
                if cols: return f'background:{cols[0]};'
            except Exception: pass
    except Exception:
        pass
    return ''

ALIGN = {1: 'left', 2: 'center', 3: 'right', 4: 'justify'}

def text_html(tf):
    out = []
    for p in tf.paragraphs:
        al = ALIGN.get(p.alignment, 'left')
        runs = []
        for r in p.runs:
            c = rgb_of(r.font.color) or '#1c1917'
            sz = (r.font.size.pt if r.font.size else 14)
            fpx = round(sz * 12700 * SCALE)
            w = '700' if r.font.bold else '400'
            it = 'italic' if r.font.italic else 'normal'
            fam = r.font.name or 'Geist, Arial, sans-serif'
            runs.append(f'<span style="color:{c};font-size:{fpx}px;font-weight:{w};'
                        f'font-style:{it};font-family:\'{fam}\',Arial,sans-serif">{html.escape(r.text)}</span>')
        out.append(f'<div style="text-align:{al};line-height:1.15">{"".join(runs) or "&nbsp;"}</div>')
    return ''.join(out)

def emit(shape, dx=0, dy=0):
    parts = []
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            gx, gy = px(shape.left) + dx, px(shape.top) + dy
            for sub in shape.shapes:
                parts.append(emit(sub, gx, gy))
            return ''.join(parts)
    except Exception:
        pass
    try:
        l, t, w, h = px(shape.left) + dx, px(shape.top) + dy, px(shape.width), px(shape.height)
    except Exception:
        return ''
    style = f'position:absolute;left:{l}px;top:{t}px;width:{w}px;height:{h}px;overflow:visible;'
    inner = ''
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            blob = shape.image.blob
            b64 = base64.b64encode(blob).decode()
            ct = shape.image.content_type
            inner = f'<img src="data:{ct};base64,{b64}" style="width:100%;height:100%;object-fit:contain">'
        elif shape.has_table:
            rows = []
            for row in shape.table.rows:
                cells = ''.join(f'<td style="border:1px solid #ccc;padding:2px 4px;font-size:{round(9*12700*SCALE)}px">{html.escape(c.text)}</td>' for c in row.cells)
                rows.append(f'<tr>{cells}</tr>')
            inner = f'<table style="width:100%;border-collapse:collapse">{"".join(rows)}</table>'
            style += fill_css(shape)
        else:
            style += fill_css(shape)
            if shape.has_text_frame and shape.text_frame.text.strip():
                inner = f'<div style="padding:4px">{text_html(shape.text_frame)}</div>'
    except Exception:
        pass
    return f'<div style="{style}">{inner}</div>'

slide = prs.slides[idx]
body = ''.join(emit(s) for s in slide.shapes)
print(f'''<!doctype html><html><head><meta charset="utf-8"><style>*{{margin:0;box-sizing:border-box}}</style></head>
<body><div style="position:relative;width:{W}px;height:{H}px;background:#fff;">{body}</div></body></html>''')
