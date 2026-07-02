"""
build_architecture.py
Generates a 12-slide technical architecture report for the Santander Business
Banking prototype using python-pptx.

Rebuilt to PowerPoint best-practice standards (REFERENCE.md §30–33, §92):
  · Panels, stat boxes, and cards — never raw bullet dumps
  · Dark / warm alternation for visual rhythm
  · One clear idea per slide; generous whitespace
  · Consistent slide furniture: red bar, title, red rule, footer chrome
  · INTERNAL · CONFIDENTIAL classification on every slide
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

# ── Colour tokens ───────────────────────────────────────────────────────────────
RED      = RGBColor(0xDA, 0x29, 0x1C)
RED_DARK = RGBColor(0xA0, 0x0D, 0x24)
DARK     = RGBColor(0x1C, 0x19, 0x17)
WARM     = RGBColor(0xFA, 0xF6, 0xEF)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
STONE1   = RGBColor(0xF5, 0xF5, 0xF4)
STONE2   = RGBColor(0xE7, 0xE5, 0xE4)
STONE3   = RGBColor(0xD6, 0xD3, 0xD1)
STONE5   = RGBColor(0x78, 0x71, 0x6C)
STONE7   = RGBColor(0x44, 0x40, 0x3C)
LIGHTRED = RGBColor(0xFE, 0xCA, 0xCA)
LIGHTRED2= RGBColor(0xFE, 0xF2, 0xF2)
GREEN    = RGBColor(0x05, 0x96, 0x69)
AMBER    = RGBColor(0xD9, 0x77, 0x06)
BLUE     = RGBColor(0x1D, 0x4E, 0xD8)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── Drawing helpers ─────────────────────────────────────────────────────────────

def R(s, l, t, w, h, fill=None, line=None, lw=0.75):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.line.fill.background()
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    return sh

def T(s, text, l, t, w, h, size=12, bold=False, color=DARK,
      align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    return tb

def TML(s, lines, l, t, w, h, size=11, color=DARK, leading=4, font="Calibri"):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False; p.space_after = Pt(leading)
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = font


# ── Slide furniture ─────────────────────────────────────────────────────────────

def warm_bg(s):
    R(s, 0, 0, 13.33, 7.5, fill=WARM)
    R(s, 0, 0, 13.33, 0.07, fill=RED)

def dark_bg(s):
    R(s, 0, 0, 13.33, 7.5, fill=DARK)
    R(s, 0, 0, 13.33, 0.07, fill=RED)

def title(s, text, sub=None, eyebrow=None, on_dark=False):
    tc = WHITE if on_dark else DARK
    y = 0.18
    if eyebrow:
        T(s, eyebrow.upper(), 0.5, 0.16, 12, 0.26, size=9.5, bold=True, color=RED)
        y = 0.40
    T(s, text, 0.5, y, 12.3, 0.55, size=26, bold=True, color=tc)
    if sub:
        T(s, sub, 0.5, y+0.52, 12.3, 0.3, size=11, color=STONE3 if on_dark else STONE5, italic=True)
    R(s, 0.5, (y+0.92 if sub else y+0.62), 12.33, 0.018, fill=RED if on_dark else STONE3)

def footer(s, n, is_cover=False):
    if is_cover:
        R(s, 10.95, 0.30, 1.95, 0.34, fill=RED)
        T(s, "INTERNAL · CONFIDENTIAL", 10.95, 0.34, 1.95, 0.26, size=7.5,
          bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        return
    R(s, 0, 7.22, 13.33, 0.28, fill=DARK)
    T(s, "INTERNAL · CONFIDENTIAL", 0.5, 7.24, 3.0, 0.22, size=8, bold=True,
      color=LIGHTRED, align=PP_ALIGN.LEFT)
    T(s, "Alan Davidson · Business Banking Advisor · Self-initiated · Own time · June 2026",
      3.5, 7.24, 8.6, 0.22, size=8, color=STONE3, align=PP_ALIGN.RIGHT, italic=True)
    T(s, str(n), 12.85, 7.24, 0.4, 0.22, size=8, color=STONE3, align=PP_ALIGN.RIGHT)


# ── Component helpers ───────────────────────────────────────────────────────────

def stat_box(s, num, label, x, y, w=2.7, h=1.55, bg=DARK, nc=WHITE, lc=STONE3,
             ns=40, ls=9.5):
    R(s, x, y, w, h, fill=bg)
    R(s, x, y, w, 0.05, fill=RED)
    T(s, num, x+0.1, y+0.12, w-0.2, h*0.56, size=ns, bold=True, color=nc, align=PP_ALIGN.CENTER)
    T(s, label, x+0.12, y+h*0.58, w-0.24, h*0.4, size=ls, color=lc, align=PP_ALIGN.CENTER)

def panel(s, title_txt, body, x, y, w, h, bg=WHITE, tc=DARK, bc=STONE5,
          accent=RED, line=STONE3, tsize=11, bsize=9):
    R(s, x, y, w, h, fill=bg, line=line)
    R(s, x, y, w, 0.055, fill=accent)
    T(s, title_txt, x+0.16, y+0.13, w-0.32, 0.32, size=tsize, bold=True, color=tc)
    if body:
        TML(s, body if isinstance(body, list) else [body],
            x+0.16, y+0.5, w-0.32, h-0.6, size=bsize, color=bc, leading=3)

def dark_panel(s, title_txt, body, x, y, w, h, accent=RED, tsize=11, bsize=9):
    panel(s, title_txt, body, x, y, w, h, bg=STONE7, tc=WHITE, bc=STONE3,
          accent=accent, line=None, tsize=tsize, bsize=bsize)

def numbered_card(s, num, label, desc, x, y, w, h, on_dark=False):
    bg = STONE7 if on_dark else WHITE
    R(s, x, y, w, h, fill=bg, line=None if on_dark else STONE3)
    R(s, x, y, 0.5, h, fill=RED)
    T(s, num, x, y+h/2-0.22, 0.5, 0.44, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    T(s, label, x+0.62, y+0.1, w-0.72, 0.3, size=10.5, bold=True,
      color=WHITE if on_dark else DARK)
    T(s, desc, x+0.62, y+0.40, w-0.72, h-0.46, size=7.8,
      color=STONE3 if on_dark else STONE5)

def flow_box(s, lines, x, y, w, h, accent=RED):
    R(s, x, y, w, h, fill=WHITE, line=STONE3)
    R(s, x, y, 0.06, h, fill=accent)
    TML(s, lines, x+0.22, y+0.12, w-0.34, h-0.2, size=10.5, color=DARK, leading=2)

def arrow_down(s, cx, y):
    T(s, "▼", cx-0.2, y, 0.4, 0.3, size=16, bold=True, color=RED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ════════════════════════════════════════════════════════════════════════════════
def slide_01_cover(prs):
    s = prs.slides.add_slide(BLANK)
    dark_bg(s)
    R(s, 0, 6.95, 13.33, 0.55, fill=RED_DARK)
    R(s, 0.55, 1.30, 0.5, 0.5, fill=RED)

    T(s, "Architecture Report", 0.55, 1.95, 11, 0.9, size=50, bold=True, color=WHITE)
    T(s, "Santander Business Banking Prototype", 0.55, 2.95, 11, 0.55, size=22, color=LIGHTRED)
    R(s, 0.55, 3.62, 8.5, 0.04, fill=RED)

    chips = [("5,700", "lines · 1 file"), ("12", "workflows"),
             ("7", "entity types"), ("16:9", "deck")]
    for i, (num, lbl) in enumerate(chips):
        cx = 0.55 + i * 2.7
        R(s, cx, 3.85, 2.5, 0.9, fill=STONE7)
        R(s, cx, 3.85, 2.5, 0.05, fill=RED)
        T(s, num, cx+0.12, 3.92, 2.3, 0.5, size=24, bold=True, color=WHITE)
        T(s, lbl, cx+0.14, 4.42, 2.2, 0.28, size=9, color=STONE3)

    T(s, "Technical deep-dive  ·  single-file React architecture  ·  state model  ·  entity & mandate logic  ·  design system",
      0.55, 5.05, 11.5, 0.4, size=11, color=STONE3, italic=True)
    T(s, "Alan Davidson  ·  Business Banking Advisor  ·  Self-initiated  ·  June 2026",
      0.55, 7.05, 11, 0.34, size=10, color=LIGHTRED)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE SINGLE-FILE ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════════
def slide_02_single_file(prs):
    s = prs.slides.add_slide(BLANK)
    warm_bg(s)
    title(s, "The Single-File Architecture", eyebrow="Why one file",
          sub="App.jsx — one React function component, read top-to-bottom")

    stats = [("5,700", "lines of code"), ("870 KB", "self-contained HTML"),
             ("190 KB", "gzipped — fast on 3G")]
    for i, (num, lbl) in enumerate(stats):
        stat_box(s, num, lbl, 0.5 + i*4.28, 1.35, w=4.05, h=1.55,
                 bg=DARK, ns=34, ls=10)

    panels = [
        ("Intentional monolith", "One file means the entire product reads top-to-bottom with no module navigation — ideal for a prototype meant to be understood at a glance."),
        ("No backend", "No real API calls except Google Fonts. All data is in-memory, fictional, and front-end only."),
        ("Single HTML output", "viteSingleFile inlines all JS and CSS into one portable file — runs from any browser, no install."),
        ("GitHub Pages", "Static hosting, zero server cost. Deployed and live at imscots-ui.github.io/santander."),
        ("The trade-off", "Monolith size vs. navigation simplicity — a deliberate choice for a demo/research context, not a production pattern."),
        ("Composition", "Sub-components are closures inside App — they share state directly without prop-drilling or context."),
    ]
    pw, ph = 4.05, 1.5
    for i, (t, b) in enumerate(panels):
        col, row = i % 3, i // 3
        panel(s, t, b, 0.5 + col*4.28, 3.18 + row*1.66, pw, ph)

    footer(s, 2)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — INTERNAL COMPONENT ORDER
# ════════════════════════════════════════════════════════════════════════════════
def slide_03_component_order(prs):
    s = prs.slides.add_slide(BLANK)
    dark_bg(s)
    title(s, "Internal Component Order", eyebrow="Internal structure",
          sub="The fixed top-to-bottom order inside the App() component", on_dark=True)

    items = [
        ("1", "Hooks", "120+ state variables, lines ~28–685"),
        ("2", "Static data", "ENTITY_INFO, accounts/mtdData/statementsData (useMemo)"),
        ("3", "Inline CSS", "css template literal: brand, animations, glass effects"),
        ("4", "Primitives", "ProgressDots, StepFrame, Input, Field, Toggle (closures)"),
        ("5", "Workflow renderers", "renderClosure … renderMtdSubmit (12 closures)"),
        ("6", "Screen components", "Home, Approve, Audit, Statements, MTD"),
        ("7", "Sheets / modals", "15 overlay sheets (OTP, Compliance, VoiceId, A11y …)"),
        ("8", "Main render", "desktop sidebar + mobile bottom-nav code paths"),
    ]
    cw, ch = 6.1, 1.18
    for i, (num, lbl, desc) in enumerate(items):
        col, row = i % 2, i // 2
        numbered_card(s, num, lbl, desc, 0.5 + col*6.28, 1.55 + row*1.3, cw, ch, on_dark=True)

    footer(s, 3)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — HOOK DISCIPLINE
# ════════════════════════════════════════════════════════════════════════════════
def slide_04_hook_discipline(prs):
    s = prs.slides.add_slide(BLANK)
    warm_bg(s)
    title(s, "Hook Discipline", eyebrow="The one hard rule",
          sub="Why closures cannot hold hooks — and how the boundary is enforced")

    # Big rule callout
    R(s, 0.5, 1.4, 12.33, 1.0, fill=RED)
    T(s, "No hook may appear after line ~744 — the closeWorkflow() boundary.",
      0.7, 1.52, 9.5, 0.45, size=17, bold=True, color=WHITE)
    T(s, "Everything below that line is a closure. Closures can't call useState / useEffect / useMemo.",
      0.7, 2.0, 11.8, 0.35, size=11, color=LIGHTRED)

    panels = [
        ("State at the top", "120+ state variables declared at the top of App (lines ~28–685), before any logic."),
        ("Closures, not components", "renderXxx functions and screen/sheet components are closures inside App — they read/write top-level state directly."),
        ("closeWorkflow() resets all", "Every workflow-specific state must be reset in closeWorkflow() — ghost state between workflows is a known failure mode."),
        ("useEffect, sparingly", "Only: 30s tick timer, font loading, scroll resets. No data fetching — there is no backend."),
        ("useMemo for derived data", "accounts and statementsData are memoised from entityType and payment history — recompute only when inputs change."),
        ("The awk caveat", "The /ship-ready hook flags hooks past line 120 — a known false positive at this file size; the real boundary is line ~744."),
    ]
    pw, ph = 4.05, 1.78
    for i, (t, b) in enumerate(panels):
        col, row = i % 3, i // 3
        panel(s, t, b, 0.5 + col*4.28, 2.7 + row*1.94, pw, ph)

    footer(s, 4)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — NAVIGATION MODEL
# ════════════════════════════════════════════════════════════════════════════════
def slide_05_navigation(prs):
    s = prs.slides.add_slide(BLANK)
    warm_bg(s)
    title(s, "Navigation Model", eyebrow="State drives the screen",
          sub="Three layers of state compose every view")

    # Flow column (left)
    flow_box(s, ["tab  —  main screen",
                 "'home' · 'approve' · 'audit' · 'mtd' · 'statements'"],
             0.5, 1.5, 7.7, 1.05, accent=RED)
    arrow_down(s, 4.35, 2.58)
    flow_box(s, ["workflow  —  full-screen overlay",
                 "null · closure · biz · mandate · wages · lending · fx ·",
                 "dormancy · unlink · ringfence · idcheck · mtd-submit · complaint"],
             0.5, 2.95, 7.7, 1.5, accent=AMBER)
    arrow_down(s, 4.35, 4.48)
    flow_box(s, ["step  —  0-based position within the active workflow"],
             0.5, 4.85, 7.7, 0.85, accent=GREEN)

    # Side panel (right)
    panel(s, "viewMode",
          ["'mobile' (default) | 'desktop'",
           "Toggles bottom-nav vs. sidebar layout.",
           "",
           "All state is shared between shells —",
           "switching mobile ↔ desktop preserves any",
           "workflow already in progress."],
          8.5, 1.5, 4.33, 2.95, bsize=10.5)
    panel(s, "Why this matters",
          ["A workflow overlay renders on top of the",
           "current screen; closing it returns you",
           "exactly where you were. One state tree,",
           "two render paths, zero duplication."],
          8.5, 4.6, 4.33, 1.1, accent=BLUE, bsize=10.5)

    footer(s, 5)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ENTITY SYSTEM
# ════════════════════════════════════════════════════════════════════════════════
def slide_06_entity_system(prs):
    s = prs.slides.add_slide(BLANK)
    warm_bg(s)
    title(s, "Entity System", eyebrow="One switch reshapes everything",
          sub="ENTITY_INFO (~line 264) — entityType drives all derived config")

    entities = ["Sole Trader", "Partnership", "Limited", "LLP", "Charity", "Club", "Society"]
    for i, e in enumerate(entities):
        x = 0.5 + i*1.79
        R(s, x, 1.45, 1.65, 0.6, fill=DARK)
        R(s, x, 1.45, 1.65, 0.05, fill=RED)
        T(s, e, x+0.05, 1.58, 1.55, 0.4, size=10.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    T(s, "Switching entityType re-derives all of the following:", 0.5, 2.3, 12, 0.3,
      size=11, bold=True, color=STONE7)

    panels = [
        ("Entity object", "Labels (‘director’ vs ‘trustee’), required documents, and the principal terminology used throughout the UI."),
        ("Registration numbers", "Companies House number, LLP number, or Charity Commission number — shown and validated per type."),
        ("accounts array", "useMemo (~line 334) returns a different set of accounts and mandate rules for each entity type."),
        ("Mandate defaults", "Sole trader → any-1; limited/LLP → any-2; charity/club/society → all signatories required."),
        ("Board minutes", "Clubs, charities, and societies require board-minute upload on every change; companies do not."),
        ("UI copy", "Workflow steps, compliance prompts, and document checklists all re-render from the active entity."),
    ]
    pw, ph = 4.05, 1.5
    for i, (t, b) in enumerate(panels):
        col, row = i % 3, i // 3
        panel(s, t, b, 0.5 + col*4.28, 2.7 + row*1.66, pw, ph)

    footer(s, 6)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — MANDATE RULES
# ════════════════════════════════════════════════════════════════════════════════
def slide_07_mandate_rules(prs):
    s = prs.slides.add_slide(BLANK)
    warm_bg(s)
    title(s, "Mandate Rules", eyebrow="Who must authorise",
          sub="getMandateFor(accountNos) — line ~420 — picks the strictest rule across selected accounts")

    rules = [
        ("any-1", "Single authoriser", "One signatory is sufficient. No co-signer, no cooling-off.", GREEN),
        ("any-2", "Two authorisers", "Two signatories required → triggers the co-signer workflow and 48h window.", AMBER),
        ("all", "Every signatory", "All signatories must approve → 24h cooling-off period is enforced.", RED),
    ]
    cw = 4.05
    for i, (code, name, desc, col) in enumerate(rules):
        x = 0.5 + i*4.28
        R(s, x, 1.5, cw, 2.5, fill=WHITE, line=STONE3)
        R(s, x, 1.5, cw, 0.7, fill=col)
        T(s, code, x+0.2, 1.6, cw-0.4, 0.5, size=22, bold=True, color=WHITE)
        T(s, name, x+0.2, 2.35, cw-0.4, 0.4, size=13, bold=True, color=DARK)
        T(s, desc, x+0.2, 2.85, cw-0.4, 1.0, size=10.5, color=STONE5)

    # Strictness + impact callout
    R(s, 0.5, 4.3, 12.33, 0.7, fill=STONE1, line=STONE3)
    T(s, "Strictness order:   all  >  any-2  >  any-1", 0.72, 4.45, 6, 0.4,
      size=13, bold=True, color=DARK)
    T(s, "Select multiple accounts → the toughest rule wins.", 6.8, 4.47, 5.8, 0.4,
      size=11, color=STONE5, italic=True)

    R(s, 0.5, 5.25, 12.33, 0.95, fill=RED)
    T(s, "Impact", 0.72, 5.36, 2, 0.35, size=12, bold=True, color=WHITE)
    T(s, "The chosen rule determines whether a workflow requires a co-signer and/or a 24-hour cooling-off period before execution — the core control behind dual authorisation.",
      0.72, 5.68, 12, 0.45, size=11, color=LIGHTRED)

    footer(s, 7)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — COOLING-OFF & STALLED
# ════════════════════════════════════════════════════════════════════════════════
def slide_08_cooling_stalled(prs):
    s = prs.slides.add_slide(BLANK)
    warm_bg(s)
    title(s, "Cooling-off & Stalled Requests", eyebrow="Time-based safeguards",
          sub="Two live mechanisms that keep high-risk actions reversible")

    # Left — Cooling-off
    R(s, 0.5, 1.5, 6.05, 4.9, fill=WHITE, line=STONE3)
    R(s, 0.5, 1.5, 6.05, 0.6, fill=RED)
    T(s, "Cooling-off", 0.7, 1.6, 5.6, 0.42, size=16, bold=True, color=WHITE)
    TML(s, [
        "cooling array — requests pending a 24-hour",
        "execution delay (working days only, bank",
        "holidays skipped).",
        "",
        "tick state refreshes every 30 seconds —",
        "keeps the progress bars live and accurate.",
        "",
        "The user can Cancel or Fast-Track any",
        "request while it sits in the cooling window.",
        "",
        "Regulatory basis: FCA BCOBS 4A.",
    ], 0.72, 2.3, 5.6, 3.9, size=11.5, color=STONE7, leading=4)

    # Right — Stalled
    R(s, 6.78, 1.5, 6.05, 4.9, fill=WHITE, line=STONE3)
    R(s, 6.78, 1.5, 6.05, 0.6, fill=STONE7)
    T(s, "Stalled", 6.98, 1.6, 5.6, 0.42, size=16, bold=True, color=WHITE)
    TML(s, [
        "stalled array — a co-signer missed the",
        "approval window.",
        "",
        "The system auto-escalates to the",
        "Relationship Manager (Priya Desai).",
        "",
        "The UI shows a “Priya’s on it” card with a",
        "direct ‘Call Priya’ action button.",
        "",
        "Escalation path:",
        "failed co-sign → RM notification → manual",
        "override and resolution.",
    ], 7.0, 2.3, 5.6, 3.9, size=11.5, color=STONE7, leading=4)

    footer(s, 8)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — WORKFLOW ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════════════
def slide_09_workflows(prs):
    s = prs.slides.add_slide(BLANK)
    warm_bg(s)
    title(s, "Workflow Architecture", eyebrow="12 step-based wizards",
          sub="Each is a renderXxx closure reading and writing top-level App state")

    wf = [
        ("renderClosure", "~1021", "Account closure + compliance"),
        ("renderBiz", "~1332", "Business KYC / KYB verification"),
        ("renderMandate", "~1428", "Signatory mandate changes"),
        ("renderWages", "~1630", "Bulk wages (CSV + approval)"),
        ("renderLending", "~1915", "Business lending application"),
        ("renderFX", "~2018", "Foreign exchange payments"),
        ("renderDormancy", "~2149", "Dormant account reactivation"),
        ("renderUnlink", "~2177", "Personal/business unlinking"),
        ("renderRingfence", "~2310", "Credit ring-fencing"),
        ("renderIdCheck", "~2390", "Identity / Voice ID checks"),
        ("renderMtdSubmit", "~4998", "Making Tax Digital submission"),
        ("renderComplaint", "new", "FCA DISP complaint handling"),
    ]
    cw, ch = 4.05, 1.12
    for i, (name, line, desc) in enumerate(wf):
        col, row = i % 3, i // 3
        x = 0.5 + col*4.28
        y = 1.5 + row*1.27
        R(s, x, y, cw, ch, fill=WHITE, line=STONE3)
        R(s, x, y, 0.055, ch, fill=RED)
        T(s, name, x+0.18, y+0.12, cw-1.1, 0.3, size=11, bold=True, color=DARK, font="Consolas")
        T(s, line, x+cw-0.95, y+0.12, 0.85, 0.26, size=8.5, bold=True, color=RED, align=PP_ALIGN.RIGHT)
        T(s, desc, x+0.18, y+0.5, cw-0.3, 0.5, size=9, color=STONE5)

    footer(s, 9)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — DESIGN SYSTEM
# ════════════════════════════════════════════════════════════════════════════════
def slide_10_design_system(prs):
    s = prs.slides.add_slide(BLANK)
    warm_bg(s)
    title(s, "Design System", eyebrow="Bosun's Law",
          sub="Brand tokens, spacing scale, and the colour rules enforced on every commit")

    # Colour swatches
    swatches = [("#EC0000", "Brand red", RED, WHITE),
                ("#FBF1EA", "Warm bg", WARM, DARK),
                ("#1c1917", "Stone-900", DARK, WHITE),
                ("#78716c", "Stone-500", STONE5, WHITE)]
    for i, (hexv, name, col, tc) in enumerate(swatches):
        x = 0.5 + i*3.09
        R(s, x, 1.45, 2.95, 0.95, fill=col, line=STONE3)
        T(s, hexv, x+0.15, 1.55, 2.6, 0.35, size=13, bold=True, color=tc)
        T(s, name, x+0.15, 1.95, 2.6, 0.3, size=9.5, color=tc)

    # Spacing scale strip
    R(s, 0.5, 2.6, 12.33, 0.75, fill=STONE1, line=STONE3)
    T(s, "Spacing scale", 0.7, 2.7, 2.2, 0.3, size=11, bold=True, color=RED)
    T(s, "4 · 8 · 12 · 16 · 24 · 32 · 48 · 64 · 96 · 128 px      →      Tailwind  1 · 2 · 3 · 4 · 6 · 8 · 12 · 16 · 24 · 32",
      2.9, 2.78, 9.8, 0.45, size=12, bold=True, color=DARK, font="Consolas")

    panels = [
        ("Typography", "Fraunces — display headings (Google Fonts). Geist — body. Geist Mono — account numbers and figures."),
        ("Colour law", "Never text-gray-* or text-zinc-* — always text-stone-*. Never stone-400/500 on dark or red surfaces."),
        ("One primary CTA", "Exactly one bg-[#EC0000] or bg-stone-900 primary button per view — hierarchy stays clear."),
        ("Tabular figures", "num-tab class (font-variant-numeric: tabular-nums) on every monetary amount for clean alignment."),
        ("Inline CSS", "css template literal (~line 779): hero-card, cool-card, anim-fade, anim-slide, shimmer, stagger-1…7."),
        ("Enforcement", "Bosun station in /ship-ready scans for spacing, colour, and CTA violations before every push."),
    ]
    pw, ph = 4.05, 1.4
    for i, (t, b) in enumerate(panels):
        col, row = i % 3, i // 3
        panel(s, t, b, 0.5 + col*4.28, 3.55 + row*1.56, pw, ph)

    footer(s, 10)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — TECH STACK
# ════════════════════════════════════════════════════════════════════════════════
def slide_11_tech_stack(prs):
    s = prs.slides.add_slide(BLANK)
    warm_bg(s)
    title(s, "Tech Stack", eyebrow="Build & deploy",
          sub="Minimal, modern, front-end only — no backend, no database, no API server")

    # Left column — Frontend
    T(s, "FRONTEND", 0.5, 1.45, 6, 0.3, size=10, bold=True, color=RED)
    fe = [
        ("React 18.3", "Single large function component; closures as sub-components."),
        ("Tailwind CSS 3", "Utility-first styling with a custom stone-* palette."),
        ("Lucide React", "Icon library used throughout the UI."),
    ]
    for i, (t, b) in enumerate(fe):
        panel(s, t, b, 0.5, 1.8 + i*1.55, 6.05, 1.4)

    # Right column — Build & deploy
    T(s, "BUILD & DEPLOY", 6.78, 1.45, 6, 0.3, size=10, bold=True, color=RED)
    be = [
        ("Vite 5 + viteSingleFile", "Fast dev server; inlines all JS/CSS into one HTML file."),
        ("Output", "~870 KB HTML · ~190 KB gzip · GitHub Pages static hosting."),
        ("Commands", "npm run dev · npm run build · npm run preview."),
    ]
    for i, (t, b) in enumerate(be):
        panel(s, t, b, 6.78, 1.8 + i*1.55, 6.05, 1.4, accent=BLUE)

    # Bottom note
    R(s, 0.5, 6.45, 12.33, 0.6, fill=STONE7)
    T(s, "python-pptx is documentation tooling only — it builds these decks, it is not part of the app.",
      0.7, 6.55, 12, 0.4, size=10.5, color=WHITE, italic=True)

    footer(s, 11)


# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — SECURITY & ACCESSIBILITY
# ════════════════════════════════════════════════════════════════════════════════
def slide_12_security(prs):
    s = prs.slides.add_slide(BLANK)
    warm_bg(s)
    title(s, "Security & Accessibility Decisions", eyebrow="Standing orders",
          sub="The non-negotiable rules enforced before every commit")

    panels = [
        ("No dangerous HTML", "No dangerouslySetInnerHTML on dynamic content — hardcoded static strings only."),
        ("Focus visible", "No focus:outline-none without a focus-visible: companion — every bare removal is an a11y breach."),
        ("WCAG 2.1 AA", "Contrast ≥ 4.5:1, focus rings preserved for keyboard navigation, ARIA labels where needed."),
        ("No secrets in source", "No real credentials, account numbers, API keys, or emails beyond the author credit."),
        ("No force push", "git push --force is forbidden — ever. History stays intact."),
        ("/ship-ready gate", "Seven-station pre-commit checklist — one RED finding and the ship stands down to fix."),
        ("Fictional data only", "Every account number, name, and amount in the prototype is invented."),
        ("One external call", "Google Fonts is the only external network request in production — nothing else leaves the browser."),
    ]
    pw, ph = 6.05, 1.18
    for i, (t, b) in enumerate(panels):
        col, row = i % 2, i // 2
        panel(s, t, b, 0.5 + col*6.28, 1.5 + row*1.3, pw, ph, bsize=9.5)

    footer(s, 12)


# ── Main ────────────────────────────────────────────────────────────────────────
def main():
    output_path = "/home/user/santander/Santander_Architecture.pptx"
    builders = [
        ("Slide 1  — Cover",                    slide_01_cover),
        ("Slide 2  — Single-File Architecture", slide_02_single_file),
        ("Slide 3  — Internal Component Order",  slide_03_component_order),
        ("Slide 4  — Hook Discipline",           slide_04_hook_discipline),
        ("Slide 5  — Navigation Model",          slide_05_navigation),
        ("Slide 6  — Entity System",             slide_06_entity_system),
        ("Slide 7  — Mandate Rules",             slide_07_mandate_rules),
        ("Slide 8  — Cooling-off & Stalled",     slide_08_cooling_stalled),
        ("Slide 9  — Workflow Architecture",     slide_09_workflows),
        ("Slide 10 — Design System",             slide_10_design_system),
        ("Slide 11 — Tech Stack",                slide_11_tech_stack),
        ("Slide 12 — Security & Accessibility",  slide_12_security),
    ]
    for label, builder in builders:
        try:
            builder(prs)
            print(f"  [OK] {label}")
        except Exception as exc:
            print(f"  [ERROR] {label}: {exc}")
            import traceback; traceback.print_exc()

    # Cover classification tag
    footer(prs.slides[0], 1, is_cover=True)

    prs.save(output_path)
    print(f"\nSaved: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
