"""
build_future_v2.py
Best-quality rebuild of Santander_Digital_Banking_Future.pptx

This is the EXECUTIVE PITCH DECK. Audience: senior management, product directors,
potential stakeholders. Story: problem → prototype → evidence → business case → ask.

Every headline is an assertion that works alone.
Every slide has one point. No bullet dumps.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

RED    = RGBColor(0xC8, 0x10, 0x2E)
DARK   = RGBColor(0x1C, 0x19, 0x17)
DARK2  = RGBColor(0x29, 0x24, 0x20)
WARM   = RGBColor(0xFA, 0xF6, 0xEF)
STONE1 = RGBColor(0xF5, 0xF0, 0xE8)
STONE2 = RGBColor(0xE7, 0xE5, 0xE4)
STONE5 = RGBColor(0x78, 0x71, 0x6C)
STONE7 = RGBColor(0x44, 0x40, 0x3C)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
AMBER  = RGBColor(0xD9, 0x77, 0x06)
AMBERLT= RGBColor(0xFE, 0xF3, 0xC7)
BLUE   = RGBColor(0x25, 0x63, 0xEB)
EMRLD  = RGBColor(0x05, 0x96, 0x69)
EMRLDLT= RGBColor(0xD1, 0xFA, 0xE5)
REDLT  = RGBColor(0xFE, 0xF2, 0xF2)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def rect(slide, l, t, w, h, fill=None, lc=None, lpt=0.75):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else:    s.fill.background()
    if lc:   s.line.color.rgb = lc; s.line.width = Pt(lpt)
    else:    s.line.fill.background()
    return s

def txt(slide, s, l, t, w, h, sz=14, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = s
    r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = "Calibri"
    return tb

def ml(slide, lines, l, t, w, h, sz=12, color=DARK, sp=5):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line
        r.font.size = Pt(sz); r.font.color.rgb = color
        r.font.name = "Calibri"; p.space_after = Pt(sp)
    return tb

def bars(slide):
    rect(slide, 0, 0, 13.33, 0.09, fill=RED)
    rect(slide, 0, 7.41, 13.33, 0.09, fill=RED)

def rule_line(slide, l=0.55, t=1.08, w=12.23):
    rect(slide, l, t, w, 0.035, fill=RED)

# ── SLIDE 1 — COVER ───────────────────────────────────────────────────────────
def s1_cover():
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
    bars(slide)
    # Large red circle (decorative orb)
    rect(slide, 8.5, -0.5, 5.5, 5.5, fill=RGBColor(0xA0, 0x0A, 0x24))
    rect(slide, 9.5, 0.5, 4.0, 4.0, fill=DARK2)
    rect(slide, 0.55, 1.7, 0.18, 3.8, fill=RED)
    txt(slide, "SANTANDER BUSINESS BANKING", 0.9, 1.8, 10, 0.4,
        sz=9, bold=True, color=RED)
    txt(slide, "Digital.\nPersonal.\nBuilt from within.", 0.9, 2.25, 9.5, 2.6,
        sz=52, bold=True, color=WHITE)
    txt(slide, "An executive briefing on the prototype — and what comes next.",
        0.9, 5.05, 10, 0.55, sz=15, color=RGBColor(0xC0, 0xBB, 0xB5))
    txt(slide, "Alan Davidson  ·  Alan.Davidson@santander.co.uk  ·  June 2026",
        0.9, 6.85, 11, 0.38, sz=10, color=STONE5)

# ── SLIDE 2 — THE NUMBER ──────────────────────────────────────────────────────
def s2_number():
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
    bars(slide)
    txt(slide, "THE OPPORTUNITY", 0.55, 0.18, 12, 0.35, sz=9, bold=True, color=RED)
    txt(slide, "£137M", 0.55, 0.55, 12, 2.4, sz=108, bold=True, color=WHITE,
        align=PP_ALIGN.LEFT)
    txt(slide, "Annualised revenue opportunity from digitising business banking admin\nfor Santander's 280,000 UK SMB customers.",
        0.55, 3.1, 11, 0.8, sz=18, color=RGBColor(0xC0, 0xBB, 0xB5))
    rect(slide, 0.55, 4.1, 12.23, 0.04, fill=RED)
    stats = [
        ("280,000", "UK SMB customers currently held at Santander"),
        ("£490", "average annual revenue per business customer"),
        ("5–7×", "expected uplift per customer from self-service digitisation"),
    ]
    for i, (val, lbl) in enumerate(stats):
        x = 0.55 + i * 4.26
        rect(slide, x + 0.05, 4.25, 4.1, 2.65, fill=DARK2)
        rect(slide, x + 0.05, 4.25, 0.09, 2.65, fill=RED)
        txt(slide, val, x + 0.3, 4.4, 3.7, 1.0, sz=40, bold=True, color=WHITE)
        txt(slide, lbl, x + 0.3, 5.55, 3.7, 0.9, sz=11,
            color=RGBColor(0x90, 0x8A, 0x84))

# ── SLIDE 3 — THE PROBLEM ─────────────────────────────────────────────────────
def s3_problem():
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, fill=WARM)
    bars(slide)
    txt(slide, "THE PROBLEM", 0.55, 0.18, 12, 0.35, sz=9, bold=True, color=RED)
    txt(slide, "Four things a director needs to do every quarter.\nFour things that still require a letter to Sunderland.",
        0.55, 0.42, 12.2, 0.95, sz=26, bold=True, color=DARK)
    rule_line(slide)
    pains = [
        ("01", "Change a mandate", "Add or remove a signatory. Form goes to Sunderland by post.\nAverage: 5 working days. No app. No visibility.", AMBER),
        ("02", "File a VAT return", "Making Tax Digital requires bridging software.\nThe bank holds the transaction data — but is not connected to HMRC.", RED),
        ("03", "Pay the team", "Bulk BACS payments need a separate system or a phone call.\nThe bank has the payees — but doesn't offer the workflow.", BLUE),
        ("04", "Renew KYC / KYB", "Annual identity checks and Companies House confirmations.\nAnother form. Another envelope. Another wait.", EMRLD),
    ]
    W = 12.23 / 4
    for i, (num, title, body, accent) in enumerate(pains):
        x = 0.55 + i * W
        rect(slide, x + 0.06, 1.25, W - 0.12, 5.6, fill=WHITE, lc=STONE2)
        rect(slide, x + 0.06, 1.25, W - 0.12, 0.09, fill=accent)
        txt(slide, num, x + 0.18, 1.45, W - 0.3, 0.6,
            sz=32, bold=True, color=accent)
        txt(slide, title, x + 0.18, 2.1, W - 0.3, 0.65,
            sz=14, bold=True, color=DARK)
        ml(slide, body.split('\n'), x + 0.18, 2.82, W - 0.3, 3.8,
           sz=11, color=STONE7, sp=6)

# ── SLIDE 4 — THE INSIGHT ─────────────────────────────────────────────────────
def s4_insight():
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
    bars(slide)
    txt(slide, "THE INSIGHT", 0.55, 0.18, 12, 0.35, sz=9, bold=True, color=RED)
    txt(slide, "We already own the relationship.\nWe just need to digitise it.",
        0.55, 0.55, 12, 2.1, sz=44, bold=True, color=WHITE)
    rect(slide, 0.55, 2.8, 12.23, 0.04, fill=RED)
    points = [
        ("The data is already there",
         "Santander holds every transaction, every payee, every mandate.\n"
         "The prototype uses that data — showing what's possible when the\n"
         "bank becomes the interface, not just the ledger."),
        ("The regulation is already met",
         "FCA BCOBS 4A cooling-off is built in.\n"
         "HMRC MTD API integration is modelled.\n"
         "Companies House verification is in the workflow.\n"
         "Compliance is not an afterthought — it's the product."),
        ("The technology is already proven",
         "React. Vite. Tailwind. GitHub Pages.\n"
         "The prototype runs on any mobile browser — right now.\n"
         "No new infrastructure. No new vendor relationships.\n"
         "One self-contained HTML file."),
    ]
    for i, (title, body) in enumerate(points):
        x = 0.55 + i * 4.26
        rect(slide, x + 0.05, 3.0, 4.1, 3.9, fill=DARK2)
        txt(slide, title, x + 0.22, 3.15, 3.75, 0.55,
            sz=12, bold=True, color=RED)
        ml(slide, body.split('\n'), x + 0.22, 3.76, 3.75, 3.0,
           sz=11, color=RGBColor(0xB0, 0xAA, 0xA4), sp=5)

# ── SLIDE 5 — THE PROTOTYPE ───────────────────────────────────────────────────
def s5_prototype():
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, fill=WARM)
    bars(slide)
    txt(slide, "THE PROTOTYPE  ·  LIVE TODAY", 0.55, 0.18, 12, 0.35,
        sz=9, bold=True, color=RED)
    txt(slide, "Phase 1 is done. The prototype is live.\nOpen it on your phone right now.",
        0.55, 0.42, 12.2, 0.95, sz=26, bold=True, color=DARK)
    rule_line(slide)
    txt(slide, "imscots-ui.github.io/santander",
        0.55, 1.22, 12.23, 0.55, sz=20, bold=True, color=RED)
    # Stats boxes
    stats = [
        ("5,700", "lines of production-quality\nReact code", RED),
        ("11",    "complete paperless\nworkflows", AMBER),
        ("7",     "business entity types\nsupported", BLUE),
        ("0",     "paper forms needed\nfor any covered task", EMRLD),
        ("1",     "engineer\nself-initiated, spare time", STONE5),
    ]
    W = 12.23 / 5
    for i, (val, lbl, ac) in enumerate(stats):
        x = 0.55 + i * W
        rect(slide, x + 0.04, 2.05, W - 0.08, 2.7, fill=WHITE, lc=STONE2)
        rect(slide, x + 0.04, 2.05, W - 0.08, 0.08, fill=ac)
        txt(slide, val, x + 0.15, 2.22, W - 0.25, 1.1,
            sz=44, bold=True, color=DARK)
        ml(slide, lbl.split('\n'), x + 0.15, 3.38, W - 0.25, 1.25,
           sz=10, color=STONE5, sp=3)
    rect(slide, 0.55, 5.0, 12.23, 2.1, fill=DARK)
    txt(slide, "WHAT'S IN THE PROTOTYPE", 0.75, 5.12, 12, 0.32,
        sz=9, bold=True, color=RED)
    ml(slide, [
        "Home Command Centre  ·  15 proactive signals  ·  personalized greeting  ·  total balance hero  ·  13-week forecast chart",
        "Account Closure  ·  Mandate Change  ·  Bulk Payments  ·  FX International  ·  MTD VAT Submission  ·  KYC/KYB ID Check",
        "Cooling-off (FCA BCOBS 4A)  ·  Stalled request / RM escalation  ·  OTP confirmation  ·  Biometric authentication",
        "Voice memo expense tagging  ·  Receipt scanning  ·  Payment sequencer optimiser  ·  Accessibility settings (A11y)",
    ], 0.75, 5.5, 12.0, 1.5, sz=10, color=RGBColor(0xB0, 0xAA, 0xA4), sp=4)

# ── SLIDE 6 — THE COMMAND CENTRE ─────────────────────────────────────────────
def s6_command_centre():
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, fill=WARM)
    bars(slide)
    txt(slide, "THE COMMAND CENTRE  ·  THE FLAGSHIP FEATURE", 0.55, 0.18, 12, 0.35,
        sz=9, bold=True, color=RED)
    txt(slide, "The Home screen that tells you what to do — before you think to ask.",
        0.55, 0.42, 12.2, 0.65, sz=26, bold=True, color=DARK)
    rule_line(slide, t=1.15)
    # Left dark panel
    rect(slide, 0.55, 1.28, 5.6, 5.85, fill=DARK)
    txt(slide, "PREDICT", 0.75, 1.45, 5.2, 0.38, sz=10, bold=True, color=RED)
    ml(slide, [
        "Payroll run · in 3 days · £42,180",
        "VAT return due · in 22 days · £4,118",
        "Companies House confirmation · 14 days",
        "",
        "Colour-coded by urgency.",
        "Amber · Red · Blue.",
        "Each is a live nudge card.",
    ], 0.75, 1.88, 5.2, 2.0, sz=11, color=RGBColor(0xB0, 0xAA, 0xA4), sp=5)
    txt(slide, "EXPLAIN", 0.75, 4.0, 5.2, 0.38, sz=10, bold=True, color=AMBER)
    ml(slide, [
        "'Balance covers it — review if payees changed'",
        "'Q3 Making Tax Digital · HMRC submission'",
        "'Companies House · Northgate Systems Ltd'",
    ], 0.75, 4.43, 5.2, 1.0, sz=11, color=RGBColor(0xB0, 0xAA, 0xA4), sp=5)
    txt(slide, "ACT", 0.75, 5.55, 5.2, 0.38, sz=10, bold=True, color=EMRLD)
    ml(slide, [
        "One tap → opens the relevant workflow directly.",
        "No navigation. No menu. No searching.",
    ], 0.75, 5.98, 5.2, 0.7, sz=11, color=RGBColor(0xB0, 0xAA, 0xA4), sp=5)
    # Right content
    txt(slide, "Why this is different",
        6.4, 1.38, 6.6, 0.5, sz=15, bold=True, color=DARK)
    ml(slide, [
        "Conventional banking apps show you data.",
        "This one tells you what to do with it.",
        "",
        "The three forecast cards replace:",
        "  •  The payroll reminder email from the accountant",
        "  •  The tax calendar the director keeps in their head",
        "  •  The Companies House letter that gets filed and forgotten",
        "",
        "Each card surfaces a specific, time-bounded, actionable signal.",
        "Each taps directly into the resolution — no navigation required.",
        "",
        "'Predict → Explain → Act' — applied to banking.",
        "It's not a feature. It's a design principle.",
        "",
        "The Home screen knows:",
        "  •  When payroll is due (wages workflow history)",
        "  •  How much VAT is owed (transaction categorisation)",
        "  •  When Companies House confirmation is due (entity setup)",
        "",
        "This information already exists in the bank.",
        "The question is whether the bank surfaces it.",
    ], 6.4, 1.95, 6.5, 5.2, sz=10.5, color=STONE7, sp=4)

# ── SLIDE 7 — MTD / HMRC ─────────────────────────────────────────────────────
def s7_mtd():
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
    bars(slide)
    txt(slide, "MAKING TAX DIGITAL  ·  THE BANK AS TAX AGENT", 0.55, 0.18, 12, 0.35,
        sz=9, bold=True, color=RED)
    txt(slide, "From receipt scan to HMRC submission — without leaving the bank app.",
        0.55, 0.52, 12, 0.7, sz=28, bold=True, color=WHITE)
    rect(slide, 0.55, 1.35, 12.23, 0.04, fill=RED)
    steps = [
        ("01\nCAPTURE", "Transactions auto-categorised\nfrom bank feed.\nVoice memo records expenses\nby speaking.\nReceipt scan via camera."),
        ("02\nREVIEW", "Director reviews each transaction.\nCategories pre-filled.\nManual override available.\n47 transactions · Q3 2026."),
        ("03\nCALCULATE", "Output tax: £18,472.40\nInput tax: £14,354.40\nVAT due: £4,118.00\nPeriod: Q3 Jul–Sep 2026."),
        ("04\nDECLARE", "Statutory HMRC declaration.\nCheckbox. Timestamp.\nFull audit trail entry.\nDirector's name recorded."),
        ("05\nSUBMIT", "5-second countdown overlay.\nCalm grey Cancel button.\nAt zero: files to HMRC.\nReceipt # in audit log."),
    ]
    W = 12.23 / 5
    for i, (title, body) in enumerate(steps):
        x = 0.55 + i * W
        rect(slide, x + 0.04, 1.52, W - 0.08, 5.35, fill=DARK2)
        if i < 4:
            rect(slide, x + W - 0.04, 2.9, 0.12, 0.5, fill=RED)
            txt(slide, "→", x + W - 0.04, 2.85, 0.3, 0.5, sz=14, bold=True, color=RED)
        txt(slide, title, x + 0.18, 1.68, W - 0.28, 0.9,
            sz=13, bold=True, color=WHITE)
        ml(slide, body.split('\n'), x + 0.18, 2.65, W - 0.28, 4.0,
           sz=10, color=RGBColor(0xB0, 0xAA, 0xA4), sp=5)
    txt(slide, "No spreadsheet. No bridging software. No separate app. The bank files the VAT return.",
        0.55, 7.08, 12.23, 0.3, sz=10, italic=True, color=STONE5)

# ── SLIDE 8 — COMPLIANCE ─────────────────────────────────────────────────────
def s8_compliance():
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, fill=WARM)
    bars(slide)
    txt(slide, "REGULATORY COMPLIANCE  ·  BUILT IN, NOT BOLTED ON",
        0.55, 0.18, 12, 0.35, sz=9, bold=True, color=RED)
    txt(slide, "Every feature is tied to a specific regulation. No vague compliance.",
        0.55, 0.42, 12.2, 0.65, sz=26, bold=True, color=DARK)
    rule_line(slide, t=1.15)
    regs = [
        ("FCA BCOBS 4A", "Cooling-Off",
         "High-risk requests — account closures, certain mandate changes —\n"
         "carry a mandatory 24-hour execution delay.\n"
         "The Home screen shows a live progress bar.\n"
         "The customer can cancel at any point.", RED),
        ("HMRC MTD API", "VAT submission",
         "Making Tax Digital — mandatory for VAT-registered businesses.\n"
         "The prototype models the API contract:\n"
         "transaction categorisation → calculation → declaration → submission.\n"
         "Receipt #VAT2026Q3-9482 in audit log.", AMBER),
        ("Companies House", "KYC / KYB",
         "Annual confirmation statement.\n"
         "Director identity verification.\n"
         "Charity Commission integration for charity entity type.\n"
         "ID register: Lists 1, 2 & 3 (JMLSG guidance).", BLUE),
        ("PSD2 / Open Banking", "Account unlinking",
         "Personal/business account separation.\n"
         "Full FCA BCOBS disclosure before unlinking.\n"
         "Ring-fencing opt-in with regulatory declaration.\n"
         "All separation actions require OTP.", EMRLD),
    ]
    W = 12.23 / 4
    for i, (reg, title, body, ac) in enumerate(regs):
        x = 0.55 + i * W
        rect(slide, x + 0.04, 1.28, W - 0.08, 5.85, fill=WHITE, lc=STONE2)
        rect(slide, x + 0.04, 1.28, W - 0.08, 0.08, fill=ac)
        txt(slide, reg, x + 0.16, 1.45, W - 0.25, 0.38,
            sz=10, bold=True, color=ac)
        txt(slide, title, x + 0.16, 1.88, W - 0.25, 0.5,
            sz=13, bold=True, color=DARK)
        ml(slide, body.split('\n'), x + 0.16, 2.45, W - 0.25, 4.55,
           sz=10, color=STONE7, sp=4)

# ── SLIDE 9 — INTELLIGENCE ────────────────────────────────────────────────────
def s9_intelligence():
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
    bars(slide)
    txt(slide, "FINANCIAL INTELLIGENCE  ·  BEYOND BALANCES", 0.55, 0.18, 12, 0.35,
        sz=9, bold=True, color=RED)
    txt(slide, "The prototype surfaces insights that existing banking apps cannot.",
        0.55, 0.52, 12, 0.65, sz=24, bold=True, color=WHITE)
    rect(slide, 0.55, 1.28, 12.23, 0.04, fill=RED)
    features = [
        ("13-Week Cash Flow Forecast",
         "Smooth gradient area chart — Catmull-Rom curve.\n"
         "Zoomed to actual min–max band so movement is legible.\n"
         "Event markers: outflows, receipts, lowest week.\n"
         "'Largest outflow: Payroll £42,180 in week 3 (15 Jul).'\n"
         "'Balance dips below £80k floor — consider Reserve transfer.'",
         RED),
        ("Business Health Score",
         "Live 0–100 score computed from 5 weighted factors:\n"
         "  · Liquidity ratio\n  · Tax compliance\n  · Cash flow volatility\n"
         "  · Payment timeliness\n  · Mandate completeness\n"
         "Shown as a circular SVG gauge on the Home screen.",
         AMBER),
        ("Supplier Radar",
         "Counterparty intelligence based on payment history.\n"
         "Shows top payees, payment reliability, and flags.\n"
         "Identifies new counterparties for KYC review.\n"
         "Surfaces supplier payment terms vs actual payment dates.",
         BLUE),
        ("Pre-Approved Lending",
         "Based on 6-month transaction analysis, not a credit application.\n"
         "£45,000 from 6.9% APR shown on the Home screen.\n"
         "Draw-down in 4 steps — no call, no branch visit.\n"
         "Replaces 'arrange an appointment' with 'tap to accept'.",
         EMRLD),
    ]
    W = 12.23 / 4
    for i, (title, body, ac) in enumerate(features):
        x = 0.55 + i * W
        rect(slide, x + 0.04, 1.42, W - 0.08, 5.6, fill=DARK2)
        txt(slide, title, x + 0.16, 1.58, W - 0.25, 0.65,
            sz=11, bold=True, color=ac)
        ml(slide, body.split('\n'), x + 0.16, 2.3, W - 0.25, 4.65,
           sz=10, color=RGBColor(0xB0, 0xAA, 0xA4), sp=4)

# ── SLIDE 10 — THE BUSINESS CASE ─────────────────────────────────────────────
def s10_business_case():
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, fill=WARM)
    bars(slide)
    txt(slide, "THE BUSINESS CASE", 0.55, 0.18, 12, 0.35, sz=9, bold=True, color=RED)
    txt(slide, "Phase 1 is live. Phase 2 is a product decision.",
        0.55, 0.42, 12.2, 0.65, sz=26, bold=True, color=DARK)
    rule_line(slide, t=1.15)
    # Left: business case numbers
    rect(slide, 0.55, 1.28, 5.9, 5.85, fill=DARK)
    txt(slide, "THE NUMBERS", 0.75, 1.42, 5.5, 0.35, sz=9, bold=True, color=RED)
    numbers = [
        ("£137M", "annualised opportunity"),
        ("280,000", "UK SMB customers"),
        ("£490", "avg revenue per customer today"),
        ("5–7×", "expected uplift from self-service"),
        ("0 days", "to readiness — prototype is live"),
    ]
    for i, (val, lbl) in enumerate(numbers):
        y = 1.88 + i * 0.95
        rect(slide, 0.72, y, 1.8, 0.78, fill=DARK2)
        txt(slide, val, 0.79, y + 0.04, 1.7, 0.55, sz=22, bold=True, color=WHITE)
        txt(slide, lbl, 2.65, y + 0.18, 3.5, 0.45, sz=11, color=RGBColor(0xB0, 0xAA, 0xA4))
    # Right: phases
    rect(slide, 6.7, 1.28, 6.25, 5.85, fill=WHITE, lc=STONE2)
    txt(slide, "WHAT COMES NEXT", 6.9, 1.42, 5.85, 0.35, sz=9, bold=True, color=RED)
    phases = [
        ("PHASE 1  ·  COMPLETE", "Prototype built and live.\nAll 11 workflows functional.\n7 entity types supported.\nDesign system documented.", EMRLD),
        ("PHASE 2  ·  PILOT", "Select 500 business customers.\nConnect to real Santander accounts.\nRun for 90 days.\nMeasure: task completion, NPS, cost-to-serve.", AMBER),
        ("PHASE 3  ·  SCALE", "Roll out to all 280,000 SMB customers.\nIntegrate HMRC MTD API (production).\nEnable real-time mandate changes.\nTarget: £137M annual uplift.", RED),
    ]
    y = 1.88
    for title, body, ac in phases:
        rect(slide, 6.88, y, 5.88, 1.6, fill=WARM)
        rect(slide, 6.88, y, 0.08, 1.6, fill=ac)
        txt(slide, title, 7.1, y + 0.1, 5.5, 0.38, sz=10, bold=True, color=ac)
        ml(slide, body.split('\n'), 7.1, y + 0.52, 5.5, 1.0,
           sz=10, color=STONE7, sp=4)
        y += 1.72

# ── SLIDE 11 — ROADMAP ────────────────────────────────────────────────────────
def s11_roadmap():
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
    bars(slide)
    txt(slide, "ROADMAP  ·  TWELVE MONTHS TO PRODUCTION AT SCALE",
        0.55, 0.18, 12, 0.35, sz=9, bold=True, color=RED)
    txt(slide, "Four phases. Three decisions.",
        0.55, 0.52, 12, 0.65, sz=28, bold=True, color=WHITE)
    rect(slide, 0.55, 1.28, 12.23, 0.04, fill=RED)
    quarters = [
        ("Q3 2026\nJUL–SEP", "PHASE 1\nCOMPLETE",
         "✓  Prototype live\n✓  11 workflows built\n✓  7 entity types\n✓  Design system\n✓  Architecture documented",
         EMRLD, True),
        ("Q3–Q4 2026\nSEP–DEC", "PHASE 2\nPILOT",
         "500 business customers\nReal account integration\n90-day pilot\nNPS baseline\nCost-to-serve measurement",
         AMBER, False),
        ("Q1 2027\nJAN–MAR", "PHASE 3\nPRODUCTION",
         "Full Santander integration\nHMRC MTD live\nMandate API live\nComm strategy\nTraining rollout",
         RED, False),
        ("Q2 2027\nAPR–JUN", "PHASE 4\nSCALE",
         "All 280,000 SMB customers\nFull feature parity\nInternational expansion\nTarget: £137M uplift\nPhase 1 cost: recovered",
         BLUE, False),
    ]
    W = 12.23 / 4
    for i, (period, title, body, ac, done) in enumerate(quarters):
        x = 0.55 + i * W
        bg = DARK2 if not done else RGBColor(0x07, 0x32, 0x1A)
        rect(slide, x + 0.04, 1.42, W - 0.08, 5.6, fill=bg)
        rect(slide, x + 0.04, 1.42, W - 0.08, 0.09, fill=ac)
        txt(slide, period, x + 0.16, 1.56, W - 0.25, 0.7,
            sz=10, bold=True, color=ac)
        txt(slide, title, x + 0.16, 2.32, W - 0.25, 0.7,
            sz=13, bold=True, color=WHITE)
        ml(slide, body.split('\n'), x + 0.16, 3.1, W - 0.25, 3.8,
           sz=11, color=RGBColor(0xB0, 0xAA, 0xA4), sp=5)
        if done:
            txt(slide, "LIVE", x + 0.16, 6.6, W - 0.25, 0.35,
                sz=11, bold=True, color=EMRLD)

# ── SLIDE 12 — DESIGN SYSTEM ─────────────────────────────────────────────────
def s12_design():
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, fill=WARM)
    bars(slide)
    txt(slide, "DESIGN SYSTEM  ·  BUILT FOR EVERY CUSTOMER", 0.55, 0.18, 12, 0.35,
        sz=9, bold=True, color=RED)
    txt(slide, "Production-grade design tokens — consistent across all 11 workflows.",
        0.55, 0.42, 12.2, 0.65, sz=22, bold=True, color=DARK)
    rule_line(slide, t=1.15)
    # Colour swatches
    swatches = [
        ("#C8102E", "Brand Red\nCTAs · Active states\nTop bar · Alerts", RED),
        ("#1C1917", "Stone-900\nPrimary text\nDark cards · Filled buttons", DARK),
        ("#FAF6EF", "Warm Off-White\nPage background\nCard backgrounds", RGBColor(0x78, 0x71, 0x6C)),
        ("#D97706", "Amber-600\nCooling-off · Warnings\nPayroll nudge cards", AMBER),
        ("#2563EB", "Blue-600\nCompanies House\nUnfreeze button", BLUE),
        ("#059669", "Emerald-600\nSuccess states\nLoan confirmed\nHealth score", EMRLD),
    ]
    W = 12.23 / 6
    for i, (hex_val, lbl, color) in enumerate(swatches):
        x = 0.55 + i * W
        rect(slide, x + 0.04, 1.28, W - 0.08, 1.6, fill=color)
        txt(slide, hex_val, x + 0.1, 2.98, W - 0.12, 0.35, sz=9, bold=True, color=DARK)
        ml(slide, lbl.split('\n'), x + 0.1, 3.35, W - 0.12, 1.2, sz=9, color=STONE5, sp=3)
    # Typography
    rect(slide, 0.55, 4.7, 5.9, 2.43, fill=WHITE, lc=STONE2)
    txt(slide, "TYPOGRAPHY", 0.75, 4.85, 5.5, 0.32, sz=9, bold=True, color=RED)
    ml(slide, [
        "Fraunces  — editorial headings (display)",
        "Geist Sans  — body text, UI copy",
        "Geist Mono  — account numbers, financial figures",
        "num-tab class  — tabular figures on all monetary values",
    ], 0.75, 5.22, 5.5, 1.7, sz=11, color=STONE7, sp=5)
    # Accessibility
    rect(slide, 6.7, 4.7, 6.05, 2.43, fill=WHITE, lc=STONE2)
    txt(slide, "ACCESSIBILITY  ·  WCAG 2.1 AA", 6.9, 4.85, 5.65, 0.32, sz=9, bold=True, color=RED)
    ml(slide, [
        "focus-visible ring on every interactive element (keyboard nav)",
        "Dyslexia-friendly font  ·  Reduce motion  ·  High contrast  ·  Large text",
        "SVG charts: role='img' + aria-label on every chart",
        "Colour contrast checked against warm-white (#FAF6EF) background",
    ], 6.9, 5.22, 5.65, 1.7, sz=11, color=STONE7, sp=5)

# ── SLIDE 13 — THE ASK ────────────────────────────────────────────────────────
def s13_ask():
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
    bars(slide)
    rect(slide, 0.55, 1.8, 0.18, 4.0, fill=RED)
    txt(slide, "£137M. Live today. Zero days to readiness.",
        0.9, 1.9, 11.8, 2.2, sz=44, bold=True, color=WHITE)
    rect(slide, 0.9, 4.2, 11.5, 0.04, fill=RED)
    txt(slide, "Three decisions for senior management:",
        0.9, 4.4, 11.5, 0.5, sz=16, color=RGBColor(0xC0, 0xBB, 0xB5))
    decisions = [
        ("01", "Approve Phase 2 pilot",
         "500 customers · 90 days · real accounts · measure NPS and cost-to-serve"),
        ("02", "Assign a product team",
         "The prototype needs an engineer, a PM, and a designer to move to production"),
        ("03", "Name a sponsor",
         "A senior stakeholder to own the business case through to Phase 4"),
    ]
    for i, (num, title, body) in enumerate(decisions):
        x = 0.9 + i * 4.12
        rect(slide, x, 4.98, 3.85, 2.0, fill=DARK2)
        txt(slide, num, x + 0.15, 5.1, 0.6, 0.55, sz=24, bold=True, color=RED)
        txt(slide, title, x + 0.15, 5.7, 3.55, 0.45, sz=12, bold=True, color=WHITE)
        txt(slide, body, x + 0.15, 6.2, 3.55, 0.68, sz=10,
            color=RGBColor(0xA8, 0xA2, 0x9A))

# ── SLIDE 14 — BACK COVER ─────────────────────────────────────────────────────
def s14_cover():
    slide = prs.slides.add_slide(BLANK)
    rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
    bars(slide)
    rect(slide, 0.55, 2.8, 0.18, 2.0, fill=RED)
    txt(slide, "See it. Try it. Build it.", 0.9, 2.9, 11, 1.1,
        sz=52, bold=True, color=WHITE)
    txt(slide, "imscots-ui.github.io/santander",
        0.9, 4.15, 11, 0.52, sz=20, bold=True, color=RED)
    txt(slide, "Alan Davidson  ·  Alan.Davidson@santander.co.uk",
        0.9, 4.82, 11, 0.45, sz=14, color=RGBColor(0xC0, 0xBB, 0xB5))
    txt(slide, "Self-initiated concept piece. Prototype only — no backend, no real data. Built to demonstrate what's possible.",
        0.9, 6.75, 11.5, 0.45, sz=9, italic=True, color=STONE5)

# ── BUILD ─────────────────────────────────────────────────────────────────────
s1_cover()
s2_number()
s3_problem()
s4_insight()
s5_prototype()
s6_command_centre()
s7_mtd()
s8_compliance()
s9_intelligence()
s10_business_case()
s11_roadmap()
s12_design()
s13_ask()
s14_cover()

out = "/home/user/santander/Santander_Digital_Banking_Future.pptx"
prs.save(out)
from pptx import Presentation as _P
import os
p = _P(out)
print(f"Saved: {out}")
print(f"Slides: {len(p.slides)}")
print(f"Size: {os.path.getsize(out)//1024} KB")
