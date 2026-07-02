#!/usr/bin/env python3
"""
build_changes_july.py — generates Santander_July2026_Changes.pptx

A board-facing "What Changed" deck covering every material update to the
Santander Business Banking prototype SINCE the June 2026 review (the frozen
Santander_June2026_Changes.pptx is left untouched as the historical record).

House style mirrors the June deck:
  red top stripe (#EC0000) · Lisbon off-white body (#FBF1EA) · dark charcoal
  titles · red section labels · footer strip with credit/classification.
Red full-bleed section dividers separate the four themes.

Run:  python3 build_changes_july.py
"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Official palette ──────────────────────────────────────────────────────────
RED      = RGBColor(0xEC, 0x00, 0x00)   # Santander Red (Pantone 485 C)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LISBON   = RGBColor(0xFB, 0xF1, 0xEA)   # official page background
DARK     = RGBColor(0x1C, 0x19, 0x17)
MUTED    = RGBColor(0x78, 0x71, 0x6C)
SUBTLE   = RGBColor(0xD6, 0xD3, 0xD1)
FOOTERBG = RGBColor(0xF1, 0xEA, 0xE1)

# ── Geometry (matches June deck) ──────────────────────────────────────────────
W, H       = 12192000, 6858000
STRIPE_H   = 73152
MARGIN     = 457200
BODY_W     = W - 2 * MARGIN
TEXT_W     = BODY_W - 182880
INDENT_W   = BODY_W - 365760
FOOTER_TOP = 6217920
FOOTER_H   = 457200

CREDIT = "INTERNAL · CONFIDENTIAL   ·   A. Davidson · July 2026"


def add_rect(slide, l, t, w, h, rgb):
    sp = slide.shapes.add_shape(1, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = rgb
    sp.line.fill.background(); sp.shadow.inherit = False
    return sp


def add_text(slide, l, t, w, h, text, size, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, wrap=True, anchor=None):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tb


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def cover(prs):
    s = blank(prs)
    add_rect(s, 0, 0, W, H, DARK)
    add_rect(s, 0, 0, 137160, H, RED)           # red left stripe
    s.shapes.add_picture('brand/santander-logo-white.png',
                         MARGIN, 548640, Emu(2377440), Emu(int(2377440 * 519 / 3000)))
    add_text(s, MARGIN, 2286000, 10000000, 900000,
             "What Changed — July 2026", 40, bold=True, color=WHITE)
    add_text(s, MARGIN, 3520440, 9600000, 700000,
             "Everything that moved on the Business Banking prototype since the June "
             "review — product, compliance, brand and the board pack.", 15, color=SUBTLE)
    add_text(s, MARGIN, 6217920, 6000000, 360000,
             "July 2026  ·  Confidential", 10, color=MUTED)


def divider(prs, number, title, subtitle):
    s = blank(prs)
    add_rect(s, 0, 0, W, H, RED)
    add_text(s, MARGIN, 2560320, 10500000, 800000,
             f"{number}  {title}", 34, bold=True, color=WHITE)
    add_rect(s, MARGIN, 3520440, 2286000, 27432, WHITE)
    add_text(s, MARGIN, 3703320, 9600000, 500000, subtitle, 14, color=WHITE)


def content(prs, num, total, title, sections, footer):
    s = blank(prs)
    add_rect(s, 0, 0, W, STRIPE_H, RED)
    add_rect(s, 0, STRIPE_H, W, H - STRIPE_H, LISBON)
    add_text(s, MARGIN, STRIPE_H + 18000, 4000000, 300000,
             "SINCE THE JUNE 2026 REVIEW", 9, bold=True, color=RED)
    add_text(s, W - MARGIN - 800000, STRIPE_H + 18000, 760000, 300000,
             f"{num} / {total}", 9, color=MUTED, align=PP_ALIGN.RIGHT)
    add_text(s, MARGIN, 228600, TEXT_W, 620000, title, 26, bold=True, color=DARK)
    add_rect(s, MARGIN, 868680, BODY_W, 36576, RED)

    y = 1097280
    HEAD_H, BULLET_H, GAP = 350000, 288000, 66000
    for sec in sections:
        add_text(s, MARGIN, y, TEXT_W, HEAD_H, sec['heading'], 14, bold=True, color=RED)
        y += HEAD_H
        for b in sec['bullets']:
            add_text(s, MARGIN + 182880, y, INDENT_W, BULLET_H, "▪  " + b, 12, color=DARK)
            y += BULLET_H
        y += GAP

    add_rect(s, MARGIN, FOOTER_TOP, BODY_W, FOOTER_H, FOOTERBG)
    add_text(s, MARGIN + 182880, FOOTER_TOP + 100000, 6600000, 320000, footer, 10, color=MUTED)
    add_text(s, MARGIN + BODY_W - 4500000, FOOTER_TOP + 100000, 4318000, 320000,
             CREDIT, 9, bold=True, color=RED, align=PP_ALIGN.RIGHT)


# ── Deck ──────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = W; prs.slide_height = H
TOTAL = 14

cover(prs)

content(prs, 2, TOTAL, "Since the last review — at a glance", [
    {'heading': 'Product', 'bullets': [
        "Workflows: 11 → 17 — six new end-to-end paperless journeys",
        "Features: 21 → 36 across the four product categories"]},
    {'heading': 'Quality & compliance', 'bullets': [
        "WCAG 2.1 AA contrast pass — 54 issues resolved to 0 on the core screens",
        "Runtime drive-through: 20/20 automated checks, with an axe audit on every build"]},
    {'heading': 'Brand', 'bullets': [
        "Official Santander brand kit adopted from the guidelines — red, full palette, typefaces, logo"]},
    {'heading': 'Board pack', 'bullets': [
        "New one-page Executive Brief; Technical Report & Position Paper rebuilt to v2.0"]},
], "Baseline: the frozen 'What Changed — June 2026' record. This deck covers everything since.")

divider(prs, "01", "Product", "Six new workflows and a truthful feature count")

content(prs, 4, TOTAL, "Six new workflows  (11 → 17)", [
    {'heading': 'Added this cycle', 'bullets': [
        "Payment dispute & chargeback — provisional refund; fraud / DD-Guarantee routes",
        "International beneficiary onboarding — IBAN/SWIFT validation, sanctions screening",
        "Balance certificate — sealed proof of balance, issued instantly (no postal wait)",
        "Trusted devices & sessions — SCA-protected sign-out and device revocation",
        "Complaint handling — FCA DISP intake → triage → outcome, with FOS route",
        "Standing orders & Direct Debits — self-serve set-up and cancellation"]},
    {'heading': 'Result', 'bullets': [
        "Every routine business-banking service event now has a paperless path — seventeen in total"]},
], "App.jsx · renderDispute / renderBeneficiary / renderCertificate / renderTrusted / renderComplaint / renderRecurring")

content(prs, 5, TOTAL, "Feature set  (21 → 36)", [
    {'heading': 'What changed', 'bullets': [
        "The pitch deck's feature grid was rebuilt to reflect what has actually shipped",
        "Fifteen intelligence, workflow and security features were previously undocumented",
        "e.g. balance certificate, beneficiary screening, dispute triage, trusted devices, ID register"]},
    {'heading': 'Why it matters', 'bullets': [
        "The headline number now matches the build — no gap between the claim and the product",
        "Every current-state figure across app, decks, reports and docs traces to a single source"]},
], "Santander_Digital_Banking_Future.pptx s5 · Santander_Project_Record.pptx")

divider(prs, "02", "Quality & Compliance", "Accessibility proven, not asserted")

content(prs, 7, TOTAL, "Accessibility — WCAG 2.1 AA contrast pass", [
    {'heading': 'Before', 'bullets': [
        "54 colour-contrast failures across screens — small grey text on warm and red surfaces",
        "Contrast was assumed from the palette, never measured"]},
    {'heading': 'After', 'bullets': [
        "0 actionable violations on the core screens; text tokens moved to AA-compliant shades",
        "An automated axe audit now runs on every build — regressions are caught, not shipped"]},
    {'heading': 'Why it matters', 'bullets': [
        "A bank's digital estate must meet WCAG 2.1 AA; this makes conformance demonstrable"]},
], "scripts/verify-flows.mjs · axe-core · every core screen re-measured")

content(prs, 8, TOTAL, "Runtime verification harness", [
    {'heading': 'What was added', 'bullets': [
        "The built app is launched headless (Chromium) and driven through every workflow",
        "The sign-and-send flow and all accordion actions are exercised as a user would",
        "A deck-render + full-height screenshot pipeline now catches slide overflow before it ships"]},
    {'heading': 'Result', 'bullets': [
        "20/20 checks pass, including the WCAG audit — behaviour is proven on every build",
        "Two renderer blind spots (clipped overflow; under-painted screenshots) found and closed"]},
], "scripts/verify-flows.mjs · scripts/deck_shot.py · scripts/deck-qa.py")

divider(prs, "03", "Brand", "The genuine article, from the guidelines")

content(prs, 10, TOTAL, "Official Santander brand kit adopted", [
    {'heading': 'From the Basic Guidelines for Santander', 'bullets': [
        "Santander Red #EC0000 (Pantone 485 C) — replaces the earlier #DA291C / #c8102e approximations",
        "Full palette: Boston, London, Madrid, Mexico City, Rio and Lisbon; Lisbon is the page background",
        "Typefaces named — Santander Text / Headline / Logo (Geist as the open fallback)",
        "Official flame logo in red, white and black; correct clear-space and minimum sizes"]},
    {'heading': 'Accessibility-safe & applied everywhere', 'bullets': [
        "Small red text uses Boston #CC0000 for AA contrast; verified against the palette",
        "Applied across the app, all five decks, both reports and the logo files"]},
], "DECISIONS.md · CITATIONS.md · brand/ · verified via axe on the new palette")

divider(prs, "04", "Board Pack & Accuracy", "One consistent, board-ready set")

content(prs, 12, TOTAL, "New & refreshed deliverables", [
    {'heading': 'New', 'bullets': [
        "One-page Executive Brief — the ask, the numbers and the evidence base (board pre-read)"]},
    {'heading': 'Rebuilt', 'bullets': [
        "Technical Report & Position Paper at v2.0 — senior-management register, real citations",
        "Clean-room content: original structure, no third-party derivation"]},
    {'heading': 'Refreshed', 'bullets': [
        "All decks trued to 17 workflows / 36 features; dates advanced to July 2026",
        "Illegible regulation-tag pills fixed; covers and stat tiles corrected"]},
], "Santander_Executive_Brief.pdf · *_Technical_Report / *_Position_Paper (v2.0)")

content(prs, 13, TOTAL, "What has NOT changed", [
    {'heading': 'Architecture', 'bullets': [
        "Still a single App.jsx (~7,400 lines) — no modularisation; deliberate and documented"]},
    {'heading': 'Product invariants', 'bullets': [
        "All 17 workflows intact; all 7 entity types with correct mandate and compliance divergence",
        "Front-end only — no backend, no real data, no live connections"]},
    {'heading': 'Deployment', 'bullets': [
        "GitHub Pages, unchanged — imscots-ui.github.io/santander — runs from any browser"]},
], "Continuity: the foundations held; the work above sits on top of them.")

content(prs, 14, TOTAL, "Where we are now", [
    {'heading': 'The build', 'bullets': [
        "17 workflows  ·  36 features  ·  7 entity types  ·  8 regulatory frameworks",
        "WCAG 2.1 AA verified  ·  official Santander brand throughout  ·  zero build cost"]},
    {'heading': 'The pack', 'bullets': [
        "Executive Brief, three decks, two reports and the walkthrough — internally consistent",
        "Every number traces to one source of truth"]},
    {'heading': 'Ready', 'bullets': [
        "Board pack complete and verified — ready for the Friday 11:00 briefing"]},
], "Live prototype: imscots-ui.github.io/santander  ·  open it in the room")

prs.save('Santander_July2026_Changes.pptx')
print("WROTE Santander_July2026_Changes.pptx", len(prs.slides._sldIdLst), "slides")
